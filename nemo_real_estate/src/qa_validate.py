"""
QA 검증 스크립트
- CSV 데이터 검증
- DOCX 보고서 검증
- PPTX 프레젠테이션 검증
- HTML 대시보드 검증
"""
import sys
import os
import json
from datetime import datetime

# 프로젝트 루트 기준 상대경로
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

results = {
    "검증_일시": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    "csv": {},
    "docx": {},
    "pptx": {},
    "dashboard": {},
}

print("=" * 60)
print("QA 검증 시작")
print("=" * 60)

# ─────────────────────────────────────────────
# Step 1: CSV 데이터 검증
# ─────────────────────────────────────────────
print("\n[Step 1] CSV 데이터 검증")
try:
    import pandas as pd
    csv_path = os.path.join(BASE_DIR, "data", "nemo_real_estate_bestseller.csv")
    df = pd.read_csv(csv_path)
    print(f"  - 파일 경로: {csv_path}")
    print(f"  - 데이터 shape: {df.shape}")
    print(f"  - 컬럼: {df.columns.tolist()}")
    print(f"  - 결측값 현황:\n{df.isnull().sum()}")
    print(f"\n  상위 3행:\n{df.head(3)}")
    print(f"\n  기술통계:\n{df.describe()}")
    
    missing = df.isnull().sum().to_dict()
    missing_cols = {k: int(v) for k, v in missing.items() if v > 0}
    
    results["csv"] = {
        "건수": int(df.shape[0]),
        "컬럼수": int(df.shape[1]),
        "컬럼목록": df.columns.tolist(),
        "결측값": missing_cols,
        "중복행": int(df.duplicated().sum()),
        "상태": "정상",
    }
    print("  ✅ CSV 데이터 검증 완료")
except Exception as e:
    print(f"  ❌ CSV 검증 오류: {e}")
    results["csv"] = {"상태": "오류", "오류내용": str(e)}

# ─────────────────────────────────────────────
# Step 2: DOCX 파일 검증
# ─────────────────────────────────────────────
print("\n[Step 2] DOCX 보고서 검증")
try:
    from docx import Document
    from docx.oxml.ns import qn
    
    docx_path = os.path.join(BASE_DIR, "reports", "real_estate_report.docx")
    doc = Document(docx_path)
    
    # 단락 수
    paragraph_count = len(doc.paragraphs)
    # 표 수
    table_count = len(doc.tables)
    # 이미지 수 (관계에서 이미지 타입 필터)
    rels = doc.part.rels
    image_rels = [r for r in rels.values() if "image" in r.reltype]
    image_count = len(image_rels)
    
    print(f"  - 단락 수: {paragraph_count}")
    print(f"  - 표 수: {table_count}")
    print(f"  - 이미지 수: {image_count}")
    
    # 첫 10개 단락 출력
    print("  - 상위 단락:")
    non_empty_paras = [p for p in doc.paragraphs if p.text.strip()]
    for p in non_empty_paras[:10]:
        print(f"    [{p.style.name}] {p.text[:80]}")
    
    # 문제점 판단
    issues = []
    if image_count == 0:
        issues.append("이미지가 없음")
    if table_count == 0:
        issues.append("표가 없음")
    if paragraph_count < 20:
        issues.append(f"단락 수({paragraph_count})가 너무 적음(기준: 20개 이상)")
    
    results["docx"] = {
        "단락수": paragraph_count,
        "표수": table_count,
        "이미지수": image_count,
        "문제점": issues,
        "상태": "문제있음" if issues else "정상",
    }
    
    if issues:
        print(f"  ⚠️  문제점: {issues}")
    else:
        print("  ✅ DOCX 검증 완료 (정상)")

except Exception as e:
    print(f"  ❌ DOCX 검증 오류: {e}")
    results["docx"] = {"상태": "오류", "오류내용": str(e)}

# ─────────────────────────────────────────────
# Step 3: PPTX 파일 검증
# ─────────────────────────────────────────────
print("\n[Step 3] PPTX 프레젠테이션 검증")
try:
    from pptx import Presentation
    
    pptx_path = os.path.join(BASE_DIR, "reports", "real_estate_presentation.pptx")
    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    print(f"  - 슬라이드 수: {slide_count}")
    
    slide_details = []
    total_images = 0
    for i, slide in enumerate(prs.slides):
        shapes = slide.shapes
        texts = [s.text_frame.text[:50] for s in shapes if s.has_text_frame and s.text_frame.text.strip()]
        imgs = [s for s in shapes if s.shape_type == 13]  # Picture = 13
        total_images += len(imgs)
        slide_details.append({
            "슬라이드": i + 1,
            "텍스트수": len(texts),
            "이미지수": len(imgs),
        })
        print(f"  슬라이드 {i+1}: 텍스트 {len(texts)}개, 이미지 {len(imgs)}개")
        for t in texts[:3]:
            print(f"    - {t}")
    
    issues = []
    if slide_count < 8:
        issues.append(f"슬라이드 수({slide_count})가 8개 미만")
    if total_images == 0:
        issues.append("이미지가 없음")
    
    results["pptx"] = {
        "슬라이드수": slide_count,
        "전체이미지수": total_images,
        "슬라이드별상세": slide_details,
        "문제점": issues,
        "상태": "문제있음" if issues else "정상",
    }
    
    if issues:
        print(f"  ⚠️  문제점: {issues}")
    else:
        print("  ✅ PPTX 검증 완료 (정상)")

except Exception as e:
    print(f"  ❌ PPTX 검증 오류: {e}")
    results["pptx"] = {"상태": "오류", "오류내용": str(e)}

# ─────────────────────────────────────────────
# Step 4: HTML 대시보드 검증
# ─────────────────────────────────────────────
print("\n[Step 4] HTML 대시보드 검증")
try:
    html_path = os.path.join(BASE_DIR, "src", "dashboard.html")
    js_path = os.path.join(BASE_DIR, "src", "dashboard_data.js")
    
    with open(html_path, "r", encoding="utf-8") as f:
        html_content = f.read()
    
    # 검증 항목
    checks = {
        "외부_JS_참조": "dashboard_data.js" in html_content,
        "DASHBOARD_LISTINGS_참조": "DASHBOARD_LISTINGS" in html_content,
        "window_DASHBOARD_LISTINGS": "window.DASHBOARD_LISTINGS" in html_content or "DASHBOARD_LISTINGS" in html_content,
        "ChartJS": "chart.js" in html_content.lower() or "Chart" in html_content,
        "ECharts": "echarts" in html_content.lower(),
        "필터_기능": "filter" in html_content.lower() or "필터" in html_content,
        "역_필터": "station" in html_content.lower() or "역" in html_content,
        "거래유형_필터": "deal_type" in html_content.lower() or "거래" in html_content,
        "금액_포맷_함수": "formatPrice" in html_content or "억" in html_content or "만원" in html_content,
        "다크_라이트_테마": "dark" in html_content.lower() or "theme" in html_content.lower(),
    }
    
    for key, val in checks.items():
        status = "✅" if val else "❌"
        print(f"  {status} {key}: {val}")
    
    # dashboard_data.js 검증
    if os.path.exists(js_path):
        with open(js_path, "r", encoding="utf-8") as f:
            js_content = f.read()
        js_size_kb = os.path.getsize(js_path) / 1024
        has_listings = "DASHBOARD_LISTINGS" in js_content
        has_summary = "DASHBOARD_SUMMARY" in js_content or "summary" in js_content.lower()
        print(f"\n  dashboard_data.js:")
        print(f"    - 크기: {js_size_kb:.1f} KB")
        print(f"    - DASHBOARD_LISTINGS 포함: {has_listings}")
        print(f"    - DASHBOARD_SUMMARY 포함: {has_summary}")
        
        # 레코드 수 추정 (줄 수로 대략적 파악)
        lines = js_content.count("\n")
        print(f"    - 줄 수: {lines}")
    else:
        print("  ❌ dashboard_data.js 파일이 없음")
        has_listings = False
        has_summary = False
        js_size_kb = 0
        lines = 0
    
    issues = []
    if not checks["외부_JS_참조"] and not checks["DASHBOARD_LISTINGS_참조"]:
        issues.append("데이터 파일 참조 없음")
    if not checks["ChartJS"] and not checks["ECharts"]:
        issues.append("차트 라이브러리 감지 안됨")
    if not checks["필터_기능"]:
        issues.append("필터 기능 코드 없음")
    if not checks["금액_포맷_함수"]:
        issues.append("금액 포맷 함수 없음")
    
    results["dashboard"] = {
        "HTML_크기KB": round(os.path.getsize(html_path) / 1024, 1),
        "JS_크기KB": round(js_size_kb, 1),
        "검증항목": {k: bool(v) for k, v in checks.items()},
        "문제점": issues,
        "상태": "문제있음" if issues else "정상",
    }
    
    if issues:
        print(f"  ⚠️  문제점: {issues}")
    else:
        print("  ✅ 대시보드 검증 완료 (정상)")

except Exception as e:
    print(f"  ❌ 대시보드 검증 오류: {e}")
    results["dashboard"] = {"상태": "오류", "오류내용": str(e)}

# ─────────────────────────────────────────────
# 결과 저장
# ─────────────────────────────────────────────
result_path = os.path.join(BASE_DIR, "reports", "qa_validate_result.json")
with open(result_path, "w", encoding="utf-8") as f:
    json.dump(results, f, ensure_ascii=False, indent=2)

print("\n" + "=" * 60)
print(f"✅ 검증 결과 저장: {result_path}")
print("=" * 60)

# 결과 요약 출력
print("\n📊 QA 검증 요약:")
for section, data in results.items():
    if isinstance(data, dict) and "상태" in data:
        status_icon = "✅" if data["상태"] == "정상" else "⚠️ " if data["상태"] == "문제있음" else "❌"
        print(f"  {status_icon} {section}: {data['상태']}")
        if "문제점" in data and data["문제점"]:
            for issue in data["문제점"]:
                print(f"      → {issue}")
