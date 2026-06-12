# -*- coding: utf-8 -*-
"""
광화문역 vs 강남역 부동산 매물 비교 분석 발표 슬라이드(.pptx) 생성 스크립트

이 스크립트는 수집된 663개의 부동산 매물 데이터(보증금, 월세, 지역, 상세 텍스트)를 활용하여
8슬라이드 분량의 깔끔한 상권 분석 발표 장표를 nemo_real_estate/reports/real_estate_presentation.pptx에 자동 생성합니다.
pptx 및 pptx-design-styles 가이드라인을 참고하여 세련된 'Nordic Minimalism' 테마를 적용합니다.
(배경색: Off-white, 주색: Midnight Blue, 강조색: Warm Terracotta, 보조색: Sage)

동적 데이터 분석 연동, 겹침 방지 레이아웃, 대형 KPI 카드 시각화, 하단 캡션 및 시그니처 3-dot 도트 장식 포함.

작성자: Antigravity AI Data Pipeline Framework
작성일: 2026-06-12
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ==================== 색상 및 폰트 테마 설정 (Nordic Minimalism) ====================
COLOR_BG_DARK = RGBColor(0x1E, 0x27, 0x61)      # Midnight Blue (표지/종결 슬라이드용)
COLOR_BG_LIGHT = RGBColor(0xF4, 0xF1, 0xEC)     # Off-white / Warm Cream (본문 슬라이드용)
COLOR_TEXT_DARK = RGBColor(0x3D, 0x35, 0x30)    # 거의 검은색 (본문 기본 텍스트)
COLOR_TEXT_MUTED = RGBColor(0x8A, 0x7A, 0x6A)   # 회색/Taupe (보조 텍스트)
COLOR_TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)   # 흰색 (어두운 배경용 텍스트)

COLOR_PRIMARY = RGBColor(0x1E, 0x27, 0x61)      # 주색 (Midnight Blue)
COLOR_SECONDARY = RGBColor(0x50, 0x80, 0x8E)    # 보조색 (Sage/Slate)
COLOR_ACCENT = RGBColor(0xB8, 0x50, 0x42)       # 강조색 (Warm Terracotta/Coral)

COLOR_CARD_BG = RGBColor(0xEA, 0xE6, 0xDF)      # 카드 배경색 (밝은 그레이/베이지)
COLOR_ACCENT_BG = RGBColor(0xF2, 0xD8, 0xD5)    # 옅은 강조색 배경 (Terracotta Tint)
COLOR_ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)     # 아이스블루

FONT_TITLE = "Georgia"
FONT_BODY = "Calibri"

# ==================== 데이터 분석 및 통계 계산 ====================
def load_statistics():
    """CSV 데이터를 로드하여 동적으로 실시간 상권 통계 산출"""
    csv_path = "nemo_real_estate/data/nemo_real_estate_bestseller.csv"
    if os.path.exists(csv_path):
        try:
            df = pd.read_csv(csv_path)
            
            # 전체 통계
            total_count = len(df)
            avg_deposit = df['deposit'].mean()
            avg_rent = df['monthlyRent'].mean()
            
            # 강남역 통계
            df_gn = df[df['region'] == '강남역']
            gn_count = len(df_gn)
            gn_avg_deposit = df_gn['deposit'].mean()
            gn_med_deposit = df_gn['deposit'].median()
            gn_avg_rent = df_gn['monthlyRent'].mean()
            gn_med_rent = df_gn['monthlyRent'].median()
            
            # 광화문역 통계
            df_gw = df[df['region'] == '광화문역']
            gw_count = len(df_gw)
            gw_avg_deposit = df_gw['deposit'].mean()
            gw_med_deposit = df_gw['deposit'].median()
            gw_avg_rent = df_gw['monthlyRent'].mean()
            gw_med_rent = df_gw['monthlyRent'].median()
            
            return {
                "total_count": total_count,
                "avg_deposit": avg_deposit,
                "avg_rent": avg_rent,
                "gn_count": gn_count,
                "gn_avg_deposit": gn_avg_deposit,
                "gn_med_deposit": gn_med_deposit,
                "gn_avg_rent": gn_avg_rent,
                "gn_med_rent": gn_med_rent,
                "gw_count": gw_count,
                "gw_avg_deposit": gw_avg_deposit,
                "gw_med_deposit": gw_med_deposit,
                "gw_avg_rent": gw_avg_rent,
                "gw_med_rent": gw_med_rent
            }
        except Exception as e:
            print(f"[Warning] CSV 로드 중 오류 발생: {e}. 하드코딩 수치로 Fallback합니다.")
            
    # Fallback 기본 데이터 (nemo_real_estate_bestseller.csv 분석 기반)
    return {
        "total_count": 663,
        "avg_deposit": 63071.0,
        "avg_rent": 4676.0,
        "gn_count": 400,
        "gn_avg_deposit": 66957.0,
        "gn_med_deposit": 50000.0,
        "gn_avg_rent": 5114.0,
        "gn_med_rent": 3520.0,
        "gw_count": 263,
        "gw_avg_deposit": 57161.0,
        "gw_med_deposit": 40000.0,
        "gw_avg_rent": 4010.0,
        "gw_med_rent": 2900.0
    }

# ==================== 금액 단위 포맷팅 헬퍼 함수 ====================
def format_money(val_manwon):
    """만원 단위 수치를 'X억 Y만 원' 또는 'Y만 원' 포맷으로 변경"""
    val = round(val_manwon)
    if val >= 10000:
        uk = val // 10000
        man = val % 10000
        if man > 0:
            return f"{uk:,}억 {man:,}만 원"
        else:
            return f"{uk:,}억 원"
    else:
        return f"{val:,}만 원"

# ==================== 디자인 및 레이아웃 헬퍼 함수 ====================
def set_slide_background(slide, color):
    """슬라이드 단색 배경색 지정"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_textbox(slide, left, top, width, height):
    """여백을 제거하여 정렬의 일치성과 오버플로우를 줄이는 텍스트박스 생성"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    return tf

def draw_card(slide, left, top, width, height, bg_color):
    """시각적 카드 레이아웃을 위한 둥근 사각형 도형 그리기"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background() # 테두리 제거 (미니멀리즘 스타일)
    return shape

def draw_accent_dots(slide, is_dark_bg=False):
    """슬라이드 상단 구석에 3가지 시그니처 테마 색상 도트 배치"""
    dot_y = Inches(0.4)
    dot_size = Inches(0.1)
    dot_spacing = Inches(0.18)
    
    if is_dark_bg:
        colors = [RGBColor(0xFF, 0xFF, 0xFF), COLOR_ACCENT, COLOR_SECONDARY]
        start_x = Inches(12.0)
    else:
        colors = [COLOR_PRIMARY, COLOR_SECONDARY, COLOR_ACCENT]
        start_x = Inches(0.8)
        
    for i, color in enumerate(colors):
        x = start_x + i * dot_spacing
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, dot_y, dot_size, dot_size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

def draw_organic_bg(slide, is_dark_bg=False):
    """배경에 옅은 유기적 큰 원을 배치하여 미니멀한 공간감 연출 (가장 먼저 호출하여 맨 아래 배치)"""
    if is_dark_bg:
        color = RGBColor(0x22, 0x2B, 0x68) # 어두운 네이비
        x, y = Inches(9.5), Inches(3.8)
        size = Inches(4.5)
    else:
        color = RGBColor(0xEA, 0xE7, 0xE1) # 옅은 그레이베이지
        x, y = Inches(10.0), Inches(3.8)
        size = Inches(4.5)
        
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def draw_bottom_caption(slide, is_dark_bg=False):
    """하단 얇은 분할선과 모노스페이스 스타일 캡션 및 출처 표기"""
    line_color = RGBColor(0x3D, 0x35, 0x30) if is_dark_bg else RGBColor(0xDF, 0xD9, 0xD2)
    text_color = RGBColor(0xCA, 0xDC, 0xFC) if is_dark_bg else COLOR_TEXT_MUTED
    
    # 얇은 가로 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.8), Inches(11.73), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = line_color
    line.line.fill.background()
    
    # 출처 캡션
    tf = create_textbox(slide, Inches(0.8), Inches(6.9), Inches(11.73), Inches(0.3))
    p = tf.paragraphs[0]
    p.text = "DATA SOURCE: NEMOAPP REAL ESTATE DATA  |  PREPARED BY ANTIGRAVITY DATA PIPELINE  |  STYLE: NORDIC MINIMALISM"
    p.font.name = FONT_BODY
    p.font.size = Pt(8.5)
    p.font.color.rgb = text_color

def add_header(slide, title_text, is_dark_bg=False):
    """슬라이드 상단 제목 추가 (강조선 없는 현대적 타이포그래피)"""
    tf = create_textbox(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_TITLE
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT if is_dark_bg else COLOR_PRIMARY

def add_bullet_point(tf, bold_prefix, desc_text, size_pt=13, color=COLOR_TEXT_DARK):
    """세련된 본문 단락 포인트 추가 (자간, 줄바꿈 고려)"""
    p = tf.add_paragraph()
    p.space_before = Pt(10)
    p.line_spacing = 1.3
    
    # 강조 프리픽스
    run_bold = p.add_run()
    run_bold.text = "•  " + bold_prefix + ": "
    run_bold.font.name = FONT_BODY
    run_bold.font.size = Pt(size_pt)
    run_bold.font.bold = True
    run_bold.font.color.rgb = color
    
    # 일반 설명
    run_desc = p.add_run()
    run_desc.text = desc_text
    run_desc.font.name = FONT_BODY
    run_desc.font.size = Pt(size_pt)
    run_desc.font.color.rgb = color

def add_image_auto_ratio(slide, img_path, left, top, width):
    """종횡비를 유지하며 이미지를 안전하게 삽입 (height 생략)"""
    if os.path.exists(img_path):
        return slide.shapes.add_picture(img_path, left, top, width=width)
    print(f"[Warning] 이미지가 존재하지 않습니다: {img_path}")
    return None

# ==================== 메인 빌드 로직 ====================
def build_presentation():
    output_pptx_path = "nemo_real_estate/reports/real_estate_presentation.pptx"
    os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
    
    # 데이터 로드
    stats = load_statistics()
    
    prs = Presentation()
    # 16:9 와이드스크린 규격 지정 (13.333" x 7.5")
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 공통 빈 슬라이드 레이아웃
    slide_layout = prs.slide_layouts[6]
    
    # ---------------- SLIDE 1: 표지 슬라이드 (Dark Background) ----------------
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1, COLOR_BG_DARK)
    draw_organic_bg(slide1, is_dark_bg=True)
    draw_accent_dots(slide1, is_dark_bg=True)
    
    # 메인 타이틀 블록
    tf1 = create_textbox(slide1, Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.0))
    p1 = tf1.paragraphs[0]
    p1.text = "광화문역 vs 강남역 상업 상권 비교 분석"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_LIGHT

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "네모앱 수집 매물 663건 기반 상업용 부동산 임대 가격 및 입지 제언"
    p1_sub.font.name = FONT_BODY
    p1_sub.font.size = Pt(16)
    p1_sub.font.color.rgb = COLOR_ICE_BLUE
    p1_sub.space_before = Pt(12)

    # [신규] 핵심 메시지 3가지 카드 — 의사결정자를 위한 Bottom Line
    rent_premium_pct_str = f"{(stats['gn_avg_rent'] - stats['gw_avg_rent']) / stats['gw_avg_rent'] * 100:.1f}%"

    # 카드 1
    draw_card(slide1, Inches(1.0), Inches(3.8), Inches(3.7), Inches(1.2), RGBColor(0x22, 0x2B, 0x68))
    tf1_c1 = create_textbox(slide1, Inches(1.15), Inches(3.92), Inches(3.4), Inches(1.0))
    p_c1_kw = tf1_c1.paragraphs[0]
    p_c1_kw.text = f"강남역 월세 프리미엄 {rent_premium_pct_str}"
    p_c1_kw.font.name = FONT_TITLE
    p_c1_kw.font.size = Pt(16)
    p_c1_kw.font.bold = True
    p_c1_kw.font.color.rgb = COLOR_ACCENT
    p_c1_desc = tf1_c1.add_paragraph()
    p_c1_desc.text = "광화문 대비 고비용, 고마진 업종에만 유효"
    p_c1_desc.font.name = FONT_BODY
    p_c1_desc.font.size = Pt(11)
    p_c1_desc.font.color.rgb = COLOR_ICE_BLUE
    p_c1_desc.space_before = Pt(4)

    # 카드 2
    draw_card(slide1, Inches(4.8), Inches(3.8), Inches(3.7), Inches(1.2), RGBColor(0x22, 0x2B, 0x68))
    tf1_c2 = create_textbox(slide1, Inches(4.95), Inches(3.92), Inches(3.4), Inches(1.0))
    p_c2_kw = tf1_c2.paragraphs[0]
    p_c2_kw.text = "광화문역 월세 절감 ~20%"
    p_c2_kw.font.name = FONT_TITLE
    p_c2_kw.font.size = Pt(16)
    p_c2_kw.font.bold = True
    p_c2_kw.font.color.rgb = COLOR_SECONDARY
    p_c2_desc = tf1_c2.add_paragraph()
    p_c2_desc.text = "B2B·전문직·공유 오피스 안정 수익 최적"
    p_c2_desc.font.name = FONT_BODY
    p_c2_desc.font.size = Pt(11)
    p_c2_desc.font.color.rgb = COLOR_ICE_BLUE
    p_c2_desc.space_before = Pt(4)

    # 카드 3
    draw_card(slide1, Inches(8.6), Inches(3.8), Inches(3.7), Inches(1.2), RGBColor(0x22, 0x2B, 0x68))
    tf1_c3 = create_textbox(slide1, Inches(8.75), Inches(3.92), Inches(3.4), Inches(1.0))
    p_c3_kw = tf1_c3.paragraphs[0]
    p_c3_kw.text = "초역세권 프리미엄 19.2%"
    p_c3_kw.font.name = FONT_TITLE
    p_c3_kw.font.size = Pt(16)
    p_c3_kw.font.bold = True
    p_c3_kw.font.color.rgb = COLOR_ICE_BLUE
    p_c3_desc = tf1_c3.add_paragraph()
    p_c3_desc.text = "고층+일반역세권 조합이 은닉된 비용 기회"
    p_c3_desc.font.name = FONT_BODY
    p_c3_desc.font.size = Pt(11)
    p_c3_desc.font.color.rgb = COLOR_ICE_BLUE
    p_c3_desc.space_before = Pt(4)

    p1_meta = create_textbox(slide1, Inches(1.0), Inches(5.3), Inches(11.0), Inches(0.4))
    p1_meta_p = p1_meta.paragraphs[0]
    p1_meta_p.text = "Antigravity AI Real Estate Data Pipeline  |  McKinsey Style Consultant  |  2026. 06"
    p1_meta_p.font.name = FONT_BODY
    p1_meta_p.font.size = Pt(10)
    p1_meta_p.font.color.rgb = COLOR_TEXT_MUTED

    
    # ---------------- SLIDE 2: 데이터 수집 개요 및 기초 통계 ----------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2, COLOR_BG_LIGHT)
    draw_organic_bg(slide2, is_dark_bg=False)
    draw_accent_dots(slide2, is_dark_bg=False)
    add_header(slide2, "1. 수집 데이터 개요 및 기초 통계")
    
    # 좌측 영역: 텍스트 정보 (Inches 0.8~6.3)
    tf2_left = create_textbox(slide2, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.8))
    p2_intro = tf2_left.paragraphs[0]
    p2_intro.text = "상업용 매물 데이터 분석 개요"
    p2_intro.font.name = FONT_TITLE
    p2_intro.font.size = Pt(18)
    p2_intro.font.bold = True
    p2_intro.font.color.rgb = COLOR_PRIMARY
    
    add_bullet_point(tf2_left, "수집 채널", "부동산 상업용 매물 전문 플랫폼 네모앱 (nemoapp.kr)")
    add_bullet_point(tf2_left, "유효 표본 수", f"총 {stats['total_count']:,}개 실시간 부동산 매물 (결측치 정제 완료)")
    add_bullet_point(tf2_left, "상권 구성", f"강남역 {stats['gn_count']:,}건 ({stats['gn_count']/stats['total_count']*100:.1f}%) | 광화문역 {stats['gw_count']:,}건 ({stats['gw_count']/stats['total_count']*100:.1f}%)")
    add_bullet_point(tf2_left, "보증금 평균", f"약 {format_money(stats['avg_deposit'])} 수준")
    add_bullet_point(tf2_left, "월세 평균", f"약 {format_money(stats['avg_rent'])} 수준")
    
    # 우측 영역: 대형 KPI 카드 (Inches 7.0~12.5)
    # KPI 1: 총 매물 수
    draw_card(slide2, Inches(7.0), Inches(1.7), Inches(5.5), Inches(1.4), COLOR_CARD_BG)
    tf2_card1 = create_textbox(slide2, Inches(7.3), Inches(1.85), Inches(4.9), Inches(1.1))
    tf2_card1.paragraphs[0].text = "TOTAL SAMPLE SIZE"
    tf2_card1.paragraphs[0].font.size = Pt(10)
    tf2_card1.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED
    p_num1 = tf2_card1.add_paragraph()
    p_num1.text = f"{stats['total_count']} 개 매물"
    p_num1.font.name = FONT_TITLE
    p_num1.font.size = Pt(28)
    p_num1.font.bold = True
    p_num1.font.color.rgb = COLOR_PRIMARY
    
    # KPI 2: 평균 보증금
    draw_card(slide2, Inches(7.0), Inches(3.3), Inches(5.5), Inches(1.4), COLOR_CARD_BG)
    tf2_card2 = create_textbox(slide2, Inches(7.3), Inches(3.45), Inches(4.9), Inches(1.1))
    tf2_card2.paragraphs[0].text = "OVERALL AVERAGE DEPOSIT"
    tf2_card2.paragraphs[0].font.size = Pt(10)
    tf2_card2.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED
    p_num2 = tf2_card2.add_paragraph()
    p_num2.text = format_money(stats['avg_deposit'])
    p_num2.font.name = FONT_TITLE
    p_num2.font.size = Pt(28)
    p_num2.font.bold = True
    p_num2.font.color.rgb = COLOR_ACCENT
    
    # KPI 3: 평균 월세
    draw_card(slide2, Inches(7.0), Inches(4.9), Inches(5.5), Inches(1.4), COLOR_CARD_BG)
    tf2_card3 = create_textbox(slide2, Inches(7.3), Inches(5.05), Inches(4.9), Inches(1.1))
    tf2_card3.paragraphs[0].text = "OVERALL AVERAGE MONTHLY RENT"
    tf2_card3.paragraphs[0].font.size = Pt(10)
    tf2_card3.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED
    p_num3 = tf2_card3.add_paragraph()
    p_num3.text = format_money(stats['avg_rent'])
    p_num3.font.name = FONT_TITLE
    p_num3.font.size = Pt(28)
    p_num3.font.bold = True
    p_num3.font.color.rgb = COLOR_SECONDARY
    
    draw_bottom_caption(slide2, is_dark_bg=False)

    # ===================================================================
    # [신규] SLIDE 3: Key Findings — 데이터 기반 핵심 발견 5가지
    # ===================================================================
    slide_kf = prs.slides.add_slide(slide_layout)
    set_slide_background(slide_kf, COLOR_BG_LIGHT)
    draw_organic_bg(slide_kf, is_dark_bg=False)
    draw_accent_dots(slide_kf, is_dark_bg=False)
    add_header(slide_kf, "KEY FINDINGS — 데이터가 말하는 5가지 필수 발견")

    # 5개 발견 사항 자료 정의
    rent_diff_pct = (stats['gn_avg_rent'] - stats['gw_avg_rent']) / stats['gw_avg_rent'] * 100
    transit_premium = 19.2
    findings = [
        (
            "01",
            f"강남역 월세 {rent_diff_pct:.1f}% 프리미엄 — 고마진 업종에만 정당성 존재",
            f"강남역 평균 월세 {format_money(stats['gn_avg_rent'])} vs 광화문역 {format_money(stats['gw_avg_rent'])}. "
            f"F&B·메디컴은 트래픽 기반 수익으로 회수 가능, 힘마진율 업종은 BEP +6개월."
        ),
        (
            "02",
            f"매물 분포는 근 수적 Right-Skewed — 중간값이 진짜 시세",
            f"전체 평균 보증금 {format_money(stats['avg_deposit'])} 대비 중간값은 훨씬 낮음. "
            f"예산 수립 시 평균이 아닌 중간값 기준으로 출발해야 과도 지출 리스크를 피할 수 있다."
        ),
        (
            "03",
            f"초역세권(5분 이내) 월세 프리미엄 {transit_premium:.1f}% — 워크인 필수 업종에만 유효",
            f"초역세권 평균 월세 5,006만 월 vs 일반역세권 4,201만 월. "
            f"B2B 서비스업·전문직은 5분 프리미엄 포기 + 고층 조합으로 20~30% 절감 가능."
        ),
        (
            "04",
            "6층 이상 고층은 최고 월세 수준 — 대형 사무실 계약의 었당놈는 수익맨",
            "6층 이상 평균 월세 6,414만 월로 전체 최고. "
            "그러나 제시 평평급 빌딩의 대형 전용 면적 폴가로 인한 것으로, "
            "㎡당 월세 단가는 1층 대비 15~20% 낙다. 대형 임차 업체에게 실질 가성비 기회."
        ),
        (
            "05",
            "TF-IDF 핵심: '층수·역세권·도보 분' — 임대인이 관리하는 시세 레버 3종셋",
            "매물 홍보의 최상위 셀링포인트는 층수, 역 5분·4분 등 교통 접근성, 일반음식점 등 업종이다. "
            "임차인은 이 3가지 레버를 활용해 협상 시 실질적 비용 절감을 요구해야 한다."
        ),
    ]

    # 2열 배치 (1~3번 좌측, 4~5번 우측)
    left_col_x = Inches(0.5)
    right_col_x = Inches(6.7)
    card_w = Inches(5.9)
    card_h = Inches(1.55)
    y_positions_left = [Inches(1.6), Inches(3.3), Inches(5.0)]
    y_positions_right = [Inches(1.6), Inches(3.3)]

    for i, (num, title, desc) in enumerate(findings):
        if i < 3:
            cx = left_col_x
            cy = y_positions_left[i]
        else:
            cx = right_col_x
            cy = y_positions_right[i - 3]

        # 카드 배경
        bg_color = COLOR_ACCENT_BG if i % 2 == 0 else COLOR_CARD_BG
        draw_card(slide_kf, cx, cy, card_w, card_h, bg_color)

        # 번호 표시
        tf_num = create_textbox(slide_kf, cx + Inches(0.12), cy + Inches(0.1), Inches(0.55), Inches(0.6))
        p_num = tf_num.paragraphs[0]
        p_num.text = num
        p_num.font.name = FONT_TITLE
        p_num.font.size = Pt(24)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_ACCENT if i % 2 == 0 else COLOR_SECONDARY

        # 타이틀
        tf_t = create_textbox(slide_kf, cx + Inches(0.75), cy + Inches(0.1), card_w - Inches(0.9), Inches(0.55))
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.name = FONT_BODY
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_PRIMARY

        # 설명
        tf_d = create_textbox(slide_kf, cx + Inches(0.75), cy + Inches(0.7), card_w - Inches(0.9), Inches(0.75))
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = COLOR_TEXT_DARK
        p_d.word_wrap = True

    draw_bottom_caption(slide_kf, is_dark_bg=False)

    # ---------------- SLIDE 4: 이상치 제거 전후 분포 비교 (신규) ----------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4, COLOR_BG_LIGHT)
    draw_organic_bg(slide4, is_dark_bg=False)
    draw_accent_dots(slide4, is_dark_bg=False)
    add_header(slide4, "3. 통계적 이상치(IQR 1.5) 식별 및 전후 분포 비교")
    
    # 좌측 영역: 인사이트 기술
    tf4_left = create_textbox(slide4, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.8))
    p4_title = tf4_left.paragraphs[0]
    p4_title.text = "초고가 빌딩으로 인한 시세 착시 현상 제거"
    p4_title.font.name = FONT_TITLE
    p4_title.font.size = Pt(18)
    p4_title.font.bold = True
    p4_title.font.color.rgb = COLOR_PRIMARY
    
    add_bullet_point(tf4_left, "이상치 영향도", "보증금 최대 67억 원, 월세 최대 6,700만 원 선의 극단치 빌딩이 시장 평균 시세를 착시시킴.")
    add_bullet_point(tf4_left, "IQR 1.5 기준 정제", "정제 후 보증금 4,000만 원, 월세 3,000만 원 선에서 지극히 현실적이고 밀집도 높은 분포가 드러남.")
    add_bullet_point(tf4_left, "의사결정 리스크", "평균값에 기반한 재무 계획 수립 시 발생할 예산 오차 방지 가능.")
    add_bullet_point(tf4_left, "시장 분석 강점", "중소형 상가 및 일반 오피스 매물의 정확한 거래 밴드 식별에 기여.")
    
    # 우측 영역: 이미지 배치
    img_path4 = "nemo_real_estate/images/12_outlier_removed.png"
    add_image_auto_ratio(slide4, img_path4, Inches(7.2), Inches(1.7), Inches(5.0))
    
    tf4_img_cap = create_textbox(slide4, Inches(7.2), Inches(5.7), Inches(5.0), Inches(0.5))
    tf4_img_cap.paragraphs[0].text = "[그림 2] 이상치 정제 전후 보증금 및 월세 분포 비교"
    tf4_img_cap.paragraphs[0].font.size = Pt(10)
    tf4_img_cap.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED
    tf4_img_cap.paragraphs[0].font.italic = True
    tf4_img_cap.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    draw_bottom_caption(slide4, is_dark_bg=False)
    
    # ---------------- SLIDE 5: 단위면적당 월세 단가 분석 (신규) ----------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5, COLOR_BG_LIGHT)
    draw_organic_bg(slide5, is_dark_bg=False)
    draw_accent_dots(slide5, is_dark_bg=False)
    add_header(slide5, "4. 상권별 단위면적(㎡)당 월세 단가 분포")
    
    # 좌측 영역: 이미지 배치
    img_path5 = "nemo_real_estate/images/13_price_per_sqm.png"
    add_image_auto_ratio(slide5, img_path5, Inches(0.8), Inches(1.7), Inches(5.0))
    
    tf5_img_cap = create_textbox(slide5, Inches(0.8), Inches(5.7), Inches(5.0), Inches(0.5))
    tf5_img_cap.paragraphs[0].text = "[그림 3] 상권별 단위면적당 월세 단가 분포 비교"
    tf5_img_cap.paragraphs[0].font.size = Pt(10)
    tf5_img_cap.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED
    tf5_img_cap.paragraphs[0].font.italic = True
    tf5_img_cap.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 우측 영역: 분석 인사이트 기술
    tf5_right = create_textbox(slide5, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.8))
    p5_title = tf5_right.paragraphs[0]
    p5_title.text = "단위면적 기준 가격의 의외의 역전 현상"
    p5_title.font.name = FONT_TITLE
    p5_title.font.size = Pt(18)
    p5_title.font.bold = True
    p5_title.font.color.rgb = COLOR_PRIMARY
    
    add_bullet_point(tf5_right, "면적당 월세 역전", "㎡당 평균 월세는 광화문역(42.5만 원)이 강남역(38.3만 원) 대비 11% 더 비싼 것으로 밝혀짐.")
    add_bullet_point(tf5_right, "자산 크기 효과 착시", "강남역은 넓은 대형 매물 공급이 많아 임대료 총액이 비싸 보였던 착시 효과가 존재했음.")
    add_bullet_point(tf5_right, "공간 효율성", "중소형 콤팩트 오피스 임차 시 광화문의 실질 단가 부담이 강남보다 클 수 있음.")
    add_bullet_point(tf5_right, "협상 기준 설정", "총액 예산뿐 아니라 ㎡당 실질 단가를 활용한 임대차 조율 전략 수립 필수.")
    
    draw_bottom_caption(slide5, is_dark_bg=False)
    
    # ---------------- SLIDE 6: 지역별 4대 핵심 지표 비교 (신규) ----------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6, COLOR_BG_LIGHT)
    draw_organic_bg(slide6, is_dark_bg=False)
    draw_accent_dots(slide6, is_dark_bg=False)
    add_header(slide6, "5. 상권별 4대 핵심 임대 지표 종합 비교 (이상치 정제)")
    
    # 좌측 영역: 인사이트 기술
    tf6_left = create_textbox(slide6, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.8))
    p6_title = tf6_left.paragraphs[0]
    p6_title.text = "보증금·월세·면적·관리비 종합 비교"
    p6_title.font.name = FONT_TITLE
    p6_title.font.size = Pt(18)
    p6_title.font.bold = True
    p6_title.font.color.rgb = COLOR_PRIMARY
    
    add_bullet_point(tf6_left, "월세 및 관리비 편차", "강남역(월세 중간값 3,200만 원, 관리비 300만 원)이 광화문(2,725만 원, 200만 원)보다 고정비 부담이 확실히 큼.")
    add_bullet_point(tf6_left, "공간 규모 옵션", "전용면적 평균은 강남역이 112㎡로 광화문역(81.8㎡) 대비 38% 넓어, 대형 사무실 계약 시 다양한 매물 확보 가능.")
    add_bullet_point(tf6_left, "안정 지향 기업", "광화문역 상권은 4대 지표의 변동성이 좁고 가격 분포가 고른 안정적인 운영에 적합함.")
    add_bullet_point(tf6_left, "비용 통제 레버", "관리비 비율이 높은 강남 상권 계약 시 특약 조건을 통한 고정비 방어 중요.")
    
    # 우측 영역: 이미지 배치
    img_path6 = "nemo_real_estate/images/14_region_comparison.png"
    add_image_auto_ratio(slide6, img_path6, Inches(7.2), Inches(1.7), Inches(5.0))
    
    tf6_img_cap = create_textbox(slide6, Inches(7.2), Inches(5.7), Inches(5.0), Inches(0.5))
    tf6_img_cap.paragraphs[0].text = "[그림 4] 지역별 4대 지표 상자그림 분포 비교"
    tf6_img_cap.paragraphs[0].font.size = Pt(10)
    tf6_img_cap.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED
    tf6_img_cap.paragraphs[0].font.italic = True
    tf6_img_cap.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    draw_bottom_caption(slide6, is_dark_bg=False)
    
    # ---------------- SLIDE 7: 보증금 대비 월세 상관관계 버블 분석 (신규) ----------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7, COLOR_BG_LIGHT)
    draw_organic_bg(slide7, is_dark_bg=False)
    draw_accent_dots(slide7, is_dark_bg=False)
    add_header(slide7, "6. 보증금-월세-면적 관계의 입체적 버블 분석")
    
    # 좌측 영역: 이미지 배치
    img_path7 = "nemo_real_estate/images/15_deposit_rent_bubble.png"
    add_image_auto_ratio(slide7, img_path7, Inches(0.8), Inches(1.7), Inches(5.0))
    
    tf7_img_cap = create_textbox(slide7, Inches(0.8), Inches(5.7), Inches(5.0), Inches(0.5))
    tf7_img_cap.paragraphs[0].text = "[그림 5] 보증금-월세-면적 결합 버블 분석"
    tf7_img_cap.paragraphs[0].font.size = Pt(10)
    tf7_img_cap.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED
    tf7_img_cap.paragraphs[0].font.italic = True
    tf7_img_cap.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 우측 영역: 분석 인사이트 기술
    tf7_right = create_textbox(slide7, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.8))
    p7_title = tf7_right.paragraphs[0]
    p7_title.text = "상호보완이 아닌 동반 상승하는 임대 시세"
    p7_title.font.name = FONT_TITLE
    p7_title.font.size = Pt(18)
    p7_title.font.bold = True
    p7_title.font.color.rgb = COLOR_PRIMARY
    
    add_bullet_point(tf7_right, "양(+)의 상관관계", "보증금과 월세 간 상관계수가 강남 0.87, 광화문 0.75로 전형적인 우상향 자산 가치 동반 상승을 보여줌.")
    add_bullet_point(tf7_right, "자산 규모 효과", "보증금을 증액해 월세를 삭감하는 주거용 오피스텔 성격과 달리, 공간 가치가 커질수록 두 항목이 동시에 비례하여 치솟음.")
    add_bullet_point(tf7_right, "버블 크기의 시사점", "면적(버블 크기)이 대형화될수록 강남역 상권은 가파른 경사로 월세와 보증금이 수직 상승하는 고위험 예산 구간 도달.")
    add_bullet_point(tf7_right, "자산 포트폴리오", "통합 예산 운용 시 두 변수를 독립적으로 설계하지 말고 총 보유 현금흐름으로 접근해야 함.")
    
    draw_bottom_caption(slide7, is_dark_bg=False)
    
    # ---------------- SLIDE 8: 지역별 권리금 및 비매몰 초기비용 분석 (신규) ----------------
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide8, COLOR_BG_LIGHT)
    draw_organic_bg(slide8, is_dark_bg=False)
    draw_accent_dots(slide8, is_dark_bg=False)
    add_header(slide8, "7. 상권별 권리금 분포 및 유권리 매물 비율")
    
    # 좌측 영역: 인사이트 기술
    tf8_left = create_textbox(slide8, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.8))
    p8_title = tf8_left.paragraphs[0]
    p8_title.text = "권리금 유무에 따른 초기 매몰 비용 진단"
    p8_title.font.name = FONT_TITLE
    p8_title.font.size = Pt(18)
    p8_title.font.bold = True
    p8_title.font.color.rgb = COLOR_PRIMARY
    
    add_bullet_point(tf8_left, "유권리 비중 48.7%", "전체 등록 매물의 절반에 가까운 상가들이 권리금을 요구하여 창업 시 초기 예산 압박의 주된 원인이 됨.")
    add_bullet_point(tf8_left, "광화문의 권리금 역설", "평균 권리금은 광화문(9,518만 원)이 강남(7,014만 원)보다 약 35.7% 높아 실질 영업 기반 인수에 장벽이 큼.")
    add_bullet_point(tf8_left, "상권 배후 안정성", "광화문은 직장인/공공기관의 단단한 주중 고정 수요로 기존 상가의 시설·영업 가치를 더 높게 평가받는 경향이 있음.")
    add_bullet_point(tf8_left, "인수 의사결정", "권리금 회수 가능성을 정밀하게 검증하여 매몰 리스크가 큰 플래그십 업종의 광화문 진입 시 신중해야 함.")
    
    # 우측 영역: 이미지 배치
    img_path8 = "nemo_real_estate/images/16_premium_analysis.png"
    add_image_auto_ratio(slide8, img_path8, Inches(7.2), Inches(1.7), Inches(5.0))
    
    tf8_img_cap = create_textbox(slide8, Inches(7.2), Inches(5.7), Inches(5.0), Inches(0.5))
    tf8_img_cap.paragraphs[0].text = "[그림 6] 상권별 권리금 평균값 및 유권리 매물 빈도"
    tf8_img_cap.paragraphs[0].font.size = Pt(10)
    tf8_img_cap.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED
    tf8_img_cap.paragraphs[0].font.italic = True
    tf8_img_cap.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    draw_bottom_caption(slide8, is_dark_bg=False)
    
    # ---------------- SLIDE 9: 매물 상세 설명 텍스트 마이닝 분석 ----------------
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide9, COLOR_BG_LIGHT)
    draw_organic_bg(slide9, is_dark_bg=False)
    draw_accent_dots(slide9, is_dark_bg=False)
    add_header(slide9, "8. 매물 상세 설명 텍스트 마이닝 분석")
    
    # 좌측 영역: 이미지 배치
    img_path9 = "nemo_real_estate/images/11_tfidf_keywords_bar.png"
    add_image_auto_ratio(slide9, img_path9, Inches(0.8), Inches(1.7), Inches(5.0))
    
    tf9_img_cap = create_textbox(slide9, Inches(0.8), Inches(5.7), Inches(5.0), Inches(0.5))
    tf9_img_cap.paragraphs[0].text = "[그림 7] TF-IDF 중요 키워드 상위 15개 핵심 가중치 시각화"
    tf9_img_cap.paragraphs[0].font.size = Pt(10)
    tf9_img_cap.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED
    tf9_img_cap.paragraphs[0].font.italic = True
    tf9_img_cap.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 우측 영역: 분석 인사이트 기술
    tf9_right = create_textbox(slide9, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.8))
    p9_title = tf9_right.paragraphs[0]
    p9_title.text = "가중치 키워드를 통한 중개 실무 해석"
    p9_title.font.name = FONT_TITLE
    p9_title.font.size = Pt(18)
    p9_title.font.bold = True
    p9_title.font.color.rgb = COLOR_PRIMARY
    
    add_bullet_point(tf9_right, "층수 가치 최우선", "최상위 키워드 '1층', '2층', '층수'의 강세는 가시성 및 도보 접근성이 상업적 집객의 핵심 척도임을 입증.")
    add_bullet_point(tf9_right, "도보 역세권 입지", "'5분', '4분', '강남역', '역삼역', '종각역' 등 역사 접근거리가 임대인들의 가장 주된 셀링 포인트로 활용됨.")
    add_bullet_point(tf9_right, "권리금 및 관리비", "'000만원', '500만원' 등 권리금 관련 키워드의 빈번한 노출은 초기 창업 비용에 미치는 영향력을 방증.")
    add_bullet_point(tf9_right, "업종 및 용도 매칭", "'일반음식점', '기타업종', '휴게음식점'의 강세로 식음료(F&B) 업계 중심의 임차 수요 활성도가 드러남.")
    
    draw_bottom_caption(slide9, is_dark_bg=False)
    
    # ---------------- SLIDE 10: 종합 결론 및 입지 제언 ----------------
    slide10 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide10, COLOR_BG_LIGHT)
    draw_organic_bg(slide10, is_dark_bg=False)
    draw_accent_dots(slide10, is_dark_bg=False)
    add_header(slide10, "9. 상권별 맞춤형 입지 추천 가이드라인")
    
    # 2열 배치 카드 (Inches 0.8~6.3 / Inches 7.0~12.5)
    # 1. 강남역 추천 입지 전략
    draw_card(slide10, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.6), COLOR_ACCENT_BG)
    tf10_left = create_textbox(slide10, Inches(1.1), Inches(1.9), Inches(4.9), Inches(4.2))
    p10_l_title = tf10_left.paragraphs[0]
    p10_l_title.text = "■ 강남역 상권 (고집객 - 고위험 전략)"
    p10_l_title.font.name = FONT_TITLE
    p10_l_title.font.size = Pt(17)
    p10_l_title.font.bold = True
    p10_l_title.font.color.rgb = COLOR_ACCENT
    
    add_bullet_point(tf10_left, "타깃 업종", "대형 프랜차이즈, F&B, 플래그십 샵, 피부과/안과 의원 등", size_pt=12)
    add_bullet_point(tf10_left, "핵심 메리트", "막강한 트래픽 기반 높은 집객 시너지와 젊은 층 수요 독식 구조", size_pt=12)
    add_bullet_point(tf10_left, "재무적 시사점", f"높은 월세 수준(평균 {format_money(stats['gn_avg_rent'])})을 버텨낼 초기 현금 체력 및 높은 마진율 확보가 최우선 조건임", size_pt=12)
    
    # 2. 광화문역 추천 입지 전략
    draw_card(slide10, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.6), COLOR_CARD_BG)
    tf10_right = create_textbox(slide10, Inches(7.3), Inches(1.9), Inches(4.9), Inches(4.2))
    p10_r_title = tf10_right.paragraphs[0]
    p10_r_title.text = "■ 광화문역 상권 (안정성 - 내실화 전략)"
    p10_r_title.font.name = FONT_TITLE
    p10_r_title.font.size = Pt(17)
    p10_r_title.font.bold = True
    p10_r_title.font.color.rgb = COLOR_PRIMARY
    
    add_bullet_point(tf10_right, "타깃 업종", "한식/일식 등 직장인 타깃 중고가 요식업, 전문직 사무소, 공유 오피스 등", size_pt=12)
    add_bullet_point(tf10_right, "핵심 메리트", "정부청사/대기업 본사 중심의 안정적인 주중 상주인구와 구매력 높은 오피스 배후", size_pt=12)
    add_bullet_point(tf10_right, "재무적 시사점", f"강남 상권 대비 월평균 {100 - (stats['gw_avg_rent']/stats['gn_avg_rent']*100):.1f}%의 임대 관리비 절약이 가능하여, 리스크가 적고 효율 지향적인 중장기 비즈니스에 강추", size_pt=12)
    
    draw_bottom_caption(slide10, is_dark_bg=False)
    
    # ---------------- SLIDE 11: 감사합니다 (Closing Slide - Dark Theme) ----------------
    slide11 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide11, COLOR_BG_DARK)
    draw_organic_bg(slide11, is_dark_bg=True)
    draw_accent_dots(slide11, is_dark_bg=True)
    
    tf11 = create_textbox(slide11, Inches(1.0), Inches(2.6), Inches(11.33), Inches(3.0))
    p11 = tf11.paragraphs[0]
    p11.alignment = PP_ALIGN.CENTER
    p11.text = "경청해 주셔서 감사합니다"
    p11.font.name = FONT_TITLE
    p11.font.size = Pt(38)
    p11.font.bold = True
    p11.font.color.rgb = COLOR_TEXT_LIGHT
    
    p11_sub = tf11.add_paragraph()
    p11_sub.alignment = PP_ALIGN.CENTER
    p11_sub.text = "Q & A"
    p11_sub.font.name = FONT_BODY
    p11_sub.font.size = Pt(18)
    p11_sub.font.color.rgb = COLOR_ICE_BLUE
    p11_sub.space_before = Pt(18)
    
    # 저장
    prs.save(output_pptx_path)
    print(f"[Pptx Generator] Nordic Minimalism 테마 발표 슬라이드 생성이 완료되었습니다: {output_pptx_path}")

if __name__ == "__main__":
    build_presentation()
