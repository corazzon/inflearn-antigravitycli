"""
Yes24 IT/컴퓨터 베스트셀러 EDA PPTX 보고서 자동 생성 스크립트

이 스크립트는 yes24/docs/eda_report.md의 34KB 분량 데이터 분석 리포트를 바탕으로,
총 30페이지 분량의 전문적인 PPTX 슬라이드를 생성합니다.
노르딕 미니멀리즘 스타일(배경: 웜크림 #F4F1EC, 포인트: 벽돌색 #B85042)을 따르며,
제목 폰트(Gmarket Sans Bold) 및 본문 폰트(나눔고딕)를 적용합니다.
각 슬라이드별로 다양한 인포그래픽 레이아웃(카드형, 분할형, 데이터 강조형) 및
Icons8에서 다운로드한 플랫 아이콘, 실제 EDA 분석 그래프 이미지들을 포함하고,
슬라이드 하단에 2분 분량의 상세 발표자 노트(Speaker Notes)를 포함하여 완성도 높은 산출물을 제작합니다.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # 프레젠테이션 초기화 및 16:9 슬라이드 크기 설정
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 디자인 시스템 상수 정의 (노르딕 미니멀리즘 테마)
    COLOR_BG = RGBColor(244, 241, 236)         # 웜 크림
    COLOR_BRICK = RGBColor(184, 80, 66)        # 포인트 벽돌색
    COLOR_PRIMARY_TEXT = RGBColor(61, 53, 48)  # 어두운 웜브라운 (나머지 텍스트)
    COLOR_SECONDARY_TEXT = RGBColor(138, 122, 106) # 타우페 (보조 텍스트)
    COLOR_CARD_BG = RGBColor(255, 255, 255)    # 카드 배경 (순백색)

    FONT_TITLE = "Gmarket Sans Bold"
    FONT_BODY = "NanumGothic"

    IMAGE_DIR = "yes24/images"
    ICON_DIR = "yes24/images/icons"

    # 공통 헬퍼 함수들
    def apply_bg(slide):
        """모든 슬라이드에 웜크림색 배경 사각형을 깔고, 노르딕 3-dot 시그니처를 추가합니다."""
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        
        # 3-dot 시그니처 (좌측 상단 여백에 배치)
        dot_y = Inches(0.4)
        dot_size = Inches(0.12)
        colors = [COLOR_BRICK, COLOR_SECONDARY_TEXT, COLOR_PRIMARY_TEXT]
        for i, color in enumerate(colors):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6 + i * 0.2), dot_y, dot_size, dot_size)
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

    def add_header(slide, title_text, category_text="YES24 IT/COMPUTER BESTSELLER EDA"):
        """상단 카테고리와 메인 타이틀을 추가합니다."""
        cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12.0), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_bottom = tf_cat.margin_right = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = FONT_BODY
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_SECONDARY_TEXT
        
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12.0), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_bottom = tf_title.margin_right = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY_TEXT

    def add_footer(slide, page_num):
        """하단 얇은 구분선과 페이지 번호를 추가합니다."""
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.8), Inches(12.133), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_SECONDARY_TEXT
        line.line.fill.background()
        
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(10.0), Inches(0.3))
        tf = footer_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = "Yes24 IT/컴퓨터 베스트셀러 데이터 심층 EDA 리포트  |  수석 데이터 분석가 제언"
        p.font.name = FONT_BODY
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_SECONDARY_TEXT
        
        page_box = slide.shapes.add_textbox(Inches(11.733), Inches(6.9), Inches(1.0), Inches(0.3))
        tf_page = page_box.text_frame
        tf_page.word_wrap = True
        tf_page.margin_left = tf_page.margin_top = tf_page.margin_bottom = tf_page.margin_right = 0
        p_page = tf_page.paragraphs[0]
        p_page.text = f"{page_num} / 30"
        p_page.alignment = PP_ALIGN.RIGHT
        p_page.font.name = FONT_BODY
        p_page.font.size = Pt(9)
        p_page.font.bold = True
        p_page.font.color.rgb = COLOR_BRICK

    def add_icon(slide, icon_name, x, y, size=Inches(0.5)):
        """지정된 위치에 다운로드한 벽돌색 플랫 아이콘을 추가합니다."""
        icon_path = os.path.join(ICON_DIR, f"{icon_name}.png")
        if os.path.exists(icon_path):
            slide.shapes.add_picture(icon_path, x, y, width=size, height=size)

    def add_card(slide, x, y, w, h, fill_color=COLOR_CARD_BG):
        """컨텐츠 영역 구분을 위해 옅은 테두리가 있는 순백색 카드 도형을 생성합니다."""
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        card.line.color.rgb = RGBColor(220, 215, 205)
        card.line.width = Pt(1)
        return card

    def set_speaker_notes(slide, notes_text):
        """슬라이드에 2분 분량의 발표자 스크립트를 삽입합니다."""
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

    def add_standard_chart_slide(page_num, title, chart_filename, stat_title, stats_list, insight_text, icon_name="chart"):
        """차트 이미지와 설명 카드가 들어간 표준 슬라이드를 생성합니다."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        apply_bg(slide)
        add_header(slide, title)
        
        # 좌측 차트 이미지 배치
        chart_path = os.path.join(IMAGE_DIR, chart_filename)
        if os.path.exists(chart_path):
            slide.shapes.add_picture(chart_path, Inches(0.6), Inches(1.8), Inches(6.0), Inches(4.5))
            
        # 우측 설명 카드 배치
        card_x, card_y = Inches(6.9), Inches(1.8)
        card_w, card_h = Inches(5.8), Inches(4.7)
        add_card(slide, card_x, card_y, card_w, card_h)
        
        # 카드 상단 포인트 바
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card_x, card_y, card_w, Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_BRICK
        bar.line.fill.background()
        
        # [아이콘 및 타이틀 가로 배치]
        icon_size = Inches(0.45)
        # 아이콘 배치 (카드 타이틀 바로 좌측에 정렬하여 텍스트 겹침 차단)
        add_icon(slide, icon_name, card_x + Inches(0.4), card_y + Inches(0.35), icon_size)
        
        # 타이틀용 텍스트박스 추가
        title_tb = slide.shapes.add_textbox(card_x + Inches(0.95), card_y + Inches(0.3), card_w - Inches(1.35), Inches(0.6))
        tf_t = title_tb.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_bottom = tf_t.margin_right = 0
        p_stat_title = tf_t.paragraphs[0]
        p_stat_title.text = stat_title
        p_stat_title.font.name = FONT_TITLE
        p_stat_title.font.size = Pt(18)
        p_stat_title.font.bold = True
        p_stat_title.font.color.rgb = COLOR_PRIMARY_TEXT
        
        # [본문 텍스트박스 하단 분리 배치]
        tb = slide.shapes.add_textbox(card_x + Inches(0.4), card_y + Inches(1.05), card_w - Inches(0.8), card_h - Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        # 통계 지표 리스트
        first = True
        for label, val in stats_list:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = f"• {label}: "
            p.font.name = FONT_BODY
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = COLOR_PRIMARY_TEXT
            p.space_after = Pt(4)
            
            run = p.add_run()
            run.text = val
            run.font.name = FONT_BODY
            run.font.size = Pt(12)
            run.font.bold = False
            run.font.color.rgb = COLOR_SECONDARY_TEXT
            
        # 비즈니스 해석 타이틀
        p_ins_title = tf.add_paragraph()
        p_ins_title.text = "[비즈니스 해석]"
        p_ins_title.font.name = FONT_BODY
        p_ins_title.font.size = Pt(13)
        p_ins_title.font.bold = True
        p_ins_title.font.color.rgb = COLOR_BRICK
        p_ins_title.space_before = Pt(10)
        p_ins_title.space_after = Pt(4)
        
        # 해석 상세
        p_ins = tf.add_paragraph()
        p_ins.text = insight_text
        p_ins.font.name = FONT_BODY
        p_ins.font.size = Pt(11)
        p_ins.font.color.rgb = COLOR_PRIMARY_TEXT
        p_ins.line_spacing = 1.3
        
        add_footer(slide, page_num)
        return slide

    def add_deep_insight_slide(page_num, title, left_title, left_points, right_title, right_points):
        """그래프가 없는 심층분석 설명 전용 2열 카드 슬라이드를 생성합니다."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        apply_bg(slide)
        add_header(slide, title)
        
        # 좌우 카드 배치
        for i, (c_title, c_points) in enumerate([(left_title, left_points), (right_title, right_points)]):
            x = Inches(0.6 + i * 6.1)
            y = Inches(1.8)
            w = Inches(5.8)
            h = Inches(4.7)
            add_card(slide, x, y, w, h)
            
            # 상단 포인트 바
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.12))
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLOR_PRIMARY_TEXT if i == 1 else COLOR_BRICK
            bar.line.fill.background()
            
            # 텍스트 상자 (너비를 늘려 부자연스러운 개행 방지: Inches(0.4) -> Inches(0.3))
            tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), h - Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
            
            p_title = tf.paragraphs[0]
            p_title.text = c_title
            p_title.font.name = FONT_TITLE
            p_title.font.size = Pt(18)
            p_title.font.bold = True
            p_title.font.color.rgb = COLOR_PRIMARY_TEXT
            p_title.space_after = Pt(12)  # 타이틀 아래 여백 확보
            
            for p_title_text, p_body in c_points:
                p = tf.add_paragraph()
                p.text = f"• {p_title_text}"
                p.font.name = FONT_BODY
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = COLOR_PRIMARY_TEXT
                p.space_after = Pt(4)
                
                p2 = tf.add_paragraph()
                p2.text = p_body
                p2.font.name = FONT_BODY
                p2.font.size = Pt(11)
                p2.font.color.rgb = COLOR_SECONDARY_TEXT
                p2.line_spacing = 1.35
                p2.space_after = Pt(14)  # 수직 여백 고루 배치하기 위해 단락 여백을 넓혀줌
                
        add_footer(slide, page_num)
        return slide

    def add_section_intro(page_num, part_name, title_text, desc_text):
        """섹션 시작을 알리는 다크 슬라이드를 생성합니다."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # 어두운 브라운 배경
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_PRIMARY_TEXT
        bg.line.fill.background()
        
        # 벽돌색 블록 데코
        block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.2), Inches(1.5), Inches(0.12))
        block.fill.solid()
        block.fill.fore_color.rgb = COLOR_BRICK
        block.line.fill.background()
        
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(3.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        p1 = tf.paragraphs[0]
        p1.text = part_name.upper()
        p1.font.name = FONT_TITLE
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_BRICK
        
        p2 = tf.add_paragraph()
        p2.text = title_text
        p2.font.name = FONT_TITLE
        p2.font.size = Pt(36)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_BG
        
        p3 = tf.add_paragraph()
        p3.text = f"\n{desc_text}"
        p3.font.name = FONT_BODY
        p3.font.size = Pt(14)
        p3.font.color.rgb = COLOR_SECONDARY_TEXT
        p3.line_spacing = 1.3
        
        # 하단 라인 및 페이지
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(6.8), Inches(10.333), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_SECONDARY_TEXT
        line.line.fill.background()
        
        page_box = slide.shapes.add_textbox(Inches(10.833), Inches(6.9), Inches(1.0), Inches(0.3))
        tf_page = page_box.text_frame
        p_page = tf_page.paragraphs[0]
        p_page.text = f"{page_num} / 30"
        p_page.alignment = PP_ALIGN.RIGHT
        p_page.font.name = FONT_BODY
        p_page.font.size = Pt(9)
        p_page.font.bold = True
        p_page.font.color.rgb = COLOR_BRICK
        
        return slide

    # ==============================================================================
    # 30개 슬라이드 한장 한장 상세 생성 시작
    # ==============================================================================

    slide_layout = prs.slide_layouts[6]

    # --- Slide 1: 표지 ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.8), Inches(0.15), Inches(3.8))
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_BRICK
    rect.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.5), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Yes24 IT/컴퓨터 베스트셀러\n데이터 심층 EDA 보고서"
    p.font.name = FONT_TITLE
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_TEXT
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(11.5), Inches(1.0))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "데이터 기반 기술 실용서 시장 분석 및 비즈니스 마케팅 전략 제언"
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = COLOR_SECONDARY_TEXT
    info_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(11.5), Inches(1.0))
    tf_info = info_box.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "작성일: 2026년 06월 06일  |  작성자: 20년차 수석 데이터 분석가  |  대상: 베스트셀러 도서 1,000건"
    p_info.font.name = FONT_BODY
    p_info.font.size = Pt(13)
    p_info.font.bold = True
    p_info.font.color.rgb = COLOR_PRIMARY_TEXT
    set_speaker_notes(slide, 
        "안녕하십니까. 오늘 발표를 맡은 수석 데이터 분석가입니다. "
        "우리는 Yes24 IT/컴퓨터 분야 베스트셀러 도서 1,000건의 로우 데이터를 정밀하게 정제하고 탐색적 데이터 분석(EDA)을 수행했습니다. "
        "이를 통해 단순히 어떤 책이 잘 팔리는가를 넘어, IT 도서 시장만의 특이한 소비자 구매 패턴과 경제적 가격 장벽, "
        "그리고 매출 극대화를 가져올 수 있는 실질적인 비즈니스 마케팅 전략 제언을 종합적으로 도출했습니다. "
        "본 발표 자료는 노르딕 미니멀리즘 스타일로 구성되어 시각적 이해를 돕기 위한 다양한 인포그래픽 기법이 반영되었습니다. "
        "그럼 지금부터 자세한 분석 내용을 순차적으로 보고해 드리겠습니다."
    )

    # --- Slide 2: 목차 ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "목차 및 분석 로드맵")
    toc_titles = ["1부. 단변량 분석 및 기초 데이터 탐색", "2부. 이변량 & 다변량 연계 분석", "3부. 종합 비즈니스 제언"]
    toc_descs = [
        "판매가 분포, 도서 평점 만족도, 판매지수 형태, 분철 서비스 여부, 메이저 출판사 점유율 등 주요 단일 변수의 기초 통계와 시각화 데이터를 분석합니다.",
        "할인율 및 분철 서비스 제공 유무에 따른 판매 성과, 최근 출판 트렌드 시계열 분석, 다변량 상관관계 분석, TF-IDF 텍스트 마이닝을 통한 시장 수요 예측을 다룹니다.",
        "20년차 데이터 분석가의 관점으로 가격 포지셔닝 이원화 전략, 린(Lean) 퍼블리싱 모델, 소셜 증명 및 리뷰 평판 관리, 태깅 최적화 등의 액션 플랜을 제시합니다."
    ]
    toc_icons = ["book", "chart", "lightbulb"]
    for i in range(3):
        x = Inches(0.6 + i * 4.1)
        y = Inches(1.8)
        w = Inches(3.9)
        h = Inches(4.7)
        add_card(slide, x, y, w, h)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_BRICK if i == 0 else COLOR_SECONDARY_TEXT if i == 1 else COLOR_PRIMARY_TEXT
        bar.line.fill.background()
        add_icon(slide, toc_icons[i], x + Inches(0.4), y + Inches(0.4), Inches(0.6))
        t_box = slide.shapes.add_textbox(x + Inches(0.4), y + Inches(1.2), w - Inches(0.8), Inches(0.6))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = toc_titles[i]
        p.font.name = FONT_TITLE
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        d_box = slide.shapes.add_textbox(x + Inches(0.4), y + Inches(2.0), w - Inches(0.8), Inches(2.4))
        tf_d = d_box.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = toc_descs[i]
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_PRIMARY_TEXT
        p_d.line_spacing = 1.3
    add_footer(slide, 2)
    set_speaker_notes(slide,
        "이번 페이지는 목차 및 분석 로드맵 장표입니다. "
        "본 프레젠테이션은 총 30페이지 분량으로 유기적으로 구성되어 있으며, 세부적으로 3단계의 흐름을 밟게 됩니다. "
        "1부에서는 개별 변수들의 기초 분포를 탐색하고, 2부에서는 본격적인 상관분석과 이변량/다변량 차트, 그리고 TF-IDF 마이닝을 봅니다. "
        "3부에서는 앞선 정량적 데이터를 가공하여 실제 출판 경영에 도입 가능한 5대 종합 비즈니스 액션 플랜을 도출하겠습니다."
    )

    # --- Slide 3: 데이터 소개 및 수집 개요 ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "데이터 세트 수집 개요 및 스키마 구조")
    add_card(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.7))
    left_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.0), Inches(4.1))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "데이터 세트 기본 스펙"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_BRICK
    specs = [
        ("수집 대상", "Yes24 IT/컴퓨터 카테고리 베스트셀러"),
        ("데이터 크기", "1,000개 행(Rows) / 30개 열(Columns)"),
        ("중복 데이터", "0건 (완전 제거 완료)"),
        ("정제 변수", "clean 처리된 수치형 및 범주형 재가공 변수 다수"),
        ("분석 일자", "2026년 06월 05일 기준 최신 트렌드 반영")
    ]
    for label, val in specs:
        p = tf.add_paragraph()
        p.text = f"\n• {label}:  "
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        run = p.add_run()
        run.text = val
        run.font.name = FONT_BODY
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = COLOR_PRIMARY_TEXT
    add_card(slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.7))
    right_box = slide.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(5.0), Inches(4.1))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "주요 핵심 변수 구분"
    p_r.font.name = FONT_TITLE
    p_r.font.size = Pt(20)
    p_r.font.bold = True
    p_r.font.color.rgb = COLOR_PRIMARY_TEXT
    vars_desc = [
        ("대상 변수", "도서명, 저자, 출판사, 발행 연/월, 태그"),
        ("성과 지표", "판매지수(Sale Index), 평점(Rating), 리뷰 수"),
        ("기능/가격", "판매가, 정가, 할인율, 분철 여부(Spring Service)"),
        ("파생 변수", "정제된 가격(clean), 태그 개수, 할인율 범주")
    ]
    for label, val in vars_desc:
        p = tf_r.add_paragraph()
        p.text = f"\n• {label}:  "
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        run = p.add_run()
        run.text = val
        run.font.name = FONT_BODY
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = COLOR_SECONDARY_TEXT
    add_footer(slide, 3)
    set_speaker_notes(slide,
        "이번 페이지는 분석에 사용된 원시 데이터의 스펙과 변수 구성을 조망하는 장표입니다. "
        "Yes24의 베스트셀러 리스트에서 IT/컴퓨터 도서 1,000건을 온전히 추출했습니다. "
        "도서명과 저자 같은 메타데이터는 물론이고, 가격과 할인율, 그리고 핵심 성과 지표인 판매지수와 리뷰 수, 평점 등이 포함되어 있습니다. "
        "이를 깨끗이 정제하여 분석 모델링에 탑재했습니다."
    )

    # --- Slide 4: 요약 통계 (수치형 데이터) ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "전체 수치형 변수 통합 기술통계 분석")
    kpis = [
        ("평균 판매가", "23,480 원", "중위값: 22,500원 / 최소: 4,500원", "price"),
        ("평균 판매지수", "3,026 점", "중위값: 1,236점 / 최대: 87,480점", "trend"),
        ("평균 도서 평점", "7.50 점", "중위값: 9.70점 (0점 미평가 제외 시)", "star"),
        ("평균 도서 할인율", "8.64 %", "최빈값: 10% (도서정가제 상한 도달)", "percent")
    ]
    for idx, (label, val, subtext, icon) in enumerate(kpis):
        col = idx % 2
        row = idx // 2
        x = Inches(0.6 + col * 6.1)
        y = Inches(1.8 + row * 2.4)
        w = Inches(5.8)
        h = Inches(2.2)
        add_card(slide, x, y, w, h)
        add_icon(slide, icon, x + Inches(0.4), y + Inches(0.4), Inches(0.6))
        tb = slide.shapes.add_textbox(x + Inches(1.2), y + Inches(0.3), w - Inches(1.4), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.name = FONT_BODY
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_SECONDARY_TEXT
        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.name = FONT_TITLE
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_BRICK
        p3 = tf.add_paragraph()
        p3.text = subtext
        p3.font.name = FONT_BODY
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_PRIMARY_TEXT
    add_footer(slide, 4)
    set_speaker_notes(slide,
        "전체 수치형 변수의 평균과 중위수 분포를 일목요연하게 정리한 KPI 슬라이드입니다. "
        "평균 판매가는 23,480원이지만, 판매지수의 중위값은 1,236점으로 평균에 비해 상당히 낮습니다. "
        "이는 최상위 일부 베스트셀러가 평균 지수를 크게 끌어올린 우편향 구조를 나타냅니다. "
        "평점 또한 평가되지 않은 0점을 배제하면 실제 만족도 중위수가 9.7점에 수렴합니다."
    )

    # --- Slide 5: 요약 통계 (범주형 데이터) ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "전체 범주형 변수 요약 및 시장 쏠림 현상")
    add_card(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.7))
    tb_left = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.0), Inches(4.1))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "범주형 데이터 핵심 카운트"
    p_l.font.name = FONT_TITLE
    p_l.font.size = Pt(20)
    p_l.font.bold = True
    p_l.font.color.rgb = COLOR_PRIMARY_TEXT
    cat_items = [
        ("고유 도서 수", "1,000 권 (중복이 없는 유일 도서명)"),
        ("고유 저자 수", "869 명 (상위 스타 저자의 일부 다작 존재)"),
        ("등록 출판사 수", "187 개사 (과반의 독점 현상 관찰)"),
        ("분철 미지원 도서", "853 권 (전체의 85.3% 점유)"),
        ("분철 지원 도서", "147 권 (전체의 14.7% 점유)")
    ]
    for label, val in cat_items:
        p = tf_l.add_paragraph()
        p.text = f"\n• {label}:  "
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        run = p.add_run()
        run.text = val
        run.font.name = FONT_BODY
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = COLOR_SECONDARY_TEXT
    add_card(slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.7))
    tb_right = slide.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(5.0), Inches(4.1))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "범주형 데이터 시사점"
    p_r.font.name = FONT_TITLE
    p_r.font.size = Pt(20)
    p_r.font.bold = True
    p_r.font.color.rgb = COLOR_BRICK
    insights_l = [
        "1. 저자 집중도 완화: 1,000개 베스트셀러 중 고유 저자가 869명에 달해 특정 스타 저자가 전체를 독식하기보다는 다양한 도서가 순위에 진입합니다.",
        "2. 출판사 양극화 심각: 반면 출판사는 187개사이지만 한빛미디어 등 소수 대형 출판사가 상위권을 지배하는 독점 성향을 드러냅니다.",
        "3. 기능성 편의 격차: 분철 서비스가 제공되는 서적은 14.7%에 불과하여 독자들의 잠재적 니즈 대비 여전히 공급이 희소한 틈새 시장입니다."
    ]
    for ins in insights_l:
        p = tf_r.add_paragraph()
        p.text = f"\n{ins}"
        p.font.name = FONT_BODY
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        p.line_spacing = 1.3
    add_footer(slide, 5)
    set_speaker_notes(slide,
        "다섯 번째 장표는 범주형 변수의 통계 요약입니다. "
        "고유 도서와 저자 수 대비 등록 출판사 수가 187개사로 좁아집니다. "
        "특히 눈에 띄는 대목은 분철 서비스가 지원되는 도서가 14.7%에 그친다는 점입니다. "
        "독자들의 숨은 니즈 대비 공급이 상당히 저조한 과소 공급 구도를 띠고 있습니다. "
        "이를 기회 삼아 전략적 론칭을 고민해야 합니다."
    )

    # --- Slide 6: PART 1 인트로 ---
    add_section_intro(6, "PART 1", "단변량 분석 및 핵심 분포 탐색", 
                      "판매가, 도서 평점 만족도, 판매지수 형태, 분철 서비스 여부, 메이저 출판사 분포 등\n주요 개별 지표들의 통계적 분포와 비즈니스 기초 인사이트를 도출합니다.")

    # --- Slide 7: 판매가 빈도 분포 ---
    add_standard_chart_slide(
        page_num=7,
        title="1.1 판매가(Sale Price) 빈도 분포 분석",
        chart_filename="01_sale_price_distribution.png",
        stat_title="판매가 구간별 기초 통계",
        stats_list=[
            ("집중 구간", "20,000원 ~ 30,000원 대역 과밀 분포"),
            ("평균 판매가", "23,480 원"),
            ("중위 판매가", "22,500 원"),
            ("최소/최대가", "4,500원 / 67,000원")
        ],
        insight_text="Yes24 IT 베스트셀러 도서들은 약 2만 원에서 3만 원 사이 대역에 70% 가까운 수량이 밀집되어 있습니다. 이는 독자들이 기술 서적을 구매할 때 심리적으로 수용하기 가장 편한 표준 가격 대역대가 2만 원대 초중반에 견고하게 차단되어 있음을 방증합니다.",
        icon_name="price"
    )
    set_speaker_notes(prs.slides[-1],
        "일곱 번째 슬라이드는 도서 판매가의 빈도 분포 차트 분석입니다. "
        "약 20,000원에서 30,000원 구간의 히스토그램 기둥이 압도적으로 높이 솟아 있습니다. "
        "IT 기술 도서 시장의 표준 가격대가 2만 원 중후반에 정형화되어 있고, "
        "소비자 또한 이 가격대에 익숙해져 있으며, 출판사도 이 틀에 맞춰 제품 라인업을 기획하고 있습니다."
    )

    # --- Slide 8: 판매가 심층 분석 ---
    add_deep_insight_slide(
        page_num=8,
        title="1.1 [심층분석] 가격 심리적 저항선과 세그먼트 전략",
        left_title="심리적 마지노선 분석",
        left_points=[
            ("3만원 가격 장벽", "판매가가 30,000원 선을 돌파할 때, 독자가 즉각 결제하기 주저하는 심리적 이탈 저항선이 뚜렷이 나타납니다."),
            ("보급형 입문서 설계", "단기적 베스트셀러 양적 성장을 도모하는 입문서/문법서는 최종 할인가가 22,000원~26,000원 구간에 걸치도록 타겟 정가를 기획하는 것이 최선입니다.")
        ],
        right_title="고가 프리미엄 세그먼트 전략",
        right_points=[
            ("전문성 희소가치 수용", "RAG 아키텍처 구축이나 LLM 미세 조정 등 난이도가 높고 시의성 강한 독점적 전문서라면, 4만 원 이상으로 가격을 높여도 구매로 이어집니다."),
            ("가치 대비 가격 책정", "기술 대체 불가능성이 높다면 일괄 가격 타협보다는 분량을 대폭 보강하고 사은 혜택을 얹어 고가 프리미엄 전략을 펴는 편이 마진 방어에 유리합니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "여덟 번째 슬라이드는 판매가 분포 분석에 기인한 심층 전략 제언입니다. "
        "독자들은 3만 원 이상의 입문 도서에는 거부감을 가지지만, "
        "주제가 대체 불가능하고 희소하다면 4만 원이 넘어도 결제에 응합니다. "
        "이러한 심리를 분석하여 출판 가격 정책을 보급형 문법서(2만 원 대)와 프리미엄 실무서(4만 원 대)로 이원화 포지셔닝해야 합니다."
    )

    # --- Slide 9: 도서 평점 분포 ---
    add_standard_chart_slide(
        page_num=9,
        title="1.2 도서 평점(Rating) 빈도 분포 분석",
        chart_filename="02_rating_distribution.png",
        stat_title="평점 분포 수치 요약",
        stats_list=[
            ("만점 밀집도", "9.5 ~ 10.0 점 대역에 극단적인 우측 쏠림"),
            ("전체 평균 평점", "7.50 점 (0점 무평가 도서 혼재 효과)"),
            ("평가도서 중위수", "9.70 점 (독자 만족도 전반적 우수)"),
            ("최소/최대 평점", "0.00 점 / 10.00 점")
        ],
        insight_text="도서 평점의 히스토그램을 분석하면 대다수 평점이 만점에 밀집되어 일반 평점 점수 자체는 도서의 우열을 가리는 실질적 잣대로서의 변별력을 완전히 소실했습니다. 평균 7.5점은 리뷰가 0건인 도서들의 영점 효과로 발생한 굴절입니다.",
        icon_name="star"
    )
    set_speaker_notes(prs.slides[-1],
        "아홉 번째 슬라이드는 도서 평점의 빈도 분포 데이터입니다. "
        "평점은 9.5~10점 대역에 극단적으로 치우쳐져 있습니다. "
        "단순 평균은 7.5점이지만, 이는 아직 독자 평가를 1건도 받지 못해 0점으로 코딩된 신규 발간 도서들이 하단에 모여 있기 때문입니다. "
        "따라서 평점 점수 자체는 더 이상 도서 간 품질 변별 기준이 되지 못합니다."
    )

    # --- Slide 10: 평점 심층 분석 ---
    add_deep_insight_slide(
        page_num=10,
        title="1.2 [심층분석] 평점 인플레이션 극복 및 0점 돌파 전략",
        left_title="평점 인플레이션의 대용 지표",
        left_points=[
            ("사회적 증거(Social Proof) 작동", "독자들은 평점 10점 만점을 신임하기보다, 평점 9.6점이라도 상세 텍스트 리뷰가 100개 넘게 누적되어 있는 책을 강력히 의지합니다."),
            ("리뷰 수로 무게중심 이동", "따라서 평점 스펙 높이기보다는, 절대적인 평가의 볼륨(리뷰 건수)을 띄우는 것이 잠재 독자를 구매 전환시키는 소셜 락인입니다.")
        ],
        right_title="신간 도서 0점 신속 탈출 지침",
        right_points=[
            ("0점 도서의 치명적 노출 방해", "0점 도서는 불량이 아니라 미검증 낙인 효과를 주어 독자의 장바구니 결제를 심리적으로 차단하는 부작용을 유발합니다."),
            ("초기 런칭 1개월 마케팅", "도서 출간 즉시 초기 가용 마케팅 화력을 서평단 유치와 독자 리뷰 마일리지 보상에 집중해, 최소 10건 이상의 초기 서평을 신속 확보해야 합니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "열 번째 슬라이드는 평점 인플레이션을 극복하는 마케팅 실천 지침입니다. "
        "평점이 만점 일색이므로 독자는 평점 대신 리뷰가 얼마나 정성스레 쌓였는지를 봅니다. "
        "또한 신작이 0점에 묶여 있으면 구매를 꺼리므로, "
        "출간 직후 초기 1개월 골든타임 이내에 다양한 서평단과 인센티브를 연계해 빠르게 0점을 탈출해야 노출과 매출이 오릅니다."
    )

    # --- Slide 11: 판매지수 분포 ---
    add_standard_chart_slide(
        page_num=11,
        title="1.3 판매지수(Sale Index) 상자 그림 분포 분석",
        chart_filename="03_sale_index_boxplot.png",
        stat_title="판매지수 통계 및 사분위",
        stats_list=[
            ("중위 판매지수", "1,236 점 (현실적 성과 척도)"),
            ("평균 판매지수", "3,026 점 (상위 아웃라이어 왜곡 효과)"),
            ("최대 판매지수", "87,480 점 (독보적 메가 히트 타이틀)"),
            ("표준편차(std)", "7,077 점 (매우 큰 분포 변동성)")
        ],
        insight_text="판매지수의 박스플롯을 보면 대다수 베스트셀러가 5,000점 이하 대역에 밀집된 반면, 상위 킬러 타이틀들은 최고 8.7만 점에 육박합니다. 베스트셀러 랭킹 내에서도 엄청난 판매 점유 양극화가 일어나는 전형적인 롱테일 구조입니다.",
        icon_name="trend"
    )
    set_speaker_notes(prs.slides[-1],
        "열한 번째 장표는 판매지수 상자 그림 분석입니다. "
        "차트에서 보이듯 거의 모든 도서가 5,000점 이하의 바닥권에 수렴해 있으나 상위 점들은 우주 위로 흩어져 있습니다. "
        "소수의 메가 히트작이 시장 트래픽과 출판 유통사의 매출을 견인하고 있는 구조입니다. "
        "따라서 전체 평균인 3,026점은 신작 론칭 시 적합한 기대를 제공하지 못합니다."
    )

    # --- Slide 12: 판매지수 심층 분석 ---
    add_deep_insight_slide(
        page_num=12,
        title="1.3 [심층분석] 롱테일 시장 구조와 현실적 KPI 수립",
        left_title="파레토 법칙의 지배력 검증",
        left_points=[
            ("소수 킬러 타이틀 독점", "베스트셀러 1,000권 목록에서도 상위 20%의 서적이 총 판매지수의 압도적 지분을 장악하며 유통 플랫폼의 헤드 역할을 담당합니다."),
            ("평균 왜곡의 함정", "출판 기획 시 단순 평균값(3,026)을 목표 판매지수로 삼으면 과도한 마케팅 비용 태우기나 예산 편성 불균형을 야기할 수 있습니다.")
        ],
        right_title="중위수(Median) 벤치마킹 액션",
        right_points=[
            ("1차 마일스톤 설정", "현실적인 신작의 성공 안착 기준은 50% 중위선인 '판매지수 1,236점'을 달성하는 것으로 잡고, 2차 확산 전략을 모색하는 것이 리스크를 관리하는 지름길입니다."),
            ("업데이트 주기적 투자", "상위 극소수 아웃라이어 타이틀은 1회성 마케팅 결과가 아닙니다. 깃허브 피드백을 수용한 신속한 개정판 관리로 탄생시킨 지속적 업그레이드의 산물입니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "열두 번째 장표는 판매지수의 왜곡을 고려한 신작 경영 가이드라인입니다. "
        "메가 히트 서적들이 지수를 지배하므로, 신진 출판사나 저자는 중위값 1,236점을 1차 도달 목표로 잡는 것이 예산 손실을 방지하는 현명한 행동입니다. "
        "동시에 최상위 아웃라이어 서적을 관찰해보면, 지속적인 예제 코드 갱신과 주기적 개정판 출간을 무기로 영토를 지킨 스테디셀러들입니다."
    )

    # --- Slide 13: 분철 서비스 빈도 ---
    add_standard_chart_slide(
        page_num=13,
        title="1.4 분철 서비스 제공 여부(Spring Service) 분석",
        chart_filename="04_spring_service_count.png",
        stat_title="분철 제공 수량 및 점유비",
        stats_list=[
            ("분철 미제공 도서(N)", "853 권 (85.3% 압도적 지분)"),
            ("분철 제공 도서(Y)", "147 권 (14.7% 소수 공급)"),
            ("주요 적용 도서 특성", "수험서, 개발 실무 매뉴얼, 두꺼운 대학 전공서")
        ],
        insight_text="Yes24 베스트셀러 내 분철 지원 서적은 14.7%에 불과합니다. 제작 공정 단가와 물류 유통 편의상 대다수 출판사가 분철을 지양하고 있으나, 수험서나 실무 전공 서적을 펼쳐놓고 학습하는 독자들에게는 잠재적 수요가 거대하게 깔려 있는 영역입니다.",
        icon_name="spring"
    )
    set_speaker_notes(prs.slides[-1],
        "열세 번째 장표는 분철 서비스 제공 유무의 분포입니다. "
        "일반 무선 제본 서적이 85.3%이며 분철은 14.7%에 그칩니다. "
        "대형 출판사나 인쇄소는 공정 관리 복잡성과 배송 중 훼손 리스크 때문에 분철을 꺼려 왔지만, "
        "컴퓨터를 보며 실습을 타이핑해야 하는 IT 서적 독자층의 갈망은 매우 높습니다. "
        "이를 선점할 수 있다면 훌륭한 차별화가 됩니다."
    )

    # --- Slide 14: 상위 30개 출판사별 도서 수 ---
    add_standard_chart_slide(
        page_num=14,
        title="1.5 상위 30개 출판사별 베스트셀러 점유율",
        chart_filename="05_top_30_publishers.png",
        stat_title="상위 출판사 점유 통계",
        stats_list=[
            ("압도적 1위", "한빛미디어 (150 권, 15.0% 독점)"),
            ("2위 ~ 4위권", "길벗(7.4%), 제이펍(5.1%), 이지스퍼블리싱(5.0%)"),
            ("5위 신흥 주자", "골든래빗 (43 권, 4.3% 가파른 상승)"),
            ("상위 10개사 점유율", "과반(약 50.2%) 돌파 (나머지 177개사 분산)")
        ],
        insight_text="국내 IT 베스트셀러 시장은 대형 메이저 브랜드의 점유 집중도가 매우 극심합니다. 상위 10개 퍼블리셔가 전체 시장의 과반을 점유하고 있어, 신진 브랜드가 독자적인 힘으로 랭킹에 오르기란 매우 가파른 진입 장벽이 가로막고 있습니다.",
        icon_name="building"
    )
    set_speaker_notes(prs.slides[-1],
        "열네 번째 슬라이드는 상위 출판사별 베스트셀러 점유 데이터입니다. "
        "한빛미디어가 15%의 도서 점유율로 독보적인 1위이고, 길벗, 제이펍, 이지스퍼블리싱, 골든래빗이 그 뒤를 잇습니다. "
        "상위 10개사의 합계가 과반을 차지하므로, 신생 기획사가 단독 런칭하여 눈도장을 받기란 여간 어려운 일이 아닙니다."
    )

    # --- Slide 15: 출판사 심층 분석 ---
    add_deep_insight_slide(
        page_num=15,
        title="1.5 [심층분석] 메이저 출판사 브랜드 락인과 상생 전략",
        left_title="대형 퍼블리셔의 후광 효과",
        left_points=[
            ("베타 독자단 인프라", "한빛이나 길벗 등은 수만 명에 달하는 IT 전문 베타 리더단 풀과 충성도 높은 기술 뉴스레터 구독자 락인 체계를 가지고 있습니다."),
            ("기생/협력 투고 전략", "무명 저자나 신인 필자는 마케팅 인프라가 전무하므로, 대형 출판사의 기획 투고를 통과해 그들의 거대한 브랜드 후광을 타는 것이 성공 기회를 열어줍니다.")
        ],
        right_title="틈새 린(Lean) 론칭 전략",
        right_points=[
            ("대형사의 느린 속도 틈새", "대형 퍼블리셔들은 내부 결재와 편집 공정이 다소 길어, 신생 프레임워크나 최신 AI API의 마이크로 업데이트에 민첩하게 반응하기 힘듭니다."),
            ("마이크로 커뮤니티 타겟", "독립 출판이나 소형 기획사는 초단기 트렌드 원고를 선점하고, 개발자 단톡방이나 디스코드, 깃허브 서포트를 활용해 기동성 있게 시장을 우회 침투해야 승산이 있습니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "열다섯 번째 장표는 메이저 과점 구도에 대한 우회 액션 가이드입니다. "
        "대형사의 거대한 마케팅 리소스를 이기기 힘드므로 투고를 통한 윈윈 파트너십을 우선시하되, "
        "만약 소형 출판사가 자생하고자 한다면 대형사가 따라올 수 없는 의사결정의 날렵함을 갖춰야 합니다. "
        "최신 챗GPT나 클로드 코딩 관련 주제를 발빠르게 다루고 개발자 소셜 채널을 직공략하는 게릴라 전술이 유효합니다."
    )

    # --- Slide 16: PART 2 인트로 ---
    add_section_intro(16, "PART 2", "이변량 및 다변량 연계 분석", 
                      "할인율 및 분철 서비스와 판매지수 간의 교차 영향력 검증,\n최근 출판 트렌드 시계열 분석, 다변량 피어슨 상관관계, TF-IDF 텍스트 마이닝을 통한\nIT 도서 시장의 실질적인 수요 동력(Driver)을 도출합니다.")

    # --- Slide 17: 할인율 vs 판매지수 ---
    add_standard_chart_slide(
        page_num=17,
        title="2.1 할인율(Discount Rate) vs 판매지수 산점도",
        chart_filename="06_discount_vs_sale_index.png",
        stat_title="할인 정책별 세부 지표",
        stats_list=[
            ("10% 최대 할인 적용 도서 수", "851 권  |  평균 판매지수: 3,394 점"),
            ("0% 무할인(정가) 적용 도서 수", "123 권  |  평균 판매지수: 826 점"),
            ("5% 부분 할인 적용 도서 수", "25 권  |  평균 판매지수: 1,276 점"),
            ("도서정가제에 따른 가격 통제", "10% 할인이 실질적 판매 안착의 필수적 기본 조건화")
        ],
        insight_text="거의 모든 베스트셀러 도서가 최대 상한선인 10% 할인을 깔고 유통되고 있습니다. 정가를 지키는 무할인 도서는 판매지수가 4분의 1 토막에 그치므로, 10% 할인은 선택이 아닌 필수입니다. 추가 경쟁력은 금전적 마찰 회피가 아닌 비금전적 가치 추가에서 발굴해야 합니다.",
        icon_name="percent"
    )
    set_speaker_notes(prs.slides[-1],
        "열일곱 번째 장표는 할인율과 판매지수 간의 산점도 분석입니다. "
        "법이 허용하는 최대치인 10% 할인을 적용하는 도서가 압도적인 수치 성과를 보입니다. "
        "정가 도서는 평균 지수가 현저히 처집니다. "
        "따라서 10% 할인은 유통 진입의 공통 기본 전제조건이며, "
        "이를 넘어선 경쟁력 창출은 가격 깎기가 아닌 품질 및 기능 차별화로 달성해야 합니다."
    )

    # --- Slide 18: 분철 서비스 여부 vs 판매지수 ---
    add_standard_chart_slide(
        page_num=18,
        title="2.2 분철 여부에 따른 판매 성과 교차 분석",
        chart_filename="07_spring_service_vs_sale_index.png",
        stat_title="분철 여부별 판매지수 비교",
        stats_list=[
            ("분철 지원 그룹 (Y)", "147 권  |  평균 판매지수: 8,687 점 (최대: 87,480)"),
            ("분철 미지원 그룹 (N)", "853 권  |  평균 판매지수: 2,050 점 (최대: 71,388)"),
            ("평균 격차 배율", "약 4.2배 수준의 평균 판매지수 초과 달성"),
            ("중위 지수 비교", "Y 그룹 중위수: 4,062점  |  N 그룹 중위수: 1,047점")
        ],
        insight_text="분철 서비스를 옵션으로 장착한 도서들의 평균 판매 성과가 미지원 도서 대비 무려 4.2배가량 압도적으로 높게 도출되었습니다. 이는 단편적인 제작비 증가분 이상의 판매 견인 및 락인 효과가 작용함을 통계적 증거로 규명하는 극적인 순간입니다.",
        icon_name="spring"
    )
    set_speaker_notes(prs.slides[-1],
        "열여덟 번째 장표는 분철 여부별 판매 성과를 분석한 연계 차트입니다. "
        "평균치와 중위수 모두 분철 서비스를 장착한 도서군이 4배 가까이 앞서 나갑니다. "
        "이 강력한 격차는 기술을 실무적으로 공부할 때 링 제본이 선사하는 가독성과 실습 효율이 "
        "독자들의 최종 선택을 드라이브하는 거대한 촉매가 됨을 단적으로 드러냅니다."
    )

    # --- Slide 19: 분철 서비스 심층 분석 ---
    add_deep_insight_slide(
        page_num=19,
        title="2.2 [심층분석] 분철 가치와 출판 유통사 락인 설계",
        left_title="독자 실용 소비주의 이해",
        left_points=[
            ("핸즈프리(Hands-free) 학습 니즈", "IT 독자들은 책을 모니터 옆에 완전히 젖혀두고 키보드를 자유롭게 조작하고 싶어 합니다. 일반 제본은 자꾸 책이 접혀 지속적인 피로를 가중시킵니다."),
            ("가치 체감 대비 저렴한 투자", "독자는 1,000원~2,000원의 소액 추가금을 주더라도 책의 활용성을 높일 수 있다면 흔쾌히 구매 옵션을 선택하는 패턴을 가집니다.")
        ],
        right_title="출판-유통망 연계 마케팅",
        right_points=[
            ("출판 유통의 기본 탑재", "500페이지가 넘는 기술 서적의 인쇄 기획 시점부터 출판사는 예스24의 주문 연계 분철 라인업을 상설 채널로 확보하고 전면 홍보해야 합니다."),
            ("구매 상세 페이지 노출 강화", "상세 페이지 최상단 배너에 '독자 만족 스프링 분철 옵션 완비' 표시를 강조하여, 분철을 미지원하는 타사 도서와의 결정적 구매 우위를 확립합니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "열아홉 번째 장표는 분철 서비스 성과에 기반한 실무 액션안입니다. "
        "IT 서적 구매자들에게 있어 링 제본은 단순 옵션이 아니라 실습 피로를 줄여주는 킬러 서비스입니다. "
        "출판사와 총판, 유통사는 신간 발간 즉시 이 분철 가치를 메인에 소구해 경쟁사의 일반 서적에 대한 "
        "기능적 격차를 유도해야 가격 저항선을 우회하여 성공할 수 있습니다."
    )

    # --- Slide 20: 출판년도별 추이 ---
    add_standard_chart_slide(
        page_num=20,
        title="2.3 출판년도별 도서 수 및 평균 판매지수 추이",
        chart_filename="08_publish_year_trends.png",
        stat_title="출판 연도별 성과 통계",
        stats_list=[
            ("2025년 출간 도서", "367 권  |  평균 판매지수: 3,902 점"),
            ("2026년 출간 도서", "356 권  |  평균 판매지수: 2,404 점"),
            ("2024년 출간 도서", "138 권  |  평균 판매지수: 1,747 점"),
            ("2023년 이전 도서", "총 139 권  |  평균 지수 3,000점 이하로 급감 흐름")
        ],
        insight_text="2025~2026년 신간들이 베스트셀러 진입 권수와 판매지수의 대부분을 장악하고 있습니다. 연도가 지날수록 도서 수가 급감하여 IT 분야 서적의 감가상각과 트렌드 이탈 속도가 초고속으로 전개되고 있음을 보여줍니다.",
        icon_name="calendar"
    )
    set_speaker_notes(prs.slides[-1],
        "스무 번째 슬라이드는 도서 출판 연도별 시계열 성과 분석입니다. "
        "보시는 바와 같이 2025년과 2026년에 출간된 지 얼마 되지 않은 신간들의 판매지수가 가장 거대합니다. "
        "연차가 단 2년만 지나도 베스트셀러 점유 개수가 급격하게 축소됩니다. "
        "기술 유효 주기가 지독히 짧은 IT 산업군의 특성이 그대로 도서 시장 생명력에 투영되는 것입니다."
    )

    # --- Slide 21: 출판년도 심층 분석 ---
    add_deep_insight_slide(
        page_num=21,
        title="2.3 [심층분석] 시계열 지식 감가상각과 린 출판 대응",
        left_title="IT 도서 지식의 급격한 유통기한",
        left_points=[
            ("오픈소스 버전업 리스크", "라이브러리 마이너 버전 패치 하나만으로도 실습 예제 코드가 실행 에러를 뿜게 되며, 이는 곧장 소비자 평판 추락으로 이어집니다."),
            ("기존 기획 관행 탈피", "과거의 1년짜리 고전적 기획-집필-출판 루프에서 벗어나지 않으면 출간하는 당일 이미 구닥다리 기술서가 될 위험이 농후합니다.")
        ],
        right_title="린 퍼블리싱(Lean Publishing) 도입",
        right_points=[
            ("실시간 독자 피드백 결합", "원고 기획서 단계에서부터 온라인 깃허브 저장소를 개설하고, 독자들과 온라인에서 집필 원고를 소통하며 완성 시점에 신속히 찍어내야 합니다."),
            ("콘텐츠 두께 슬림화 전략", "트렌드가 불을 뿜는 프론트엔드나 생성형 AI 기술 서적은 분량을 200페이지 내외로 축소하여 런칭 주기를 최대한 민첩하게 회전해야 합니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "스물한 번째 장표는 지식 감가상각을 우회하는 린 퍼블리싱 대응책입니다. "
        "코드가 안 돌아가서 발생하는 평판 추락을 막기 위해 지속해서 피드백을 받는 구조를 도입하고, "
        "집필 호흡이 너무 긴 바이블 책 위주의 생산에서 벗어나 "
        "가볍고 명확한 슬림북 개정판을 기민하게 회전 유통하는 린 출판 모델이 절대적으로 요구됩니다."
    )

    # --- Slide 22: 상관관계 히트맵 ---
    add_standard_chart_slide(
        page_num=22,
        title="2.4 수치형 변수 간 상관관계 피어슨 상관행렬",
        chart_filename="09_correlation_heatmap.png",
        stat_title="피어슨 상관계수(r) 결과",
        stats_list=[
            ("판매지수 vs 리뷰 수", "r = 0.214 (유의미한 정적 상관성 성립)"),
            ("판매지수 vs 평점", "r = 0.156 (약한 정적 상관성)"),
            ("판매지수 vs 할인율", "r = 0.124 (미약한 정적 상관성)"),
            ("할인율 vs 포인트적립", "r = 0.658 (가격 마케팅 변수 간 상호 수렴)")
        ],
        insight_text="상관관계 분석 결과, 판매지수와 가장 높은 인과적 연계를 보이는 변수는 '리뷰 수'로 규명되었습니다. 평점 수치 자체보다 리뷰의 절대 수량이 온라인 구매 채널에서 소비자의 안도감을 유도하는 가장 큰 유발 요인입니다.",
        icon_name="chart"
    )
    set_speaker_notes(prs.slides[-1],
        "스물두 번째 장표는 피어슨 상관행렬 히트맵에 대한 분석 결과입니다. "
        "주요 수치 변수 간의 연계성을 보면 판매지수와 리뷰 수 간에 가장 확실한 유의 관계가 발견됩니다. "
        "반면 판매가나 정가, 적립 포인트 등은 판매지수와 거의 상관이 없었습니다. "
        "가격을 깎아주는 혜택보다 '앞서 사본 독자들의 평판이 두터운가'가 의사결정의 무게추라는 의미입니다."
    )

    # --- Slide 23: 상관관계 심층 분석 ---
    add_deep_insight_slide(
        page_num=23,
        title="2.4 [심층분석] 리뷰 자산과 소셜 증명(Social Proof)의 선순환",
        left_title="별점보다 두터운 리뷰 볼륨",
        left_points=[
            ("학습 삽질 비용 회피", "IT 독자가 지불하는 가장 큰 비용은 책값이 아니라 내 주말 시간입니다. 독자는 실패 리스크를 줄이기 위해 풍부한 리뷰 텍스트를 찾습니다."),
            ("텍스트 피드백의 힘", "구체적인 실습 삽질기나 문제 해결 방안이 기록된 리뷰가 많을수록 제품 만족에 대한 확신을 심어주어 구매 결정을 촉진합니다.")
        ],
        right_title="런칭 초기의 리뷰 빌드업 지침",
        right_points=[
            ("초기 30일 골든타임 관리", "출간 1개월 이내에 정성스러운 텍스트 및 사진 리뷰가 플랫폼 메인에 노출될 수 있도록 독자 서평단을 정교하게 스케줄링해야 합니다."),
            ("인센티브 선순환 구조", "도서 예제 코드 깃허브나 책 뒤편에 '리뷰 작성 시 추가 소스 제공 및 마일리지 프로모션' 링크를 연계해 구매자의 자발적 평판 기여를 끊임없이 조장해야 합니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "스물세 번째 장표는 상관분석에 따른 소셜 프루프 락인 마케팅입니다. "
        "독자는 학습 시간의 낭비를 가장 경계하기 때문에 먼저 사본 선배 독자들의 텍스트 후기를 열심히 파고듭니다. "
        "따라서 런칭 초기 30일 이내에 높은 만족도의 정성 서평이 예스24 상세 페이지에 쌓이도록 유치하고, "
        "도서 내부에 리뷰 인센티브 연동 구조를 필수 탑재할 것을 권고합니다."
    )

    # --- Slide 24: 할인율-분철 여부별 판매지수 ---
    add_standard_chart_slide(
        page_num=24,
        title="2.5 [다변량] 할인율-분철 여부별 평균 판매지수",
        chart_filename="10_discount_spring_vs_sale_index.png",
        stat_title="할인 x 분철 다변량 피봇 결과",
        stats_list=[
            ("10% 할인 + 분철지원 (Y)", "평균 판매지수: 9,131 점 (최상의 성과 구간)"),
            ("10% 할인 + 분철미지원 (N)", "평균 판매지수: 2,321 점 (보편적 평균 구간)"),
            ("0% 무할인 + 분철지원 (Y)", "평균 판매지수: 3,653 점 (할인 N 그룹 초과)"),
            ("0% 무할인 + 분철미지원 (N)", "평균 판매지수: 603 점 (최하의 성과 구간)")
        ],
        insight_text="10% 기본 할인에 분철 서비스를 결합한 그룹이 압도적 1위의 판매지수(9,131)를 냅니다. 흥미롭게도 할인을 0%로 고수하더라도 분철을 장착하기만 하면 평균 3,653점으로 분철이 없는 할인 도서 집단보다 우수한 매출 성적을 거두었습니다.",
        icon_name="percent"
    )
    set_speaker_notes(prs.slides[-1],
        "스물네 번째 장표는 할인과 분철을 결합한 다변량 피봇 결과입니다. "
        "당연히 할인과 분철을 다 준 그룹이 압도적 1위입니다. "
        "주목할 점은 정가를 지키는 무할인 도서 집단이라도 분철을 탑재하기만 하면 평균 지수가 3,653점까지 치솟아, "
        "분철 없이 10%를 깎아준 도서군(2,321)을 훌쩍 앞질렀다는 것입니다. "
        "분철이라는 부가 서비스가 10%의 현금성 할인 가치 장벽을 가볍게 뛰어넘는다는 극적인 결과입니다."
    )

    # --- Slide 25: 도서명 TF-IDF 중요 키워드 ---
    add_standard_chart_slide(
        page_num=25,
        title="2.6 도서명 기준 TF-IDF 중요 키워드 분석",
        chart_filename="11_goods_name_tfidf.png",
        stat_title="상위 TF-IDF 가중치 단어",
        stats_list=[
            ("핵심 기술 테마", "'AI' (가중치 0.059), '파이썬' (0.019), '인공지능' (0.012)"),
            ("인공지능 보조 도구", "'제미나이' (0.010), '클로드' (0.009), '챗GPT' (0.012)"),
            ("학습 지향 키워드", "'코딩' (0.017), '가이드' (0.016), '입문' (0.011)"),
            ("생산성 향상 키워드", "'활용' (0.011), '실무' (0.009), '실전' (0.009)")
        ],
        insight_text="도서 제목의 텍스트 마이닝 결과, 단순한 프로그래밍 기본 이론에서 벗어나 AI 도구(클로드, 제미나이 등)를 활용한 코딩 자동화와 실용적인 파이썬 기반 데이터 과학 관련 주제가 가중치 상위를 석권하고 있습니다.",
        icon_name="key"
    )
    set_speaker_notes(prs.slides[-1],
        "스물다섯 번째 장표는 도서명 텍스트에 대한 TF-IDF 키워드 가중치 마이닝입니다. "
        "시장의 수요는 이미 'ai', '파이썬', '제미나이', '클로드' 같은 인공지능 보조 도구와 코딩 연계 실용주의로 가파르게 고도화되었습니다. "
        "더 이상 고전적 문법 서적만으로는 시장에서 환영받지 못하며, "
        "모든 도서에 AI 협업 코딩 기법을 융합하는 패러다임 전환이 일어나고 있습니다."
    )

    # --- Slide 26: 저자별 평균 판매지수 ---
    add_standard_chart_slide(
        page_num=26,
        title="2.7 상위 저자별 도서 수 및 평균 판매지수 분석",
        chart_filename="12_author_mean_sale_index.png",
        stat_title="저자 세그먼트 성적표",
        stats_list=[
            ("최다 도서 등록 저자군", "놀이교육콘텐츠랩(9권, 평균 1,170점)"),
            ("최다 등록 출판사 서적", "해람북스 기획팀(5권, 평균 1,213점)"),
            ("소수 스타 저자 쏠림", "오힘찬 저자(3권, 평균 36,848점)"),
            ("다작 vs 고효율 전략", "도서 등록 권수와 평균 판매지수 간 반비례 양상")
        ],
        insight_text="단순 다작 저자들보다, 최신 기술 흐름(예: 제미나이 마스터북)을 명확하게 파헤친 스타 저자의 도서가 압도적인 평균 판매지수 스파이크를 형성합니다. 독자들은 저자의 다작 여부보다 화제성과 실용성에 충성합니다.",
        icon_name="users"
    )
    set_speaker_notes(prs.slides[-1],
        "스물여섯 번째 장표는 상위 저자들의 베스트셀러 도서 수 및 평균 판매지수입니다. "
        "기대를 달리하여, 단순 많은 수의 아동 코딩 서적을 찍어낸 다작 필자 집단보다 "
        "최신 인공지능 트렌드를 관통하여 잘 쓰여진 킬러 서적을 발간한 필자군의 평균 성과가 수십 배 높습니다. "
        "IT 전문서 시장에서는 책을 여러 권 발간하는 다작 기획보다 "
        "필자의 기술 브랜딩을 견고하게 실어 한 권의 킬러 명작을 론칭하는 편이 영리합니다."
    )

    # --- Slide 27: 출판 월별 추이 ---
    add_standard_chart_slide(
        page_num=27,
        title="2.8 출판 월별 신간 수 및 평균 판매지수 분석",
        chart_filename="13_publish_month_trends.png",
        stat_title="출판 월별 시즌 성과",
        stats_list=[
            ("상반기(1~6월) 런칭 성과", "신규 진입 도서 집중 및 높은 평균 판매지수 유지"),
            ("최대 판매지수 월", "12월 (평균 3,985점), 9월 (평균 3,743점)"),
            ("최저 판매지수 월", "7월 (평균 2,006점), 10월 (평균 2,343점)"),
            ("새해/새학기 결심 효과", "방학과 신학기 직전에 맞춘 도서 소비 본능 확인")
        ],
        insight_text="상반기에 발행된 도서들의 판매지수 강세가 도드라집니다. 연말연시나 새학기 결심(새해 다짐, 방학 공부 계획)에 부합하여 소비자들이 지갑을 여는 시즌성 효과입니다. 론칭 일정을 겨울방학이나 신학기 직전에 포커싱해야 합니다.",
        icon_name="calendar"
    )
    set_speaker_notes(prs.slides[-1],
        "스물일곱 번째 장표는 도서 출판월별 시즌 통계 분석입니다. "
        "도서 소비는 결심 시즌인 12월과 1월, 그리고 신학기 준비인 2월과 9월에 매우 극적인 매출 성과를 보입니다. "
        "따라서 신간 기획안을 도출할 때 막연히 완성되는 날에 출판하기보다는, "
        "소비자 결심 수요가 지갑을 열어젖히는 연말연시나 방학 직전에 맞춰 발매일을 포커싱해야 마케팅 시너지를 누립니다."
    )

    # --- Slide 28: 도서 태그 수 vs 판매지수 ---
    add_standard_chart_slide(
        page_num=28,
        title="2.9 도서 태그 수(Tag Count)와 판매지수 분석",
        chart_filename="14_tag_count_vs_sale_index.png",
        stat_title="태그 수에 따른 성과",
        stats_list=[
            ("적정 태그 수 대역", "2개 ~ 4개 등록 도서의 판매지수 안정적 우위"),
            ("태그 0개 등록 도서 수", "685 권  |  평균 판매지수: 1,815 점"),
            ("태그 5개 등록 도서 수", "47 권  |  평균 판매지수: 8,770 점 (고가도서 쏠림)"),
            ("검색 엔진 최적화(SEO)", "부적합한 태그 남발보다 정교한 타겟 키워드 매핑 권장")
        ],
        insight_text="태그가 아예 없는(0개) 서적들은 노출 기회 박탈로 성적이 부진합니다. 반면 2~4개의 정밀하고 세련된 태그를 매핑해 둔 도서군이 안정적인 지수 흐름을 냅니다. 검색 노출을 위한 최적의 태그 관리가 요구됩니다.",
        icon_name="tag"
    )
    set_speaker_notes(prs.slides[-1],
        "스물여덟 번째 장표는 태그 개수와 판매지수 간의 연계 분석입니다. "
        "태그가 아예 없거나 단순하게만 지정된 다수의 도서들은 온라인 몰 검색에서 차단되어 성과가 낮습니다. "
        "하지만 너무 많은 무차별적 태그 도배도 노이즈가 될 수 있습니다. "
        "독자 유입 검색 엔진 최적화(SEO)를 위해 가장 유의미한 카테고리 태그 3~4개를 엄밀하게 골라 배치하는 관리 규격이 권고됩니다."
    )

    # --- Slide 29: 3부 종합 비즈니스 제언 요약 ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "PART 3. IT 도서 시장 활성화를 위한 5대 핵심 성공 전략", "BUSINESS STRATEGY & ROADMAP")
    strategies = [
        ("1. 가격 포지셔닝 이원화", "입문 기술서는 2.3만~2.7만원 대역 책정으로 가격 저항 최소화, 특수 고난도 기술서는 4만원 대 프리미엄 책정.", "price"),
        ("2. 분철 서비스 연계 활성화", "500p 이상 및 실무 코드가 중심인 도서는 제본 옵션을 기본 제공하여 26% 이상의 추가 매출 동력 확보.", "spring"),
        ("3. 지식 감가상각 방어 프로세스", "트렌드가 민감한 응용서는 린(Lean) 집필 공정을 세팅하여 빠르게 출판하고 수명을 선제적 관리.", "calendar"),
        ("4. 독자 평판 소셜프루프 극대화", "초기 서평단 1개월 빌드업 전략에 리소스를 집중하고, 자발적 리뷰 축적 유인을 도서 내부에 설계.", "review"),
        ("5. 검색 노출 및 태깅 고도화", "무질서한 키워드 태깅 지양, 핵심 노출 키워드 3~4개를 선별해 집중 매핑함으로써 유입 전환율 향상.", "key")
    ]
    for idx, (title, desc, icon) in enumerate(strategies):
        row = idx // 3
        col = idx % 3
        w = Inches(3.9)
        h = Inches(2.2)
        x = Inches(0.6 + col * 4.1)
        y = Inches(1.8 + row * 2.4)
        if idx >= 3:
            x = Inches(2.65 + (idx - 3) * 4.1)
        add_card(slide, x, y, w, h)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_BRICK
        bar.line.fill.background()
        add_icon(slide, icon, x + Inches(0.3), y + Inches(0.3), Inches(0.45))
        tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.8), w - Inches(0.6), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_TITLE
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_PRIMARY_TEXT
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_SECONDARY_TEXT
        p2.line_spacing = 1.2
    add_footer(slide, 29)
    set_speaker_notes(slide,
        "스물아홉 번째 장표는 앞선 정량 데이터를 응축한 20년차 데이터 분석가로서의 5대 핵심 경영 지침입니다. "
        "가격을 2만 원대와 4만 원대로 이원화하고, 두껍고 실습 중심의 도서에는 분철 제본 연계 서비스를 유치하며, "
        "지식의 유통기한을 고려한 린 출판 체제를 채택하고, "
        "리뷰 자산 누적의 골든타임인 발간 후 초기 30일을 집중 마케팅하며, "
        "3~4개의 핵심 태그 선별 매핑을 통해 온라인 검색 노출을 극대화해야 합니다. "
        "이 로드맵을 지킬 때 경영 효율의 눈부신 점프를 획득할 것입니다."
    )

    # --- Slide 30: 결론 및 Q&A ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "감사합니다 (Thank You)", "CONCLUSION & Q&A")
    add_card(slide, Inches(0.6), Inches(1.8), Inches(12.133), Inches(4.7))
    
    # 텍스트 박스의 크기와 정렬 개선 (수직 정렬 여백 보완)
    mid_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.1), Inches(10.333), Inches(4.0))
    tf = mid_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    
    # 메인 타이틀
    p = tf.paragraphs[0]
    p.text = "Q&A 및 미래 출판 패러다임 예측"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_BRICK
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(24)  # 넉넉한 수직 여백 확보
    
    # 첫 번째 정보 (불릿 문자 제거, 중앙 정렬의 불릿 불일치 해결)
    p = tf.add_paragraph()
    p.text = "데이터 분석 결과 요약본 배포: PDF 보고서 및 데이터셋 시각화 템플릿 별도 제공"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_TEXT
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)
    
    # 두 번째 정보
    p = tf.add_paragraph()
    p.text = "질의응답: 분석 설계 모형 및 비즈니스 의사결정 액션 플랜 세부 질문"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_TEXT
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(28)
    
    # 맺음말 타이틀
    p = tf.add_paragraph()
    p.text = "[맺음말]"
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY_TEXT
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)
    
    # 맺음말 상세 (행간 및 문단 여백 보완)
    p = tf.add_paragraph()
    p.text = "이번 EDA 분석을 통해 검증되었듯, IT 도서 시장은 독자의 편의성(분철)과 기술 트렌드의 신속성(린 출판), 그리고 소셜 평판(리뷰 건수)에 의해 지배를 받는 특수한 비즈니스 필드입니다. 본 리포트의 전략을 바탕으로 더욱 민첩하고 유연한 비즈니스 로드맵이 수립되기를 기대합니다. 질문이 있으시다면 기꺼이 답변 드리겠습니다."
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_PRIMARY_TEXT
    p.line_spacing = 1.4
    p.alignment = PP_ALIGN.CENTER
    add_footer(slide, 30)
    set_speaker_notes(slide,
        "이상으로 Yes24 IT/컴퓨터 베스트셀러 EDA 심층 분석 발표를 모두 마무리하겠습니다. "
        "우리는 데이터를 기반으로 전통적인 도서 시장 기획에서 벗어나 "
        "독자의 실용적 편의(분철)와 최신 지식 수명(린 기획), 소셜 평판(리뷰 수)을 조망하는 "
        "진화된 마케팅 체계를 설계해야 함을 검증했습니다. "
        "경청해주셔서 진심으로 감사합니다. 궁금하신 설계 내역이나 추가 지표 분석에 대한 질의에 최선을 다해 답변해 드리겠습니다."
    )

    # 저장 및 마무리
    output_path = "yes24/docs/EDA_Slide_30Pages.pptx"
    prs.save(output_path)
    print(f"PPTX 생성 완료: {output_path}")

if __name__ == "__main__":
    create_presentation()
