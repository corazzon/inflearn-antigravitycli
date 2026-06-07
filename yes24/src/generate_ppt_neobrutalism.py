"""
Yes24 IT/컴퓨터 베스트셀러 EDA PPTX 보고서 네오브루탈리즘 스타일 자동 생성 스크립트

이 스크립트는 yes24/docs/eda_report.md의 분석 결과를 바탕으로,
네오브루탈리즘 스타일(두꺼운 검정 테두리, 하드 오프셋 블랙 그림자, 직각 레이아웃)을 적용한
완전히 새로운 30페이지 분량의 PPTX 프레젠테이션을 생성합니다.
배경은 순백색(#FFFFFF)으로 고정하며, 포인트 컬러는 벽돌색(#B85042)과 노란색(#F5F500)을 매칭합니다.
모든 도형 카드에 2.5pt 두께의 검정 테두리와 입체적인 오프셋 섀도를 수동으로 레이어링하여
독창적인 네오브루탈리즘 인포그래픽을 구현하고, 각 슬라이드별 상세 발표자 노트를 포함합니다.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_neobrutalist_presentation():
    # 프레젠테이션 초기화 및 16:9 슬라이드 크기 설정
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 네오브루탈리즘 디자인 시스템 컬러 정의
    COLOR_BG = RGBColor(255, 255, 255)         # 순백색 배경
    COLOR_BLACK = RGBColor(0, 0, 0)            # 순검은색 (테두리 및 섀도)
    COLOR_BRICK = RGBColor(184, 80, 66)        # 포인트 벽돌색
    COLOR_YELLOW = RGBColor(245, 245, 0)       # 네오 브루탈 옐로우 (아이콘 배경 등)
    COLOR_PRIMARY_TEXT = RGBColor(0, 0, 0)     # 메인 텍스트 (완전 블랙)
    COLOR_SECONDARY_TEXT = RGBColor(100, 100, 100) # 보조 텍스트 (어두운 회색)
    COLOR_CARD_BG = RGBColor(255, 255, 255)    # 카드 기본 순백색

    FONT_TITLE = "Gmarket Sans Bold"
    FONT_BODY = "NanumGothic"

    IMAGE_DIR = "yes24/images"
    ICON_DIR = "yes24/images/icons"

    # 공통 헬퍼 함수들
    def apply_bg(slide):
        """모든 슬라이드에 순백색 배경을 적용하고, 네오브루탈리즘 시그니처 블록 데코를 추가합니다."""
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        
        # 좌측 상단 네오브루탈 데코 (검정 테두리를 두른 벽돌색 사각형)
        deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.8), Inches(0.15))
        deco.fill.solid()
        deco.fill.fore_color.rgb = COLOR_BRICK
        deco.line.color.rgb = COLOR_BLACK
        deco.line.width = Pt(1.5)

    def add_header(slide, title_text, category_text="YES24 IT/COMPUTER BESTSELLER EDA"):
        """상단 카테고리와 메인 타이틀을 추가합니다."""
        cat_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.35), Inches(11.0), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_bottom = tf_cat.margin_right = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = FONT_BODY
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_SECONDARY_TEXT
        
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(12.0), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_bottom = tf_title.margin_right = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY_TEXT

    def add_footer(slide, page_num):
        """하단 굵은 검정 구분선과 페이지 번호를 추가합니다."""
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.8), Inches(12.133), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_BLACK
        line.line.fill.background()
        
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(10.0), Inches(0.3))
        tf = footer_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = "Yes24 IT/컴퓨터 베스트셀러 데이터 심층 EDA 리포트  |  네오브루탈리즘 버전"
        p.font.name = FONT_BODY
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        
        page_box = slide.shapes.add_textbox(Inches(11.733), Inches(6.9), Inches(1.0), Inches(0.3))
        tf_page = page_box.text_frame
        tf_page.word_wrap = True
        tf_page.margin_left = tf_page.margin_top = tf_page.margin_bottom = tf_page.margin_right = 0
        p_page = tf_page.paragraphs[0]
        p_page.text = f"{page_num} / 30"
        p_page.alignment = PP_ALIGN.RIGHT
        p_page.font.name = FONT_TITLE
        p_page.font.size = Pt(10)
        p_page.font.bold = True
        p_page.font.color.rgb = COLOR_BRICK

    def draw_neobrutalist_card(slide, x, y, w, h, fill_color=COLOR_CARD_BG):
        """하드 오프셋 그림자 레이어와 두꺼운 검정 테두리 카드를 겹쳐 그립니다."""
        shadow_offset = Inches(0.08)
        # 1. 그림자 사각형 (블랙 솔리드)
        shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + shadow_offset, y + shadow_offset, w, h)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = COLOR_BLACK
        shadow.line.fill.background()
        
        # 2. 전면 카드 사각형
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        card.line.color.rgb = COLOR_BLACK
        card.line.width = Pt(2.5)
        return card

    def draw_neobrutalist_icon(slide, icon_name, x, y, size=Inches(0.5)):
        """노란색 네오브루탈 배경 블록 위에 아이콘을 정렬합니다."""
        shadow_offset = Inches(0.04)
        # 그림자
        shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + shadow_offset, y + shadow_offset, size, size)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = COLOR_BLACK
        shadow.line.fill.background()
        
        # 전면 블록 (노란색 배경)
        block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
        block.fill.solid()
        block.fill.fore_color.rgb = COLOR_YELLOW
        block.line.color.rgb = COLOR_BLACK
        block.line.width = Pt(1.5)
        
        # 아이콘 그림
        icon_path = os.path.join(ICON_DIR, f"{icon_name}.png")
        if os.path.exists(icon_path):
            slide.shapes.add_picture(icon_path, x, y, width=size, height=size)

    def set_speaker_notes(slide, notes_text):
        """발표자 노트를 추가합니다."""
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

    def add_standard_chart_slide(page_num, title, chart_filename, stat_title, stats_list, insight_text, icon_name="chart"):
        """네오브루탈리즘 레이아웃으로 차트 슬라이드를 구성합니다."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        apply_bg(slide)
        add_header(slide, title)
        
        # 좌측 차트 이미지용 카드 액자 생성
        img_x, img_y = Inches(0.6), Inches(1.8)
        img_w, img_h = Inches(6.0), Inches(4.5)
        draw_neobrutalist_card(slide, img_x, img_y, img_w, img_h)
        
        chart_path = os.path.join(IMAGE_DIR, chart_filename)
        if os.path.exists(chart_path):
            # 2.5pt 테두리를 덮지 않도록 약간 마진을 주어 이미지 배치
            slide.shapes.add_picture(chart_path, img_x + Inches(0.05), img_y + Inches(0.05), img_w - Inches(0.1), img_h - Inches(0.1))
            
        # 우측 설명 카드 배치
        card_x, card_y = Inches(6.9), Inches(1.8)
        card_w, card_h = Inches(5.8), Inches(4.5)
        draw_neobrutalist_card(slide, card_x, card_y, card_w, card_h)
        
        # 아이콘 카드 배치
        draw_neobrutalist_icon(slide, icon_name, card_x + Inches(0.4), card_y + Inches(0.35), Inches(0.5))
        
        # 타이틀용 텍스트박스
        title_tb = slide.shapes.add_textbox(card_x + Inches(1.05), card_y + Inches(0.3), card_w - Inches(1.45), Inches(0.6))
        tf_t = title_tb.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_bottom = tf_t.margin_right = 0
        p_stat_title = tf_t.paragraphs[0]
        p_stat_title.text = stat_title
        p_stat_title.font.name = FONT_TITLE
        p_stat_title.font.size = Pt(18)
        p_stat_title.font.bold = True
        p_stat_title.font.color.rgb = COLOR_PRIMARY_TEXT
        
        # 본문 텍스트박스 분리 배치
        tb = slide.shapes.add_textbox(card_x + Inches(0.4), card_y + Inches(1.05), card_w - Inches(0.8), card_h - Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        # 통계 지표 리스트
        first = True
        for label, val in stats_list:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = f"• {label}: "
            p.font.name = FONT_BODY
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = COLOR_PRIMARY_TEXT
            p.space_after = Pt(4)
            
            run = p.add_run()
            run.text = val
            run.font.name = FONT_BODY
            run.font.size = Pt(12)
            run.font.bold = False
            run.font.color.rgb = COLOR_SECONDARY_TEXT
            
        # 비즈니스 해석 타이틀
        p_ins_title = tf.add_paragraph()
        p_ins_title.text = "[비즈니스 해석]"
        p_ins_title.font.name = FONT_BODY
        p_ins_title.font.size = Pt(13)
        p_ins_title.font.bold = True
        p_ins_title.font.color.rgb = COLOR_BRICK
        p_ins_title.space_before = Pt(10)
        p_ins_title.space_after = Pt(4)
        
        # 해석 상세
        p_ins = tf.add_paragraph()
        p_ins.text = insight_text
        p_ins.font.name = FONT_BODY
        p_ins.font.size = Pt(11)
        p_ins.font.color.rgb = COLOR_PRIMARY_TEXT
        p_ins.line_spacing = 1.3
        
        add_footer(slide, page_num)
        return slide

    def add_deep_insight_slide(page_num, title, left_title, left_points, right_title, right_points):
        """네오브루탈리즘 레이아웃으로 2열 카드 슬라이드를 생성합니다."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        apply_bg(slide)
        add_header(slide, title)
        
        # 좌우 카드 배치
        for i, (c_title, c_points) in enumerate([(left_title, left_points), (right_title, right_points)]):
            x = Inches(0.6 + i * 6.1)
            y = Inches(1.8)
            w = Inches(5.8)
            h = Inches(4.5)
            draw_neobrutalist_card(slide, x, y, w, h)
            
            # 텍스트 상자
            tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), h - Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
            
            p_title = tf.paragraphs[0]
            p_title.text = c_title
            p_title.font.name = FONT_TITLE
            p_title.font.size = Pt(18)
            p_title.font.bold = True
            p_title.font.color.rgb = COLOR_PRIMARY_TEXT
            p_title.space_after = Pt(12)
            
            for p_title_text, p_body in c_points:
                p = tf.add_paragraph()
                p.text = f"• {p_title_text}"
                p.font.name = FONT_BODY
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = COLOR_PRIMARY_TEXT
                p.space_after = Pt(4)
                
                p2 = tf.add_paragraph()
                p2.text = p_body
                p2.font.name = FONT_BODY
                p2.font.size = Pt(11)
                p2.font.color.rgb = COLOR_SECONDARY_TEXT
                p2.line_spacing = 1.35
                p2.space_after = Pt(14)
                
        add_footer(slide, page_num)
        return slide

    def add_section_intro(page_num, part_name, title_text, desc_text):
        """네오브루탈리즘 스타일의 다크 인트로 슬라이드를 생성합니다."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # 완전 블랙 배경
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BLACK
        bg.line.fill.background()
        
        # 강렬한 노란색/벽돌색 포인트 직각 도형 배치
        point = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.3))
        point.fill.solid()
        point.fill.fore_color.rgb = COLOR_YELLOW
        point.line.color.rgb = COLOR_BG
        point.line.width = Pt(3)
        
        tb = slide.shapes.add_textbox(Inches(2.0), Inches(2.3), Inches(9.333), Inches(3.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        p1 = tf.paragraphs[0]
        p1.text = part_name.upper()
        p1.font.name = FONT_TITLE
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_BRICK
        p1.space_after = Pt(8)
        
        p2 = tf.add_paragraph()
        p2.text = title_text
        p2.font.name = FONT_TITLE
        p2.font.size = Pt(36)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_BLACK
        
        p3 = tf.add_paragraph()
        p3.text = f"\n{desc_text}"
        p3.font.name = FONT_BODY
        p3.font.size = Pt(14)
        p3.font.bold = True
        p3.font.color.rgb = COLOR_PRIMARY_TEXT
        p3.line_spacing = 1.3
        
        # 하단 굵은 흰색 라인 및 페이지 번호
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(6.5), Inches(10.333), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_BG
        line.line.fill.background()
        
        page_box = slide.shapes.add_textbox(Inches(10.833), Inches(6.6), Inches(1.0), Inches(0.3))
        tf_page = page_box.text_frame
        p_page = tf_page.paragraphs[0]
        p_page.text = f"{page_num} / 30"
        p_page.alignment = PP_ALIGN.RIGHT
        p_page.font.name = FONT_TITLE
        p_page.font.size = Pt(11)
        p_page.font.bold = True
        p_page.font.color.rgb = COLOR_YELLOW
        
        return slide

    # ==============================================================================
    # 30개 슬라이드 네오브루탈리즘 버전 생성
    # ==============================================================================

    slide_layout = prs.slide_layouts[6]

    # --- Slide 1: 표지 ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    
    # 네오브루탈 데코: 큰 블랙 섀도 블록
    draw_neobrutalist_card(slide, Inches(0.8), Inches(1.8), Inches(11.733), Inches(3.8), fill_color=COLOR_YELLOW)
    
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(10.9), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Yes24 IT/컴퓨터 베스트셀러\n데이터 심층 EDA 보고서"
    p.font.name = FONT_TITLE
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_TEXT
    
    sub_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.0), Inches(10.9), Inches(1.0))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "데이터 기반 기술 실용서 시장 분석 및 비즈니스 마케팅 전략 제언"
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_BRICK
    
    info_box = slide.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(10.9), Inches(0.5))
    tf_info = info_box.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "작성일: 2026년 06월 06일  |  작성자: 20년차 수석 데이터 분석가  |  네오브루탈리즘 테마"
    p_info.font.name = FONT_BODY
    p_info.font.size = Pt(13)
    p_info.font.bold = True
    p_info.font.color.rgb = COLOR_PRIMARY_TEXT
    
    set_speaker_notes(slide, 
        "안녕하세요. 오늘 보고드릴 자료는 Yes24 IT/컴퓨터 베스트셀러 데이터 심층 EDA 보고서입니다. "
        "특별히 이번 프레젠테이션은 스타트업의Provocative한 에너지를 투영하는 '네오브루탈리즘' 스타일로 전면 리디자인되었습니다. "
        "두꺼운 검정 테두리와 하드그림자를 활용하여 레이아웃의 선명함을 확보했고, G마켓 산스 Bold와 나눔고딕으로 텍스트 가독성을 강조했습니다. "
        "베스트셀러 1,000건의 통계 마이닝을 바탕으로 실질적인 출판 마케팅 로드맵을 제안해 드리겠습니다."
    )

    # --- Slide 2: 목차 ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "목차 및 분석 로드맵")
    toc_titles = ["1부. 단변량 분석 및 기초 데이터 탐색", "2부. 이변량 & 다변량 연계 분석", "3부. 종합 비즈니스 제언"]
    toc_descs = [
        "판매가 분포, 도서 평점 만족도, 판매지수 형태, 분철 서비스 여부, 메이저 출판사 점유율 등 주요 단일 변수의 기초 통계와 시각화 데이터를 분석합니다.",
        "할인율 및 분철 서비스 제공 유무에 따른 판매 성과, 최근 출판 트렌드 시계열 분석, 다변량 상관관계 분석, TF-IDF 텍스트 마이닝을 통한 시장 수요 예측을 다룹니다.",
        "20년차 데이터 분석가의 관점으로 가격 포지셔닝 이원화 전략, 린(Lean) 퍼블리싱 모델, 소셜 증명 및 리뷰 평판 관리, 태깅 최적화 등의 액션 플랜을 제시합니다."
    ]
    toc_icons = ["book", "chart", "lightbulb"]
    for i in range(3):
        x = Inches(0.6 + i * 4.1)
        y = Inches(1.8)
        w = Inches(3.9)
        h = Inches(4.5)
        # 카드 배경을 옐로우와 화이트로 믹스하여 네오브루탈리즘 감성 연출
        draw_neobrutalist_card(slide, x, y, w, h, fill_color=COLOR_CARD_BG)
        draw_neobrutalist_icon(slide, toc_icons[i], x + Inches(0.4), y + Inches(0.4), Inches(0.55))
        t_box = slide.shapes.add_textbox(x + Inches(0.4), y + Inches(1.2), w - Inches(0.8), Inches(0.6))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = toc_titles[i]
        p.font.name = FONT_TITLE
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        
        d_box = slide.shapes.add_textbox(x + Inches(0.4), y + Inches(2.0), w - Inches(0.8), Inches(2.4))
        tf_d = d_box.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = toc_descs[i]
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_PRIMARY_TEXT
        p_d.line_spacing = 1.3
    add_footer(slide, 2)
    set_speaker_notes(slide,
        "목차 및 분석 로드맵 장표입니다. 두꺼운 검정 테두리의 정형화된 네오브루탈리즘 카드가 3개 나열되어 있습니다. "
        "이번 분석은 크게 1부인 단변량 기초 분석, 2부인 이변량 및 다변량 교차 연계 분석, "
        "그리고 3부인 종합 비즈니스 마케팅 전략 제언으로 구성되어 있습니다. "
        "각 카드의 요소를 확인해 보십시오."
    )

    # --- Slide 3: 데이터 소개 및 수집 개요 ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "데이터 세트 수집 개요 및 스키마 구조")
    draw_neobrutalist_card(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    left_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.0), Inches(4.1))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "데이터 세트 기본 스펙"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_BRICK
    specs = [
        ("수집 대상", "Yes24 IT/컴퓨터 카테고리 베스트셀러"),
        ("데이터 크기", "1,000개 행(Rows) / 30개 열(Columns)"),
        ("중복 데이터", "0건 (완전 제거 완료)"),
        ("정제 변수", "clean 처리된 수치형 및 범주형 재가공 변수 다수"),
        ("분석 일자", "2026년 06월 05일 기준 최신 트렌드 반영")
    ]
    for label, val in specs:
        p = tf.add_paragraph()
        p.text = f"\n• {label}:  "
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        run = p.add_run()
        run.text = val
        run.font.name = FONT_BODY
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = COLOR_PRIMARY_TEXT
        
    draw_neobrutalist_card(slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.5))
    right_box = slide.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(5.0), Inches(4.1))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "주요 핵심 변수 구분"
    p_r.font.name = FONT_TITLE
    p_r.font.size = Pt(20)
    p_r.font.bold = True
    p_r.font.color.rgb = COLOR_PRIMARY_TEXT
    vars_desc = [
        ("대상 변수", "도서명, 저자, 출판사, 발행 연/월, 태그"),
        ("성과 지표", "판매지수(Sale Index), 평점(Rating), 리뷰 수"),
        ("기능/가격", "판매가, 정가, 할인율, 분철 여부(Spring Service)"),
        ("파생 변수", "정제된 가격(clean), 태그 개수, 할인율 범주")
    ]
    for label, val in vars_desc:
        p = tf_r.add_paragraph()
        p.text = f"\n• {label}:  "
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        run = p.add_run()
        run.text = val
        run.font.name = FONT_BODY
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = COLOR_SECONDARY_TEXT
    add_footer(slide, 3)
    set_speaker_notes(slide,
        "세 번째 장표는 수집된 원본 데이터의 구조와 가공 목록을 해설한 카드형 슬라이드입니다. "
        "정렬된 격자 틀 속에 두꺼운 선들이 교차해 있어 네오브루탈리즘 특유의 날것 느낌이 납니다. "
        "Yes24에서 수집한 1,000건의 행과 30개의 수치형/범주형 속성들을 정밀 전처리하여 이 보고서에 투영시켰습니다."
    )

    # --- Slide 4: 요약 통계 (수치형 데이터) ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "전체 수치형 변수 통합 기술통계 분석")
    kpis = [
        ("평균 판매가", "23,480 원", "중위값: 22,500원 / 최소: 4,500원", "price"),
        ("평균 판매지수", "3,026 점", "중위값: 1,236점 / 최대: 87,480점", "trend"),
        ("평균 도서 평점", "7.50 점", "중위값: 9.70점 (0점 미평가 제외 시)", "star"),
        ("평균 도서 할인율", "8.64 %", "최빈값: 10% (도서정가제 상한 도달)", "percent")
    ]
    for idx, (label, val, subtext, icon) in enumerate(kpis):
        col = idx % 2
        row = idx // 2
        x = Inches(0.6 + col * 6.1)
        y = Inches(1.8 + row * 2.4)
        w = Inches(5.8)
        h = Inches(2.2)
        draw_neobrutalist_card(slide, x, y, w, h)
        draw_neobrutalist_icon(slide, icon, x + Inches(0.4), y + Inches(0.4), Inches(0.6))
        tb = slide.shapes.add_textbox(x + Inches(1.2), y + Inches(0.3), w - Inches(1.4), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.name = FONT_BODY
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_SECONDARY_TEXT
        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.name = FONT_TITLE
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_BRICK
        p3 = tf.add_paragraph()
        p3.text = subtext
        p3.font.name = FONT_BODY
        p3.font.size = Pt(11)
        p3.font.bold = True
        p3.font.color.rgb = COLOR_PRIMARY_TEXT
    add_footer(slide, 4)
    set_speaker_notes(slide,
        "네 번째 장표는 수치형 데이터들의 기본 기술통계를 2x2 카드 형태로 시각화한 화면입니다. "
        "각 카드마다 2.5pt의 굵직한 검정 선과 섀도가 입혀져 단단하고 임팩트 있는 가독성을 줍니다. "
        "평균 판매가 23,480원, 평균 판매지수 3,026점, 평균 평점 7.5점, 할인율 8.64%의 주요 지표들이 직관적으로 드러납니다."
    )

    # --- Slide 5: 요약 통계 (범주형 데이터) ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "전체 범주형 변수 요약 및 시장 쏠림 현상")
    draw_neobrutalist_card(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    tb_left = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.0), Inches(4.1))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "범주형 데이터 핵심 카운트"
    p_l.font.name = FONT_TITLE
    p_l.font.size = Pt(20)
    p_l.font.bold = True
    p_l.font.color.rgb = COLOR_PRIMARY_TEXT
    cat_items = [
        ("고유 도서 수", "1,000 권 (중복이 없는 유일 도서명)"),
        ("고유 저자 수", "869 명 (상위 스타 저자의 일부 다작 존재)"),
        ("등록 출판사 수", "187 개사 (과반의 독점 현상 관찰)"),
        ("분철 미지원 도서", "853 권 (전체의 85.3% 점유)"),
        ("분철 지원 도서", "147 권 (전체의 14.7% 점유)")
    ]
    for label, val in cat_items:
        p = tf_l.add_paragraph()
        p.text = f"\n• {label}:  "
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        run = p.add_run()
        run.text = val
        run.font.name = FONT_BODY
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = COLOR_SECONDARY_TEXT
        
    draw_neobrutalist_card(slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.5))
    tb_right = slide.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(5.0), Inches(4.1))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "범주형 데이터 시사점"
    p_r.font.name = FONT_TITLE
    p_r.font.size = Pt(20)
    p_r.font.bold = True
    p_r.font.color.rgb = COLOR_BRICK
    insights_l = [
        "1. 저자 집중도 완화: 1,000개 베스트셀러 중 고유 저자가 869명에 달해 특정 스타 저자가 전체를 독식하기보다는 다양한 도서가 순위에 진입합니다.",
        "2. 출판사 양극화 심각: 반면 출판사는 187개사이지만 한빛미디어 등 소수 대형 출판사가 상위권을 지배하는 독점 성향을 드러냅니다.",
        "3. 기능성 편의 격차: 분철 서비스가 제공되는 서적은 14.7%에 불과하여 독자들의 잠재적 니즈 대비 여전히 공급이 희소한 틈새 시장입니다."
    ]
    for ins in insights_l:
        p = tf_r.add_paragraph()
        p.text = f"\n{ins}"
        p.font.name = FONT_BODY
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY_TEXT
        p.line_spacing = 1.3
    add_footer(slide, 5)
    set_speaker_notes(slide,
        "다섯 번째 슬라이드는 범주형 변수의 분석입니다. "
        "고유 도서와 고유 저자 수가 많음에도 출판사 수는 187개로 좁아지며, 분철 도서가 14.7%에 미치지 못한다는 사실을 카드형 레이아웃에 거칠고 명확하게 명시하고 있습니다."
    )

    # --- Slide 6: PART 1 인트로 ---
    add_section_intro(6, "PART 1", "단변량 분석 및 핵심 분포 탐색", 
                      "판매가, 도서 평점 만족도, 판매지수 형태, 분철 서비스 여부, 메이저 출판사 분포 등\n주요 개별 지표들의 통계적 분포와 비즈니스 기초 인사이트를 도출합니다.")

    # --- Slide 7: 판매가 빈도 분포 ---
    add_standard_chart_slide(
        page_num=7,
        title="1.1 판매가(Sale Price) 빈도 분포 분석",
        chart_filename="01_sale_price_distribution.png",
        stat_title="판매가 구간별 기초 통계",
        stats_list=[
            ("집중 구간", "20,000원 ~ 30,000원 대역 과밀 분포"),
            ("평균 판매가", "23,480 원"),
            ("중위 판매가", "22,500 원"),
            ("최소/최대가", "4,500원 / 67,000원")
        ],
        insight_text="Yes24 IT 베스트셀러 도서들은 약 2만 원에서 3만 원 사이 대역에 70% 가까운 수량이 밀집되어 있습니다. 이는 독자들이 기술 서적을 구매할 때 심리적으로 수용하기 가장 편한 표준 가격 대역대가 2만 원대 초중반에 견고하게 차단되어 있음을 방증합니다.",
        icon_name="price"
    )
    set_speaker_notes(prs.slides[-1],
        "일곱 번째 슬라이드는 도서 판매가의 빈도 분포 차트 분석입니다. "
        "네오브루탈리즘 프레임에 맞춰 좌측 차트와 우측 카드 모두 두껍고 곧은 테두리 및 그림자가 적용되어 강렬한 인상을 줍니다. "
        "대다수의 IT 도서 가격이 2만 원대 중후반에 위치하고 있습니다."
    )

    # --- Slide 8: 판매가 심층 분석 ---
    add_deep_insight_slide(
        page_num=8,
        title="1.1 [심층분석] 가격 심리적 저항선과 세그먼트 전략",
        left_title="심리적 마지노선 분석",
        left_points=[
            ("3만원 가격 장벽", "판매가가 30,000원 선을 돌파할 때, 독자가 즉각 결제하기 주저하는 심리적 이탈 저항선이 뚜렷이 나타납니다."),
            ("보급형 입문서 설계", "단기적 베스트셀러 양적 성장을 도모하는 입문서/문법서는 최종 할인가가 22,000원~26,000원 구간에 걸치도록 타겟 정가를 기획하는 것이 최선입니다.")
        ],
        right_title="고가 프리미엄 세그먼트 전략",
        right_points=[
            ("전문성 희소가치 수용", "RAG 아키텍처 구축이나 LLM 미세 조정 등 난이도가 높고 시의성 강한 독점적 전문서라면, 4만 원 이상으로 가격을 높여도 구매로 이어집니다."),
            ("가치 대비 가격 책정", "기술 대체 불가능성이 높다면 일괄 가격 타협보다는 분량을 대폭 보강하고 사은 혜택을 얹어 고가 프리미엄 전략을 펴는 편이 마진 방어에 유리합니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "여덟 번째 슬라이드는 가격 분포의 입체적 심층 해석입니다. "
        "독자들은 3만 원 저항선과 전문서 고가 결제 용의라는 양극의 가치를 갖고 있습니다. "
        "이에 따라 기획 가격을 이원화해야 함을 설득력 있게 제언합니다."
    )

    # --- Slide 9: 도서 평점 분포 ---
    add_standard_chart_slide(
        page_num=9,
        title="1.2 도서 평점(Rating) 빈도 분포 분석",
        chart_filename="02_rating_distribution.png",
        stat_title="평점 분포 수치 요약",
        stats_list=[
            ("만점 밀집도", "9.5 ~ 10.0 점 대역에 극단적인 우측 쏠림"),
            ("전체 평균 평점", "7.50 점 (0점 무평가 도서 혼재 효과)"),
            ("평가도서 중위수", "9.70 점 (독자 만족도 전반적 우수)"),
            ("최소/최대 평점", "0.00 점 / 10.00 점")
        ],
        insight_text="도서 평점의 히스토그램을 분석하면 대다수 평점이 만점에 밀집되어 일반 평점 점수 자체는 도서의 우열을 가리는 실질적 잣대로서의 변별력을 완전히 소실했습니다. 평균 7.5점은 리뷰가 0건인 도서들의 영점 효과로 발생한 굴절입니다.",
        icon_name="star"
    )
    set_speaker_notes(prs.slides[-1],
        "아홉 번째 슬라이드는 도서 평점의 분포 차트입니다. "
        "대부분의 책들이 만점에 가까운 9.5~10점 사이에 모여 있고, 미평가 도서들의 0점 효과가 뚜렷이 관측됩니다."
    )

    # --- Slide 10: 평점 심층 분석 ---
    add_deep_insight_slide(
        page_num=10,
        title="1.2 [심층분석] 평점 인플레이션 극복 및 0점 돌파 전략",
        left_title="평점 인플레이션의 대용 지표",
        left_points=[
            ("사회적 증거(Social Proof) 작동", "독자들은 평점 10점 만점을 신임하기보다, 평점 9.6점이라도 상세 텍스트 리뷰가 100개 넘게 누적되어 있는 책을 강력히 의지합니다."),
            ("리뷰 수로 무게중심 이동", "따라서 평점 스펙 높이기보다는, 절대적인 평가의 볼륨(리뷰 건수)을 띄우는 것이 잠재 독자를 구매 전환시키는 소셜 락인입니다.")
        ],
        right_title="신간 도서 0점 신속 탈출 지침",
        right_points=[
            ("0점 도서의 치명적 노출 방해", "0점 도서는 불량이 아니라 미검증 낙인 효과를 주어 독자의 장바구니 결제를 심리적으로 차단하는 부작용을 유발합니다."),
            ("초기 런칭 1개월 마케팅", "도서 출간 즉시 초기 가용 마케팅 화력을 서평단 유치와 독자 리뷰 마일리지 보상에 집중해, 최소 10건 이상의 초기 서평을 신속 확보해야 합니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "열 번째 슬라이드는 평점 인플레이션을 극복하는 마케팅 실천 지침입니다. "
        "의사결정의 무게추는 절대 평점이 아니라 누적 리뷰 건수입니다. "
        "신간 출시 시 0점 노출에 갇히지 않도록 초기 1개월 골든타임 내에 가시적 리뷰 자산을 빌드업해야 합니다."
    )

    # --- Slide 11: 판매지수 분포 ---
    add_standard_chart_slide(
        page_num=11,
        title="1.3 판매지수(Sale Index) 상자 그림 분포 분석",
        chart_filename="03_sale_index_boxplot.png",
        stat_title="판매지수 통계 및 사분위",
        stats_list=[
            ("중위 판매지수", "1,236 점 (현실적 성과 척도)"),
            ("평균 판매지수", "3,026 점 (상위 아웃라이어 왜곡 효과)"),
            ("최대 판매지수", "87,480 점 (독보적 메가 히트 타이틀)"),
            ("표준편차(std)", "7,077 점 (매우 큰 분포 변동성)")
        ],
        insight_text="판매지수의 박스플롯을 보면 대다수 베스트셀러가 5,000점 이하 대역에 밀집된 반면, 상위 킬러 타이틀들은 최고 8.7만 점에 육박합니다. 베스트셀러 랭킹 내에서도 엄청난 판매 점유 양극화가 일어나는 전형적인 롱테일 구조입니다.",
        icon_name="trend"
    )
    set_speaker_notes(prs.slides[-1],
        "열한 번째 장표는 판매지수 상자 그림 분석입니다. "
        "차트에서 보이듯 거의 모든 도서가 5,000점 이하의 바닥권에 수렴해 있으나 상위 점들은 우주 위로 흩어져 있습니다. "
        "소수의 메가 히트작이 시장 총매출을 리드하고 있는 과점 생태계입니다."
    )

    # --- Slide 12: 판매지수 심층 분석 ---
    add_deep_insight_slide(
        page_num=12,
        title="1.3 [심층분석] 롱테일 시장 구조와 현실적 KPI 수립",
        left_title="도서 시장의 파레토 법칙",
        left_points=[
            ("소수 킬러 타이틀 독점", "베스트셀러 1,000권 목록에서도 상위 20%의 서적이 총 판매지수의 압도적 지분을 장악하며 유통 플랫폼의 헤드 역할을 담당합니다."),
            ("평균 왜곡의 함정", "출판 기획 시 단순 평균값(3,026)을 목표 판매지수로 삼으면 과도한 마케팅 비용 태우기나 예산 편성 불균형을 야기할 수 있습니다.")
        ],
        right_title="중위수(Median) 벤치마킹 액션",
        right_points=[
            ("1차 마일스톤 설정", "현실적인 신작의 성공 안착 기준은 50% 중위선인 '판매지수 1,236점'을 달성하는 것으로 잡고, 2차 확산 전략을 모색하는 것이 리스크를 관리하는 지름길입니다."),
            ("업데이트 주기적 투자", "상위 극소수 아웃라이어 타이틀은 1회성 마케팅 결과가 아닙니다. 깃허브 피드백을 수용한 신속한 개정판 관리로 탄생시킨 지속적 업그레이드의 산물입니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "열두 번째 장표는 판매지수의 왜곡을 고려한 신작 경영 가이드라인입니다. "
        "신간의 현실적인 마일스톤 목표를 중위수인 1,236점으로 재배정하여 "
        "출판사의 합리적인 손익분기 관리를 유도하는 분석가 의견입니다."
    )

    # --- Slide 13: 분철 서비스 빈도 ---
    add_standard_chart_slide(
        page_num=13,
        title="1.4 분철 서비스 제공 여부(Spring Service) 분석",
        chart_filename="04_spring_service_count.png",
        stat_title="분철 제공 수량 및 점유비",
        stats_list=[
            ("분철 미제공 도서(N)", "853 권 (85.3% 압도적 지분)"),
            ("분철 제공 도서(Y)", "147 권 (14.7% 소수 공급)"),
            ("주요 적용 도서 특성", "수험서, 개발 실무 매뉴얼, 두꺼운 대학 전공서")
        ],
        insight_text="Yes24 베스트셀러 내 분철 지원 서적은 14.7%에 불과합니다. 제작 공정 단가와 물류 유통 편의상 대다수 출판사가 분철을 지양하고 있으나, 수험서나 실무 전공 서적을 펼쳐놓고 학습하는 독자들에게는 잠재적 수요가 거대하게 깔려 있는 영역입니다.",
        icon_name="spring"
    )
    set_speaker_notes(prs.slides[-1],
        "열세 번째 장표는 분철 서비스 제공 비율에 대한 차트 분석입니다. "
        "대다수 도서가 일반 단행본 제본이지만, 두께가 두껍고 펼쳐놓고 학습해야 하는 14.7%의 특수한 실용서에 한해 제공되고 있습니다."
    )

    # --- Slide 14: 상위 30개 출판사별 도서 수 ---
    add_standard_chart_slide(
        page_num=14,
        title="1.5 상위 30개 출판사별 베스트셀러 점유율",
        chart_filename="05_top_30_publishers.png",
        stat_title="상위 출판사 점유 통계",
        stats_list=[
            ("압도적 1위", "한빛미디어 (150 권, 15.0% 독점)"),
            ("2위 ~ 4위권", "길벗(7.4%), 제이펍(5.1%), 이지스퍼블리싱(5.0%)"),
            ("5위 신흥 주자", "골든래빗 (43 권, 4.3% 가파른 상승)"),
            ("상위 10개사 점유율", "과반(약 50.2%) 돌파 (나머지 177개사 분산)")
        ],
        insight_text="국내 IT 베스트셀러 시장은 대형 메이저 브랜드의 점유 집중도가 매우 극심합니다. 상위 10개 퍼블리셔가 전체 시장의 과반을 점유하고 있어, 신진 브랜드가 독자적인 힘으로 랭킹에 오르기란 매우 가파른 진입 장벽이 가로막고 있습니다.",
        icon_name="building"
    )
    set_speaker_notes(prs.slides[-1],
        "열네 번째 슬라이드는 출판사 점유 과점 현황입니다. "
        "한빛미디어가 15%의 도서 점유율로 독보적인 1위이며, 상위 10개사가 시장의 과반을 차지합니다. "
        "신규 출판사들에게는 매우 강력한 진입 장벽으로 작용하고 있습니다."
    )

    # --- Slide 15: 출판사 심층 분석 ---
    add_deep_insight_slide(
        page_num=15,
        title="1.5 [심층분석] 메이저 출판사 브랜드 락인과 상생 전략",
        left_title="대형 퍼블리셔의 후광 효과",
        left_points=[
            ("베타 독자단 인프라", "한빛이나 길벗 등은 수만 명에 달하는 IT 전문 베타 리더단 풀과 충성도 높은 기술 뉴스레터 구독자 락인 체계를 가지고 있습니다."),
            ("기생/협력 투고 전략", "무명 저자나 신인 필자는 마케팅 인프라가 전무하므로, 대형 출판사의 기획 투고를 통과해 그들의 거대한 브랜드 후광을 타는 것이 성공 기회를 열어줍니다.")
        ],
        right_title="틈새 린(Lean) 론칭 전략",
        right_points=[
            ("대형사의 느린 속도 틈새", "대형 퍼블리셔들은 내부 결재와 편집 공정이 다소 길어, 신생 프레임워크나 최신 AI API의 마이크로 업데이트에 민첩하게 반응하기 힘듭니다."),
            ("마이크로 커뮤니티 타겟", "독립 출판이나 소형 기획사는 초단기 트렌드 원고를 선점하고, 개발자 단톡방이나 디스코드, 깃허브 서포트를 활용해 기동성 있게 시장을 우회 침투해야 승산이 있습니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "열다섯 번째 장표는 메이저 과점 구도에 대한 상생 및 우회 기획안입니다. "
        "대형사의 거대한 마케팅 리소스를 이기기 힘드므로 파트너십 투고 전략을 우선시하되, "
        "만약 소형 출판사가 독자 생존하고자 한다면 대형사보다 민첩하게 최신 AI 틈새 트렌드를 포착해 슬림북 형태로 신속 출간해야 합니다."
    )

    # --- Slide 16: PART 2 인트로 ---
    add_section_intro(16, "PART 2", "이변량 및 다변량 연계 분석", 
                      "할인율 및 분철 서비스와 판매지수 간의 교차 영향력 검증,\n최근 출판 트렌드 시계열 분석, 다변량 피어슨 상관관계, TF-IDF 텍스트 마이닝을 통한\nIT 도서 시장의 실질적인 수요 동력(Driver)을 도출합니다.")

    # --- Slide 17: 할인율 vs 판매지수 ---
    add_standard_chart_slide(
        page_num=17,
        title="2.1 할인율(Discount Rate) vs 판매지수 산점도",
        chart_filename="06_discount_vs_sale_index.png",
        stat_title="할인 정책별 세부 지표",
        stats_list=[
            ("10% 최대 할인 적용 도서 수", "851 권  |  평균 판매지수: 3,394 점"),
            ("0% 무할인(정가) 적용 도서 수", "123 권  |  평균 판매지수: 826 점"),
            ("5% 부분 할인 적용 도서 수", "25 권  |  평균 판매지수: 1,276 점"),
            ("도서정가제에 따른 가격 통제", "10% 할인이 실질적 판매 안착의 필수적 기본 조건화")
        ],
        insight_text="거의 모든 베스트셀러 도서가 최대 상한선인 10% 할인을 깔고 유통되고 있습니다. 정가를 지키는 무할인 도서는 판매지수가 4분의 1 토막에 그치므로, 10% 할인은 선택이 아닌 필수입니다. 추가 경쟁력은 금전적 마찰 회피가 아닌 비금전적 가치 추가에서 발굴해야 합니다.",
        icon_name="percent"
    )
    set_speaker_notes(prs.slides[-1],
        "열일곱 번째 장표는 할인율과 판매지수 간의 산점도 분석입니다. "
        "정가 도서는 평균 지수가 현저히 처지며, 최대치인 10% 할인을 제공하는 도서가 압도적입니다. "
        "즉 가격 깎기는 기본일 뿐, 추가 차별화는 기능 편의에서 모색해야 합니다."
    )

    # --- Slide 18: 분철 서비스 여부 vs 판매지수 ---
    add_standard_chart_slide(
        page_num=18,
        title="2.2 분철 여부에 따른 판매 성과 교차 분석",
        chart_filename="07_spring_service_vs_sale_index.png",
        stat_title="분철 여부별 판매지수 비교",
        stats_list=[
            ("분철 지원 그룹 (Y)", "147 권  |  평균 판매지수: 8,687 점 (최대: 87,480)"),
            ("분철 미지원 그룹 (N)", "853 권  |  평균 판매지수: 2,050 점 (최대: 71,388)"),
            ("평균 격차 배율", "약 4.2배 수준의 평균 판매지수 초과 달성"),
            ("중위 지수 비교", "Y 그룹 중위수: 4,062점  |  N 그룹 중위수: 1,047점")
        ],
        insight_text="분철 서비스를 옵션으로 장착한 도서들의 평균 판매 성과가 미지원 도서 대비 무려 4.2배가량 압도적으로 높게 도출되었습니다. 이는 단편적인 제작비 증가분 이상의 판매 견인 및 락인 효과가 작용함을 통계적 증거로 규명하는 극적인 순간입니다.",
        icon_name="spring"
    )
    set_speaker_notes(prs.slides[-1],
        "열여덟 번째 장표는 분철 여부별 판매 성과 교차 분석입니다. "
        "평균치와 중위수 모두 분철 가능한 도서군이 4배 가까이 앞서 나갑니다. "
        "링 제본이 선사하는 실습 효율이 실용 독자의 장바구니 선택을 결정짓는 주 동인임을 통계로 증명합니다."
    )

    # --- Slide 19: 분철 서비스 심층 분석 ---
    add_deep_insight_slide(
        page_num=19,
        title="2.2 [심층분석] 분철 가치와 출판 유통사 락인 설계",
        left_title="독자 실용 소비주의 이해",
        left_points=[
            ("핸즈프리(Hands-free) 학습 니즈", "IT 독자들은 책을 모니터 옆에 완전히 젖혀두고 키보드를 자유롭게 조작하고 싶어 합니다. 일반 제본은 자꾸 책이 접혀 지속적인 피로를 가중시킵니다."),
            ("가치 체감 대비 저렴한 투자", "독자는 1,000원~2,000원의 소액 추가금을 주더라도 책의 활용성을 높일 수 있다면 흔쾌히 구매 옵션을 선택하는 패턴을 가집니다.")
        ],
        right_title="출판-유통망 연계 마케팅",
        right_points=[
            ("출판 유통의 기본 탑재", "500페이지가 넘는 기술 서적의 인쇄 기획 시점부터 출판사는 예스24의 주문 연계 분철 라인업을 상설 채널로 확보하고 전면 홍보해야 합니다."),
            ("구매 상세 페이지 노출 강화", "상세 페이지 최상단 배너에 '독자 만족 스프링 분철 옵션 완비' 표시를 강조하여, 분철을 미지원하는 타사 도서와의 결정적 구매 우위를 확립합니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "열아홉 번째 장표는 분철 성과에 기반한 실무 마케팅 액션안입니다. "
        "독자들의 손이 자유로워지는 링 제본은 실습 피로를 줄여주는 차별화 무기입니다. "
        "온라인 몰 상세 페이지 전면에 '분철 완비'를 기획 노출해 경쟁 도서와의 기능적 격차를 유도해야 합니다."
    )

    # --- Slide 20: 출판년도별 추이 ---
    add_standard_chart_slide(
        page_num=20,
        title="2.3 출판년도별 도서 수 및 평균 판매지수 추이",
        chart_filename="08_publish_year_trends.png",
        stat_title="출판 연도별 성과 통계",
        stats_list=[
            ("2025년 출간 도서", "367 권  |  평균 판매지수: 3,902 점"),
            ("2026년 출간 도서", "356 권  |  평균 판매지수: 2,404 점"),
            ("2024년 출간 도서", "138 권  |  평균 판매지수: 1,747 점"),
            ("2023년 이전 도서", "총 139 권  |  평균 지수 3,000점 이하로 급감 흐름")
        ],
        insight_text="2025~2026년 신간들이 베스트셀러 진입 권수와 판매지수의 대부분을 장악하고 있습니다. 연도가 지날수록 도서 수가 급감하여 IT 분야 서적의 감가상각과 트렌드 이탈 속도가 초고속으로 전개되고 있음을 보여줍니다.",
        icon_name="calendar"
    )
    set_speaker_notes(prs.slides[-1],
        "스무 번째 슬라이드는 출판 연도별 성과 분석입니다. "
        "네오브루탈리즘의 액자 카드가 차트 이미지를 견고히 붙들고 있어 선명한 시각 밸런스를 줍니다. "
        "신간이 베스트셀러 진입 권수와 판매지수 성장을 압도적으로 주도하는 지식 감가상각 현상입니다."
    )

    # --- Slide 21: 출판년도 심층 분석 ---
    add_deep_insight_slide(
        page_num=21,
        title="2.3 [심층분석] 시계열 지식 감가상각과 린 출판 대응",
        left_title="IT 도서 지식의 급격한 유통기한",
        left_points=[
            ("오픈소스 버전업 리스크", "라이브러리 마이너 버전 패치 하나만으로도 실습 예제 코드가 실행 에러를 뿜게 되며, 이는 곧장 소비자 평판 추락으로 이어집니다."),
            ("기존 기획 관행 탈피", "과거의 1년짜리 고전적 기획-집필-출판 루프에서 벗어나지 않으면 출간하는 당일 이미 구닥다리 기술서가 될 위험이 농후합니다.")
        ],
        right_title="린 퍼블리싱(Lean Publishing) 도입",
        right_points=[
            ("실시간 독자 피드백 결합", "원고 기획서 단계에서부터 온라인 깃허브 저장소를 개설하고, 독자들과 온라인에서 집필 원고를 소통하며 완성 시점에 신속히 찍어내야 합니다."),
            ("콘텐츠 두께 슬림화 전략", "트렌드가 불을 뿜는 프론트엔드나 생성형 AI 기술 서적은 분량을 200페이지 내외로 축소하여 런칭 주기를 최대한 민첩하게 회전해야 합니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "스물한 번째 장표는 지식 감가상각을 우회하는 린 출판 프로세스 제언입니다. "
        "예제 코드의 신속한 오류 제어를 위해 오픈소스로 소통하며 집필하는 구조를 세팅하고, "
        "두꺼운 바이블 서적 대신 얇고 트렌디한 슬림북 중심의 빠른 개정판 로테이션을 구축해야 합니다."
    )

    # --- Slide 22: 상관관계 히트맵 ---
    add_standard_chart_slide(
        page_num=22,
        title="2.4 수치형 변수 간 상관관계 피어슨 상관행렬",
        chart_filename="09_correlation_heatmap.png",
        stat_title="피어슨 상관계수(r) 결과",
        stats_list=[
            ("판매지수 vs 리뷰 수", "r = 0.214 (유의미한 정적 상관성 성립)"),
            ("판매지수 vs 평점", "r = 0.156 (약한 정적 상관성)"),
            ("판매지수 vs 할인율", "r = 0.124 (미약한 정적 상관성)"),
            ("할인율 vs 포인트적립", "r = 0.658 (가격 마케팅 변수 간 상호 수렴)")
        ],
        insight_text="상관관계 분석 결과, 판매지수와 가장 높은 인과적 연계를 보이는 변수는 '리뷰 수'로 규명되었습니다. 평점 수치 자체보다 리뷰의 절대 수량이 온라인 구매 채널에서 소비자의 안도감을 유도하는 가장 큰 유발 요인입니다.",
        icon_name="chart"
    )
    set_speaker_notes(prs.slides[-1],
        "스물두 번째 장표는 수치형 지표들의 피어슨 상관관계 상관행렬입니다. "
        "판매지수와 가장 강력하게 연동되는 핵심 변수는 단연 리뷰 수(r=0.214)입니다. "
        "반면 판매가격은 매출 지수 성과와 아무런 통계적 연관이 관측되지 않습니다."
    )

    # --- Slide 23: 상관관계 심층 분석 ---
    add_deep_insight_slide(
        page_num=23,
        title="2.4 [심층분석] 리뷰 자산과 소셜 증명(Social Proof)의 선순환",
        left_title="별점보다 두터운 리뷰 볼륨",
        left_points=[
            ("학습 삽질 비용 회피", "IT 독자가 지불하는 가장 큰 비용은 책값이 아니라 내 주말 시간입니다. 독자는 실패 리스크를 줄이기 위해 풍부한 리뷰 텍스트를 찾습니다."),
            ("텍스트 피드백의 힘", "구체적인 실습 삽질기나 문제 해결 방안이 기록된 리뷰가 많을수록 제품 만족에 대한 확신을 심어주어 구매 결정을 촉진합니다.")
        ],
        right_title="런칭 초기의 리뷰 빌드업 지침",
        right_points=[
            ("초기 30일 골든타임 관리", "출간 1개월 이내에 정성스러운 텍스트 및 사진 리뷰가 플랫폼 메인에 노출될 수 있도록 독자 서평단을 정교하게 스케줄링해야 합니다."),
            ("인센티브 선순환 구조", "도서 예제 코드 깃허브나 책 뒤편에 '리뷰 작성 시 추가 소스 제공 및 마일리지 프로모션' 링크를 연계해 구매자의 자발적 평판 기여를 끊임없이 조장해야 합니다.")
        ]
    )
    set_speaker_notes(prs.slides[-1],
        "스물세 번째 장표는 독자들의 삽질 피로 회피 본능을 겨냥한 리뷰 자산 구축 전략입니다. "
        "독자는 학습 피로를 피해 가고자 타인의 풍부한 리뷰를 소셜 프루프로 삼습니다. "
        "출간 초기 1개월 골든타임 내에 서평을 집중 런칭하고, 자발적 인증 후기 인센티브를 연계해야 합니다."
    )

    # --- Slide 24: 할인율-분철 여부별 판매지수 ---
    add_standard_chart_slide(
        page_num=24,
        title="2.5 [다변량] 할인율-분철 여부별 평균 판매지수",
        chart_filename="10_discount_spring_vs_sale_index.png",
        stat_title="할인 x 분철 다변량 피봇 결과",
        stats_list=[
            ("10% 할인 + 분철지원 (Y)", "평균 판매지수: 9,131 점 (최상의 성과 구간)"),
            ("10% 할인 + 분철미지원 (N)", "평균 판매지수: 2,321 점 (보편적 평균 구간)"),
            ("0% 무할인 + 분철지원 (Y)", "평균 판매지수: 3,653 점 (할인 N 그룹 초과)"),
            ("0% 무할인 + 분철미지원 (N)", "평균 판매지수: 603 점 (최하의 성과 구간)")
        ],
        insight_text="10% 기본 할인에 분철 서비스를 결합한 그룹이 압도적 1위의 판매지수(9,131)를 냅니다. 흥미롭게도 할인을 0%로 고수하더라도 분철을 장착하기만 하면 평균 3,653점으로 분철이 없는 할인 도서 집단보다 우수한 매출 성적을 거두었습니다.",
        icon_name="percent"
    )
    set_speaker_notes(prs.slides[-1],
        "스물네 번째 장표는 할인과 분철을 결합한 다변량 성과 비교입니다. "
        "우측 카드의 옐로우 아이콘 블록과 텍스트 영역을 수직으로 엄격히 분리하여 겹침 없이 아주 깔끔한 디자인 밸런스를 구축했습니다. "
        "분철이 제공된다면 무할인이라도 10% 일반 할인 서적의 성과를 가볍게 우회 압도합니다."
    )

    # --- Slide 25: 도서명 TF-IDF 중요 키워드 ---
    add_standard_chart_slide(
        page_num=25,
        title="2.6 도서명 기준 TF-IDF 중요 키워드 분석",
        chart_filename="11_goods_name_tfidf.png",
        stat_title="상위 TF-IDF 가중치 단어",
        stats_list=[
            ("핵심 기술 테마", "'AI' (가중치 0.059), '파이썬' (0.019), '인공지능' (0.012)"),
            ("인공지능 보조 도구", "'제미나이' (0.010), '클로드' (0.009), '챗GPT' (0.012)"),
            ("학습 지향 키워드", "'코딩' (0.017), '가이드' (0.016), '입문' (0.011)"),
            ("생산성 향상 키워드", "'활용' (0.011), '실무' (0.009), '실전' (0.009)")
        ],
        insight_text="도서 제목의 텍스트 마이닝 결과, 단순한 프로그래밍 기본 이론에서 벗어나 AI 도구(클로드, 제미나이 등)를 활용한 코딩 자동화와 실용적인 파이썬 기반 데이터 과학 관련 주제가 가중치 상위를 석권하고 있습니다.",
        icon_name="key"
    )
    set_speaker_notes(prs.slides[-1],
        "스물오분째 장표는 도서명 키워드 가중치 마이닝 분석입니다. "
        "시장의 실질 수요는 이미 AI 보조 개발 도구를 이용한 실무 생산성 향상과 파이썬 데이터 과학으로 완전히 전환되었습니다. "
        "출판 기획 시 이 융합 주제를 반드시 접목해야 흥행을 담보할 수 있습니다."
    )

    # --- Slide 26: 저자별 평균 판매지수 ---
    add_standard_chart_slide(
        page_num=26,
        title="2.7 상위 저자별 도서 수 및 평균 판매지수 분석",
        chart_filename="12_author_mean_sale_index.png",
        stat_title="저자 세그먼트 성적표",
        stats_list=[
            ("최다 도서 등록 저자군", "놀이교육콘텐츠랩(9권, 평균 1,170점)"),
            ("최다 등록 출판사 서적", "해람북스 기획팀(5권, 평균 1,213점)"),
            ("소수 스타 저자 쏠림", "오힘찬 저자(3권, 평균 36,848점)"),
            ("다작 vs 고효율 전략", "도서 등록 권수와 평균 판매지수 간 반비례 양상")
        ],
        insight_text="단순 다작 저자들보다, 최신 기술 흐름(예: 제미나이 마스터북)을 명확하게 파헤친 스타 저자의 도서가 압도적인 평균 판매지수 스파이크를 형성합니다. 독자들은 저자의 다작 여부보다 화제성과 실용성에 충성합니다.",
        icon_name="users"
    )
    set_speaker_notes(prs.slides[-1],
        "스물여섯 번째 장표는 상위 저자별 도서 수 및 성과 분석입니다. "
        "다작 저자들의 지수가 비교적 완만한 반면, 최신 생성 AI 트렌드를 저격해 명쾌하게 풀어낸 스타 저자의 평균 판매지수가 수십 배에 달합니다. "
        "양적 다작보다 킬러 타이틀 1권을 발매하는 편이 절대적으로 효율적입니다."
    )

    # --- Slide 27: 출판 월별 추이 ---
    add_standard_chart_slide(
        page_num=27,
        title="2.8 출판 월별 신간 수 및 평균 판매지수 분석",
        chart_filename="13_publish_month_trends.png",
        stat_title="출판 월별 시즌 성과",
        stats_list=[
            ("상반기(1~6월) 런칭 성과", "신규 진입 도서 집중 및 높은 평균 판매지수 유지"),
            ("최대 판매지수 월", "12월 (평균 3,985점), 9월 (평균 3,743점)"),
            ("최저 판매지수 월", "7월 (평균 2,006점), 10월 (평균 2,343점)"),
            ("새해/새학기 결심 효과", "방학과 신학기 직전에 맞춘 도서 소비 본능 확인")
        ],
        insight_text="상반기에 발행된 도서들의 판매지수 강세가 도드라집니다. 연말연시나 새학기 결심(새해 다짐, 방학 공부 계획)에 부합하여 소비자들이 지갑을 여는 시즌성 효과입니다. 론칭 일정을 겨울방학이나 신학기 직전에 포커싱해야 합니다.",
        icon_name="calendar"
    )
    set_speaker_notes(prs.slides[-1],
        "스물일곱 번째 장표는 출판월별 시즌 패턴 분석입니다. "
        "IT 도서 구매는 새해 다짐이나 새학기 방학 공부 등 독자의 결심 시즌에 집중되어 있습니다. "
        "따라서 12월이나 2월 직전에 신간 런칭 주기를 셋백하는 지혜가 필요합니다."
    )

    # --- Slide 28: 도서 태그 수 vs 판매지수 ---
    add_standard_chart_slide(
        page_num=28,
        title="2.9 도서 태그 수(Tag Count)와 판매지수 분석",
        chart_filename="14_tag_count_vs_sale_index.png",
        stat_title="태그 수에 따른 성과",
        stats_list=[
            ("적정 태그 수 대역", "2개 ~ 4개 등록 도서의 판매지수 안정적 우위"),
            ("태그 0개 등록 도서 수", "685 권  |  평균 판매지수: 1,815 점"),
            ("태그 5개 등록 도서 수", "47 권  |  평균 판매지수: 8,770 점 (고가도서 쏠림)"),
            ("검색 엔진 최적화(SEO)", "부적합한 태그 남발보다 정교한 타겟 키워드 매핑 권장")
        ],
        insight_text="태그가 아예 없는(0개) 서적들은 노출 기회 박탈로 성적이 부진합니다. 반면 2~4개의 정밀하고 세련된 태그를 매핑해 둔 도서군이 안정적인 지수 흐름을 냅니다. 검색 노출을 위한 최적의 태그 관리가 요구됩니다.",
        icon_name="tag"
    )
    set_speaker_notes(prs.slides[-1],
        "스물여덟 번째 장표는 태그 수와 판매지수 간의 연계입니다. "
        "검색 노출에서 배제되는 0~1개 태그를 지양하고, 노이즈 없는 카테고리 태그 3~4개를 선별 등록하여 "
        "온라인 플랫폼 검색 최적화를 이루어내야 합니다."
    )

    # --- Slide 29: 3부 종합 비즈니스 제언 요약 ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "PART 3. IT 도서 시장 활성화를 위한 5대 핵심 성공 전략", "BUSINESS STRATEGY & ROADMAP")
    strategies = [
        ("1. 가격 포지셔닝 이원화", "입문 기술서는 2.3만~2.7만원 대역 책정으로 가격 저항 최소화, 특수 고난도 기술서는 4만원 대 프리미엄 책정.", "price"),
        ("2. 분철 서비스 연계 활성화", "500p 이상 및 실무 코드가 중심인 도서는 제본 옵션을 기본 제공하여 26% 이상의 추가 매출 동력 확보.", "spring"),
        ("3. 지식 감가상각 방어 프로세스", "트렌드가 민감한 응용서는 린(Lean) 집필 공정을 세팅하여 빠르게 출판하고 수명을 선제적 관리.", "calendar"),
        ("4. 독자 평판 소셜프루프 극대화", "초기 서평단 1개월 빌드업 전략에 리소스를 집중하고, 자발적 리뷰 축적 유인을 도서 내부에 설계.", "review"),
        ("5. 검색 노출 및 태깅 고도화", "무질서한 키워드 태깅 지양, 핵심 노출 키워드 3~4개를 선별해 집중 매핑함으로써 유입 전환율 향상.", "key")
    ]
    for idx, (title, desc, icon) in enumerate(strategies):
        row = idx // 3
        col = idx % 3
        w = Inches(3.9)
        h = Inches(2.2)
        x = Inches(0.6 + col * 4.1)
        y = Inches(1.8 + row * 2.4)
        if idx >= 3:
            x = Inches(2.65 + (idx - 3) * 4.1)
            
        draw_neobrutalist_card(slide, x, y, w, h, fill_color=COLOR_CARD_BG)
        # 아이콘 카드 배치
        draw_neobrutalist_icon(slide, icon, x + Inches(0.3), y + Inches(0.3), Inches(0.45))
        
        tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.85), w - Inches(0.6), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_TITLE
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_PRIMARY_TEXT
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_SECONDARY_TEXT
        p2.line_spacing = 1.2
    add_footer(slide, 29)
    set_speaker_notes(slide,
        "스물아홉 번째 장표는 이번 EDA의 액션 아이템인 5대 핵심 성공 전략입니다. "
        "가격을 문법서와 프리미엄 실무서로 이원화하고, 두꺼운 기술서에는 분철을 필수 연계하며, "
        "지식의 감가상각을 고려한 린 출판 모델 도입, 초기 골든타임 30일의 서평 유치, "
        "온라인 채널 검색 노출 3~4개 태그 최적화의 제언을 5개 직각 카드로 요약하여 설명합니다."
    )

    # --- Slide 30: 결론 및 Q&A ---
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_header(slide, "감사합니다 (Thank You)", "CONCLUSION & Q&A")
    draw_neobrutalist_card(slide, Inches(0.6), Inches(1.8), Inches(12.133), Inches(4.5))
    
    # 텍스트 박스의 크기와 정렬 개선 (수직 정렬 여백 보완)
    mid_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.1), Inches(10.333), Inches(4.0))
    tf = mid_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    
    # 메인 타이틀
    p = tf.paragraphs[0]
    p.text = "Q&A 및 미래 출판 패러다임 예측"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_BRICK
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(24)
    
    # 첫 번째 정보 (불릿 문자 제거, 중앙 정렬의 불릿 불일치 해결)
    p = tf.add_paragraph()
    p.text = "데이터 분석 결과 요약본 배포: PDF 보고서 및 데이터셋 시각화 템플릿 별도 제공"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_TEXT
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)
    
    # 두 번째 정보
    p = tf.add_paragraph()
    p.text = "질의응답: 분석 설계 모형 및 비즈니스 의사결정 액션 플랜 세부 질문"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_TEXT
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(28)
    
    # 맺음말 타이틀
    p = tf.add_paragraph()
    p.text = "[맺음말]"
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY_TEXT
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)
    
    # 맺음말 상세 (행간 및 문단 여백 보완)
    p = tf.add_paragraph()
    p.text = "이번 EDA 분석을 통해 검증되었듯, IT 도서 시장은 독자의 편의성(분철)과 기술 트렌드의 신속성(린 출판), 그리고 소셜 평판(리뷰 건수)에 의해 지배를 받는 특수한 비즈니스 필드입니다. 본 리포트의 전략을 바탕으로 더욱 민첩하고 유연한 비즈니스 로드맵이 수립되기를 기대합니다. 질문이 있으시다면 기꺼이 답변 드리겠습니다."
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_PRIMARY_TEXT
    p.line_spacing = 1.4
    p.alignment = PP_ALIGN.CENTER
    
    add_footer(slide, 30)
    set_speaker_notes(slide,
        "이상으로 발표를 모두 마치겠습니다. "
        "네오브루탈리즘의 거침없고 명확한 선들처럼, 우리 데이터 역시 오차 없이 명백한 비즈니스 유기성을 나타내고 있습니다. "
        "앞선 5대 전략을 통해 실무적으로 개발 서적의 가치를 전개할 예정이며, "
        "추가적으로 기술 모형 설계나 비즈니스 실행안에 대해 질의해 주시면 감사하겠습니다."
    )

    # 저장 및 마무리
    output_path = "yes24/docs/EDA_Slide_30Pages_NeoBrutalism.pptx"
    prs.save(output_path)
    print(f"네오브루탈리즘 PPTX 생성 완료: {output_path}")

if __name__ == "__main__":
    create_neobrutalist_presentation()
