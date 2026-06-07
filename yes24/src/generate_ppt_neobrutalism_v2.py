"""
Yes24 IT/컴퓨터 베스트셀러 EDA PPTX 보고서 - 네오브루탈리즘 v2 스타일 자동 생성 스크립트

이 스크립트는 yes24/docs/eda_report.md의 분석 결과를 바탕으로,
진정한 네오브루탈리즘 디자인(고채도 솔리드 배경, 3.5pt 두꺼운 검정 테두리,
0.15인치 하드 오프셋 블랙 섀도, 56~80pt 오버사이즈 숫자, 의도적 미스얼라인먼트)을
적용한 30페이지 PPTX 프레젠테이션을 생성합니다.

주요 개선 사항 (v1 대비):
- 슬라이드마다 다양한 고채도 솔리드 배경색 사용 (#F5F500, #CCFF00, #000000 등)
- 테두리 두께 2.5pt → 3.5pt로 강화
- 하드 섀도 오프셋 0.08인치 → 0.15인치로 대담하게 확대
- KPI 숫자 28pt → 56~72pt 오버사이즈
- 섹션 인트로 파트 번호 80pt+ 거대 표시
- 데코 블록 의도적 회전(rotation)으로 날것의 에너지 부여
- 카드 배경색 다양화 (흰색, 옐로, 벽돌색, 검정 믹스)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


def create_neobrutalist_v2():
    """네오브루탈리즘 v2 스타일의 30페이지 PPTX 프레젠테이션을 생성합니다."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # 네오브루탈리즘 v2 디자인 시스템
    # =========================================================================
    # 고채도 솔리드 컬러 팔레트
    C_YELLOW = RGBColor(245, 245, 0)       # 네오브루탈 시그니처 옐로우
    C_LIME = RGBColor(204, 255, 0)         # 라임 그린 (전략 슬라이드)
    C_SOFT_LIME = RGBColor(240, 255, 111)  # 연한 라임 (일부 배경)
    C_WHITE = RGBColor(255, 255, 255)      # 순백
    C_BLACK = RGBColor(0, 0, 0)            # 순검정
    C_BRICK = RGBColor(184, 80, 66)        # 포인트 벽돌색
    C_BLUE = RGBColor(0, 0, 255)           # 보조 액센트 파랑
    C_PINK = RGBColor(255, 45, 85)         # 핫핑크 보조
    C_DARK_GRAY = RGBColor(80, 80, 80)     # 보조 텍스트

    FONT_TITLE = "Gmarket Sans Bold"
    FONT_BODY = "NanumGothic"

    IMAGE_DIR = "yes24/images"
    ICON_DIR = "yes24/images/icons"

    # 하드 섀도 오프셋 (네오브루탈리즘 핵심: 대담한 크기)
    SHADOW_OFFSET = Inches(0.15)
    BORDER_WIDTH = Pt(3.5)
    ICON_BORDER = Pt(3)

    # =========================================================================
    # 헬퍼 함수
    # =========================================================================
    def fill_bg(slide, color):
        """슬라이드 전체를 솔리드 컬러로 채웁니다."""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()

    def add_deco_block(slide, x, y, w, h, color, rotation=0):
        """회전 가능한 데코 블록을 추가합니다. 테두리 포함."""
        deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        deco.fill.solid()
        deco.fill.fore_color.rgb = color
        deco.line.color.rgb = C_BLACK
        deco.line.width = Pt(2.5)
        if rotation:
            deco.rotation = rotation
        return deco

    def neo_card(slide, x, y, w, h, fill=None):
        """하드 오프셋 섀도 + 두꺼운 검정 테두리 카드를 생성합니다."""
        if fill is None:
            fill = C_WHITE
        # 1단계: 순검정 섀도 사각형 (뒤에 깔림)
        shadow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + SHADOW_OFFSET, y + SHADOW_OFFSET, w, h
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = C_BLACK
        shadow.line.fill.background()
        # 2단계: 전면 카드
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = fill
        card.line.color.rgb = C_BLACK
        card.line.width = BORDER_WIDTH
        return card

    def neo_icon(slide, icon_name, x, y, size=Inches(0.55)):
        """옐로우 배경 블록 + 아이콘 이미지를 배치합니다."""
        off = Inches(0.05)
        # 섀도
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + off, y + off, size, size)
        sh.fill.solid()
        sh.fill.fore_color.rgb = C_BLACK
        sh.line.fill.background()
        # 블록
        blk = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
        blk.fill.solid()
        blk.fill.fore_color.rgb = C_YELLOW
        blk.line.color.rgb = C_BLACK
        blk.line.width = ICON_BORDER
        # 아이콘
        icon_path = os.path.join(ICON_DIR, f"{icon_name}.png")
        if os.path.exists(icon_path):
            margin = Inches(0.08)
            slide.shapes.add_picture(
                icon_path, x + margin, y + margin,
                width=size - margin * 2, height=size - margin * 2
            )

    def add_header(slide, title, cat="YES24 IT/COMPUTER BESTSELLER EDA", dark=False):
        """상단 카테고리 라벨 + 메인 타이틀을 추가합니다."""
        txt_color = C_WHITE if dark else C_BLACK
        sub_color = RGBColor(180, 180, 180) if dark else C_DARK_GRAY
        # 카테고리
        cb = slide.shapes.add_textbox(Inches(1.0), Inches(0.35), Inches(11), Inches(0.3))
        tf = cb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = cat.upper()
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = sub_color
        # 타이틀
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12), Inches(0.9))
        tf2 = tb.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_bottom = tf2.margin_right = 0
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.name = FONT_TITLE
        p2.font.size = Pt(30)
        p2.font.bold = True
        p2.font.color.rgb = txt_color

    def add_footer(slide, num, dark=False):
        """하단 구분선 + 페이지 번호를 추가합니다."""
        line_color = C_WHITE if dark else C_BLACK
        # 굵은 구분선
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.85),
            Inches(12.133), Inches(0.04)
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = line_color
        ln.line.fill.background()
        # 푸터 텍스트
        fb = slide.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(10), Inches(0.3))
        tf = fb.text_frame
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = "Yes24 IT/컴퓨터 베스트셀러 EDA 리포트  |  NEO-BRUTALISM v2"
        p.font.name = FONT_BODY
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = C_WHITE if dark else C_BLACK
        # 페이지 번호
        pb = slide.shapes.add_textbox(Inches(11.733), Inches(6.95), Inches(1), Inches(0.3))
        tf2 = pb.text_frame
        tf2.margin_left = tf2.margin_top = tf2.margin_bottom = tf2.margin_right = 0
        p2 = tf2.paragraphs[0]
        p2.text = f"{num} / 30"
        p2.alignment = PP_ALIGN.RIGHT
        p2.font.name = FONT_TITLE
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = C_BRICK

    def notes(slide, text):
        """발표자 노트를 추가합니다."""
        ns = slide.notes_slide
        ns.notes_text_frame.text = text

    def chart_slide(pg, title, chart_file, stat_title, stats, insight, icon="chart"):
        """표준 차트+통계 슬라이드를 생성합니다."""
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        fill_bg(sl, C_WHITE)
        # 상단 옐로우 데코 스트라이프
        add_deco_block(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.18), C_YELLOW)
        add_header(sl, title)
        # 좌측: 차트 이미지 카드
        neo_card(sl, Inches(0.6), Inches(1.8), Inches(6.2), Inches(4.6))
        cp = os.path.join(IMAGE_DIR, chart_file)
        if os.path.exists(cp):
            sl.shapes.add_picture(
                cp, Inches(0.7), Inches(1.9),
                Inches(6.0), Inches(4.4)
            )
        # 우측: 통계/해석 카드
        neo_card(sl, Inches(7.1), Inches(1.8), Inches(5.6), Inches(4.6))
        # 아이콘
        neo_icon(sl, icon, Inches(7.5), Inches(2.15))
        # 통계 타이틀
        stb = sl.shapes.add_textbox(
            Inches(8.2), Inches(2.1), Inches(4.1), Inches(0.5)
        )
        tf = stb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = stat_title
        p.font.name = FONT_TITLE
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = C_BLACK
        # 통계 목록
        dtb = sl.shapes.add_textbox(
            Inches(7.5), Inches(2.85), Inches(4.8), Inches(3.4)
        )
        tf2 = dtb.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_bottom = tf2.margin_right = 0
        first = True
        for label, val in stats:
            p = tf2.paragraphs[0] if first else tf2.add_paragraph()
            first = False
            run1 = p.add_run()
            run1.text = f"▌ {label}: "
            run1.font.name = FONT_BODY
            run1.font.size = Pt(12)
            run1.font.bold = True
            run1.font.color.rgb = C_BLACK
            run2 = p.add_run()
            run2.text = val
            run2.font.name = FONT_BODY
            run2.font.size = Pt(12)
            run2.font.color.rgb = C_DARK_GRAY
            p.space_after = Pt(5)
        # 비즈니스 해석
        pi = tf2.add_paragraph()
        pi.text = "[ 비즈니스 해석 ]"
        pi.font.name = FONT_BODY
        pi.font.size = Pt(13)
        pi.font.bold = True
        pi.font.color.rgb = C_BRICK
        pi.space_before = Pt(10)
        pi.space_after = Pt(4)
        pb = tf2.add_paragraph()
        pb.text = insight
        pb.font.name = FONT_BODY
        pb.font.size = Pt(11)
        pb.font.color.rgb = C_BLACK
        pb.line_spacing = 1.3
        add_footer(sl, pg)
        return sl

    def dual_card_slide(pg, title, l_title, l_pts, r_title, r_pts):
        """2열 카드 비교 슬라이드를 생성합니다."""
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        fill_bg(sl, C_WHITE)
        add_deco_block(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.18), C_YELLOW)
        add_header(sl, title)
        for i, (ct, cpts) in enumerate([(l_title, l_pts), (r_title, r_pts)]):
            x = Inches(0.6 + i * 6.3)
            fill_c = C_WHITE if i == 0 else C_SOFT_LIME
            neo_card(sl, x, Inches(1.8), Inches(5.9), Inches(4.6), fill=fill_c)
            tb = sl.shapes.add_textbox(
                x + Inches(0.4), Inches(2.1),
                Inches(5.1), Inches(4.2)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
            pt = tf.paragraphs[0]
            pt.text = ct
            pt.font.name = FONT_TITLE
            pt.font.size = Pt(18)
            pt.font.bold = True
            pt.font.color.rgb = C_BRICK if i == 0 else C_BLUE
            pt.space_after = Pt(12)
            for bt, bb in cpts:
                p1 = tf.add_paragraph()
                p1.text = f"▶ {bt}"
                p1.font.name = FONT_BODY
                p1.font.size = Pt(13)
                p1.font.bold = True
                p1.font.color.rgb = C_BLACK
                p1.space_after = Pt(3)
                p2 = tf.add_paragraph()
                p2.text = bb
                p2.font.name = FONT_BODY
                p2.font.size = Pt(11)
                p2.font.color.rgb = C_DARK_GRAY
                p2.line_spacing = 1.35
                p2.space_after = Pt(14)
        add_footer(sl, pg)
        return sl

    def section_intro(pg, part_num, title, desc):
        """네오브루탈리즘 섹션 인트로 (검정 배경, 거대 파트 번호)를 생성합니다."""
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        fill_bg(sl, C_BLACK)
        # 파트 번호 거대 표시 (오버사이즈 시그니처)
        pn_box = sl.shapes.add_textbox(
            Inches(1.2), Inches(0.8), Inches(5), Inches(2.0)
        )
        tf = pn_box.text_frame
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = f"PART {part_num:02d}"
        p.font.name = FONT_TITLE
        p.font.size = Pt(80)
        p.font.bold = True
        p.font.color.rgb = C_BRICK
        # 벽돌색+옐로 데코 블록 (의도적 미스얼라인먼트)
        add_deco_block(sl, Inches(9.5), Inches(0.5), Inches(3.0), Inches(1.5), C_BRICK, rotation=-5)
        add_deco_block(sl, Inches(10.2), Inches(1.8), Inches(2.2), Inches(1.0), C_YELLOW, rotation=3)
        # 메인 카드
        neo_card(sl, Inches(1.2), Inches(3.0), Inches(10.933), Inches(3.5), fill=C_YELLOW)
        tb2 = sl.shapes.add_textbox(
            Inches(1.8), Inches(3.3), Inches(9.7), Inches(2.8)
        )
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_bottom = tf2.margin_right = 0
        pt = tf2.paragraphs[0]
        pt.text = title
        pt.font.name = FONT_TITLE
        pt.font.size = Pt(36)
        pt.font.bold = True
        pt.font.color.rgb = C_BLACK
        pt.space_after = Pt(12)
        pd = tf2.add_paragraph()
        pd.text = desc
        pd.font.name = FONT_BODY
        pd.font.size = Pt(14)
        pd.font.bold = True
        pd.font.color.rgb = C_BLACK
        pd.line_spacing = 1.4
        add_footer(sl, pg, dark=True)
        return sl

    # =========================================================================
    # 30개 슬라이드 생성 시작
    # =========================================================================
    layout = prs.slide_layouts[6]

    # --- Slide 1: 표지 (옐로우 배경) ---
    s1 = prs.slides.add_slide(layout)
    fill_bg(s1, C_YELLOW)
    # 좌측 하단 벽돌색 데코 (의도적 회전)
    add_deco_block(s1, Inches(0.3), Inches(5.5), Inches(3.0), Inches(1.8), C_BRICK, rotation=-4)
    # 우측 상단 파랑 데코 (의도적 회전)
    add_deco_block(s1, Inches(10.5), Inches(0.3), Inches(2.5), Inches(1.2), C_BLUE, rotation=3)
    # 메인 카드
    neo_card(s1, Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
    # 제목
    ttl = s1.shapes.add_textbox(Inches(1.6), Inches(1.8), Inches(10.1), Inches(2.2))
    tf = ttl.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "Yes24 IT/컴퓨터 베스트셀러\n데이터 심층 EDA 보고서"
    p.font.name = FONT_TITLE
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = C_BLACK
    # 부제
    sub = s1.shapes.add_textbox(Inches(1.6), Inches(4.0), Inches(10.1), Inches(0.6))
    tf2 = sub.text_frame
    tf2.margin_left = tf2.margin_top = tf2.margin_bottom = tf2.margin_right = 0
    p2 = tf2.paragraphs[0]
    p2.text = "데이터 기반 기술 실용서 시장 분석 및 비즈니스 마케팅 전략 제언"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = C_BRICK
    # 정보
    info = s1.shapes.add_textbox(Inches(1.6), Inches(4.8), Inches(10.1), Inches(0.5))
    tf3 = info.text_frame
    tf3.margin_left = tf3.margin_top = tf3.margin_bottom = tf3.margin_right = 0
    p3 = tf3.paragraphs[0]
    p3.text = "작성일: 2026년 06월 06일  |  작성자: 20년차 수석 데이터 분석가  |  NEO-BRUTALISM v2"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = C_BLACK
    add_footer(s1, 1)
    notes(s1,
        "안녕하세요. 오늘 발표할 자료는 Yes24 IT/컴퓨터 베스트셀러 데이터 심층 EDA 보고서입니다. "
        "이번 프레젠테이션은 스타트업의 도발적이고 거침없는 에너지를 투영하는 '네오브루탈리즘' 디자인 스타일로 전면 리디자인했습니다. "
        "두꺼운 3.5포인트 검정 테두리와 대담한 하드 오프셋 섀도를 모든 카드에 적용하여 압도적인 시각 임팩트를 확보했고, "
        "고채도 옐로우와 벽돌색 포인트가 교차하며 강렬한 대비를 만들어냅니다. "
        "G마켓 산스 Bold로 타이틀의 무게감을 더하고, 나눔고딕으로 본문의 가독성을 유지했습니다. "
        "베스트셀러 1,000건의 통계 마이닝을 바탕으로 실질적인 출판 마케팅 로드맵을 제안해 드리겠습니다."
    )

    # --- Slide 2: 목차 ---
    s2 = prs.slides.add_slide(layout)
    fill_bg(s2, C_WHITE)
    # 상단 옐로우 스트라이프 2개 (비스듬한 데코)
    add_deco_block(s2, Inches(0), Inches(0), Inches(13.333), Inches(0.25), C_YELLOW)
    add_deco_block(s2, Inches(8), Inches(0.35), Inches(5.333), Inches(0.12), C_BRICK, rotation=0)
    add_header(s2, "목차 및 분석 로드맵")
    toc_data = [
        ("1부. 단변량 분석", "판매가 분포, 평점 만족도, 판매지수, 분철 서비스, 출판사 점유율 등 개별 변수의 기초 통계와 시각화 데이터를 분석합니다.", "book", C_WHITE),
        ("2부. 이변량·다변량 분석", "할인율·분철 교차 분석, 출판 트렌드 시계열, 상관관계, TF-IDF 텍스트 마이닝으로 시장 수요 예측을 다룹니다.", "chart", C_SOFT_LIME),
        ("3부. 비즈니스 전략", "가격 이원화, 린 퍼블리싱, 소셜 증명 관리, 태깅 최적화 등 데이터 기반 액션 플랜을 제시합니다.", "lightbulb", C_YELLOW),
    ]
    for i, (t, d, ic, bg_c) in enumerate(toc_data):
        x = Inches(0.6 + i * 4.1)
        neo_card(s2, x, Inches(1.8), Inches(3.8), Inches(4.5), fill=bg_c)
        neo_icon(s2, ic, x + Inches(0.35), Inches(2.15))
        tb = s2.shapes.add_textbox(x + Inches(0.35), Inches(2.9), Inches(3.1), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = t
        p.font.name = FONT_TITLE
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = C_BLACK
        db = s2.shapes.add_textbox(x + Inches(0.35), Inches(3.6), Inches(3.1), Inches(2.5))
        tf2 = db.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_bottom = tf2.margin_right = 0
        p2 = tf2.paragraphs[0]
        p2.text = d
        p2.font.name = FONT_BODY
        p2.font.size = Pt(12)
        p2.font.color.rgb = C_BLACK
        p2.line_spacing = 1.35
    add_footer(s2, 2)
    notes(s2,
        "목차 슬라이드입니다. 네오브루탈리즘의 두꺼운 테두리와 하드 섀도가 적용된 3개 카드가 나란히 배치되어 있습니다. "
        "각 카드의 배경색이 흰색, 연한 라임, 옐로우로 다르게 설정되어 시각적 리듬감을 부여했습니다. "
        "1부는 단변량 기초 분석, 2부는 이변량 및 다변량 교차 연계 분석, 3부는 종합 비즈니스 마케팅 전략 제언으로 구성됩니다."
    )

    # --- Slide 3: 데이터 세트 개요 ---
    s3 = prs.slides.add_slide(layout)
    fill_bg(s3, C_WHITE)
    add_deco_block(s3, Inches(0), Inches(0), Inches(13.333), Inches(0.18), C_YELLOW)
    add_header(s3, "데이터 세트 수집 개요 및 스키마 구조")
    # 좌측 카드 (오버사이즈 숫자 포함)
    neo_card(s3, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.6))
    # 오버사이즈 숫자
    big = s3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.0), Inches(1.5))
    tf = big.text_frame
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "1,000"
    p.font.name = FONT_TITLE
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = C_BRICK
    p2 = tf.add_paragraph()
    p2.text = "개 행(Rows) × 30개 열(Columns)"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = C_BLACK
    p2.space_after = Pt(8)
    specs = [
        ("수집 대상", "Yes24 IT/컴퓨터 카테고리 베스트셀러"),
        ("중복 데이터", "0건 (완전 제거 완료)"),
        ("분석 일자", "2026년 06월 05일 기준"),
    ]
    for label, val in specs:
        p = tf.add_paragraph()
        run1 = p.add_run()
        run1.text = f"▌ {label}: "
        run1.font.name = FONT_BODY
        run1.font.size = Pt(13)
        run1.font.bold = True
        run1.font.color.rgb = C_BLACK
        run2 = p.add_run()
        run2.text = val
        run2.font.name = FONT_BODY
        run2.font.size = Pt(13)
        run2.font.color.rgb = C_DARK_GRAY
        p.space_after = Pt(4)
    # 우측 카드
    neo_card(s3, Inches(6.7), Inches(1.8), Inches(5.8), Inches(4.6), fill=C_SOFT_LIME)
    rb = s3.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf_r = rb.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_bottom = tf_r.margin_right = 0
    pr = tf_r.paragraphs[0]
    pr.text = "핵심 변수 구분"
    pr.font.name = FONT_TITLE
    pr.font.size = Pt(20)
    pr.font.bold = True
    pr.font.color.rgb = C_BLACK
    pr.space_after = Pt(8)
    vars_data = [
        ("대상 변수", "도서명, 저자, 출판사, 발행 연/월, 태그"),
        ("성과 지표", "판매지수(Sale Index), 평점(Rating), 리뷰 수"),
        ("기능/가격", "판매가, 정가, 할인율, 분철 여부"),
        ("파생 변수", "정제된 가격(clean), 태그 개수, 할인율 범주"),
    ]
    for label, val in vars_data:
        p = tf_r.add_paragraph()
        run1 = p.add_run()
        run1.text = f"▌ {label}: "
        run1.font.name = FONT_BODY
        run1.font.size = Pt(13)
        run1.font.bold = True
        run1.font.color.rgb = C_BLACK
        run2 = p.add_run()
        run2.text = val
        run2.font.name = FONT_BODY
        run2.font.size = Pt(13)
        run2.font.color.rgb = C_DARK_GRAY
        p.space_after = Pt(8)
    add_footer(s3, 3)
    notes(s3,
        "세 번째 장표는 수집된 원본 데이터의 구조와 변수 목록을 설명합니다. "
        "좌측 카드에는 '1,000'이라는 데이터 건수를 64포인트 오버사이즈 벽돌색 숫자로 파격 배치하여 "
        "네오브루탈리즘 특유의 '하나의 오버사이즈 숫자가 레이아웃을 깨뜨리는' 시그니처를 구현했습니다. "
        "우측 카드는 연한 라임 배경으로 변수 구분을 시각적으로 분리했습니다."
    )

    # --- Slide 4: 수치형 KPI 대시보드 (오버사이즈 숫자) ---
    s4 = prs.slides.add_slide(layout)
    fill_bg(s4, C_WHITE)
    add_deco_block(s4, Inches(0), Inches(0), Inches(13.333), Inches(0.18), C_YELLOW)
    add_header(s4, "전체 수치형 변수 통합 기술통계 분석")
    kpis = [
        ("평균 판매가", "23,480원", "중위값 22,500원 | 최소 4,500원", "price", C_WHITE, C_BRICK),
        ("평균 판매지수", "3,026점", "중위값 1,236점 | 최대 87,480점", "trend", C_YELLOW, C_BLACK),
        ("평균 도서 평점", "7.50점", "중위값 9.70점 (0점 제외 시)", "star", C_BRICK, C_WHITE),
        ("평균 할인율", "8.64%", "최빈값 10% (정가제 상한)", "percent", C_BLACK, C_YELLOW),
    ]
    for idx, (label, val, sub, icon, bg, txt_c) in enumerate(kpis):
        col = idx % 2
        row = idx // 2
        x = Inches(0.6 + col * 6.3)
        y = Inches(1.8 + row * 2.55)
        neo_card(s4, x, y, Inches(5.9), Inches(2.3), fill=bg)
        neo_icon(s4, icon, x + Inches(0.35), y + Inches(0.3))
        tb = s4.shapes.add_textbox(x + Inches(1.1), y + Inches(0.2), Inches(4.4), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        # 라벨
        pl = tf.paragraphs[0]
        pl.text = label.upper()
        pl.font.name = FONT_BODY
        pl.font.size = Pt(11)
        pl.font.bold = True
        pl.font.color.rgb = C_DARK_GRAY if bg == C_WHITE or bg == C_YELLOW else RGBColor(200, 200, 200)
        # 오버사이즈 숫자 (56pt)
        pv = tf.add_paragraph()
        pv.text = val
        pv.font.name = FONT_TITLE
        pv.font.size = Pt(56)
        pv.font.bold = True
        pv.font.color.rgb = txt_c
        # 서브텍스트
        ps = tf.add_paragraph()
        ps.text = sub
        ps.font.name = FONT_BODY
        ps.font.size = Pt(11)
        ps.font.bold = True
        ps.font.color.rgb = C_DARK_GRAY if bg == C_WHITE or bg == C_YELLOW else RGBColor(200, 200, 200)
    add_footer(s4, 4)
    notes(s4,
        "네 번째 장표는 수치형 데이터의 핵심 KPI를 2×2 카드 그리드로 표현한 대시보드입니다. "
        "각 카드의 배경색이 흰색, 옐로우, 벽돌색, 검정으로 모두 다르게 설정되어 네오브루탈리즘의 강렬한 컬러 대비를 극대화했습니다. "
        "KPI 숫자를 56포인트로 파격적으로 키워 시선을 즉각 사로잡습니다. "
        "평균 판매가 23,480원, 평균 판매지수 3,026점, 평균 평점 7.5점, 할인율 8.64%의 핵심 지표가 한눈에 들어옵니다."
    )

    # --- Slide 5: 범주형 통계 요약 ---
    s5 = prs.slides.add_slide(layout)
    fill_bg(s5, C_WHITE)
    add_deco_block(s5, Inches(0), Inches(0), Inches(13.333), Inches(0.18), C_YELLOW)
    add_header(s5, "전체 범주형 변수 요약 및 시장 쏠림 현상")
    # 좌측 카드
    neo_card(s5, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.6))
    lb = s5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.0), Inches(4.3))
    tf = lb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    # 오버사이즈
    p = tf.paragraphs[0]
    p.text = "869명"
    p.font.name = FONT_TITLE
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = C_BLUE
    ps = tf.add_paragraph()
    ps.text = "고유 저자 수 (1,000 베스트셀러 중)"
    ps.font.name = FONT_BODY
    ps.font.size = Pt(14)
    ps.font.bold = True
    ps.font.color.rgb = C_BLACK
    ps.space_after = Pt(10)
    cat_items = [
        ("고유 도서 수", "1,000권 (중복 없는 유일 도서명)"),
        ("등록 출판사 수", "187개사 (상위 과점 현상)"),
        ("분철 미지원", "853권 (85.3%)"),
        ("분철 지원", "147권 (14.7%)"),
    ]
    for label, val in cat_items:
        p = tf.add_paragraph()
        run1 = p.add_run()
        run1.text = f"▌ {label}: "
        run1.font.name = FONT_BODY
        run1.font.size = Pt(13)
        run1.font.bold = True
        run1.font.color.rgb = C_BLACK
        run2 = p.add_run()
        run2.text = val
        run2.font.name = FONT_BODY
        run2.font.size = Pt(13)
        run2.font.color.rgb = C_DARK_GRAY
        p.space_after = Pt(4)
    # 우측 카드
    neo_card(s5, Inches(6.7), Inches(1.8), Inches(5.8), Inches(4.6), fill=C_YELLOW)
    rb = s5.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf_r = rb.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_bottom = tf_r.margin_right = 0
    pr = tf_r.paragraphs[0]
    pr.text = "범주형 데이터 시사점"
    pr.font.name = FONT_TITLE
    pr.font.size = Pt(20)
    pr.font.bold = True
    pr.font.color.rgb = C_BLACK
    pr.space_after = Pt(10)
    insights = [
        "1. 저자 집중도 완화: 고유 저자 869명으로 다양한 도서가 순위에 진입하는 분산 구조입니다.",
        "2. 출판사 양극화: 한빛미디어 등 소수 대형 출판사가 상위권을 지배하는 과점 성향입니다.",
        "3. 기능 편의 격차: 분철 서비스 제공 도서는 14.7%에 불과해 잠재 수요 대비 공급이 희소합니다.",
    ]
    for ins in insights:
        p = tf_r.add_paragraph()
        p.text = ins
        p.font.name = FONT_BODY
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_BLACK
        p.line_spacing = 1.35
        p.space_after = Pt(10)
    add_footer(s5, 5)
    notes(s5,
        "다섯 번째 슬라이드는 범주형 변수 분석입니다. "
        "좌측 카드에 '869명'이라는 고유 저자 수를 56포인트 파랑 컬러로 오버사이즈 배치했습니다. "
        "우측 카드는 옐로우 배경으로 시사점 3가지를 명시합니다. "
        "출판사는 187개사이지만 한빛미디어 등 소수 대형사가 과반을 차지하며, 분철 도서는 14.7%에 불과합니다."
    )

    # --- Slide 6: PART 1 인트로 ---
    s6 = section_intro(6, 1, "단변량 분석 및 핵심 분포 탐색",
        "판매가, 도서 평점, 판매지수, 분철 서비스, 출판사 분포 등\n주요 개별 지표들의 통계적 분포와 비즈니스 인사이트를 도출합니다.")
    notes(s6,
        "PART 1 인트로 슬라이드입니다. 검정 배경 위에 'PART 01'이 80포인트 벽돌색으로 거대하게 표시되어 "
        "네오브루탈리즘의 '오버사이즈 요소가 레이아웃을 깨뜨리는' 시그니처를 구현했습니다. "
        "우측 상단의 벽돌색과 옐로우 데코 블록은 의도적으로 회전시켜 날것의 에너지를 표현합니다. "
        "하단 옐로우 카드에 이번 파트의 분석 범위를 명시했습니다."
    )

    # --- Slide 7: 판매가 분포 ---
    s7 = chart_slide(7, "1.1 판매가(Sale Price) 빈도 분포 분석",
        "01_sale_price_distribution.png", "판매가 구간별 기초 통계",
        [("집중 구간", "20,000원 ~ 30,000원 대역 과밀 분포"),
         ("평균 판매가", "23,480원"),
         ("중위 판매가", "22,500원"),
         ("최소/최대가", "4,500원 / 67,000원")],
        "Yes24 IT 베스트셀러 도서들은 약 2만~3만원 사이 대역에 70% 가까이 밀집되어 있습니다. 독자들이 기술 서적 구매 시 심리적으로 수용하는 표준 가격대가 2만원대 초중반에 견고하게 형성되어 있음을 방증합니다.",
        "price")
    notes(s7,
        "일곱 번째 슬라이드는 도서 판매가의 빈도 분포 차트 분석입니다. "
        "좌측의 차트 이미지 카드와 우측의 통계 카드 모두 3.5포인트 두꺼운 검정 테두리와 "
        "0.15인치 하드 섀도가 적용되어 네오브루탈리즘 특유의 강렬한 인상을 줍니다. "
        "대다수 IT 도서 가격이 2만원대 중후반에 위치하고 있으며, "
        "이는 독자의 심리적 가격 수용 범위와 정확히 일치합니다."
    )

    # --- Slide 8: 판매가 심층분석 ---
    s8 = dual_card_slide(8, "1.1 [심층분석] 가격 심리적 저항선과 세그먼트 전략",
        "심리적 마지노선 분석",
        [("3만원 가격 장벽", "판매가가 30,000원 선을 돌파할 때, 독자가 즉각 결제를 주저하는 심리적 이탈 저항선이 뚜렷이 나타납니다."),
         ("보급형 입문서 설계", "단기적 베스트셀러 성장을 도모하는 입문서/문법서는 최종 할인가 22,000~26,000원 구간에 걸치도록 타겟 정가를 기획하는 것이 최선입니다.")],
        "고가 프리미엄 세그먼트 전략",
        [("전문성 희소가치 수용", "RAG 아키텍처 구축이나 LLM 미세 조정 등 난이도 높은 독점적 전문서라면, 4만원 이상으로 가격을 높여도 구매로 이어집니다."),
         ("가치 대비 가격 책정", "기술 대체 불가능성이 높다면 분량 보강과 사은 혜택을 얹어 고가 프리미엄 전략이 마진 방어에 유리합니다.")])
    notes(s8,
        "여덟 번째 슬라이드는 가격 분포의 심층 해석입니다. "
        "좌측 카드는 흰색, 우측 카드는 연한 라임 배경으로 시각적 차별화를 주었습니다. "
        "독자들은 3만원 저항선과 전문서 고가 결제 용의라는 양극의 가치를 갖고 있어 기획 가격 이원화가 필요합니다."
    )

    # --- Slide 9: 도서 평점 분포 ---
    s9 = chart_slide(9, "1.2 도서 평점(Rating) 빈도 분포 분석",
        "02_rating_distribution.png", "평점 분포 수치 요약",
        [("만점 밀집도", "9.5~10.0점 대역에 극단적 우측 쏠림"),
         ("전체 평균 평점", "7.50점 (0점 무평가 도서 혼재 효과)"),
         ("평가도서 중위수", "9.70점 (독자 만족도 우수)"),
         ("최소/최대 평점", "0.00점 / 10.00점")],
        "도서 평점의 히스토그램을 분석하면 대다수 평점이 만점에 밀집되어 일반 평점 점수 자체는 도서의 우열을 가리는 변별력을 완전히 소실했습니다. 평균 7.5점은 미평가 도서의 영점 효과로 발생한 굴절입니다.",
        "star")
    notes(s9,
        "아홉 번째 슬라이드는 도서 평점의 분포 차트입니다. "
        "대부분의 책들이 만점에 가까운 9.5~10점 사이에 모여 있고, 미평가 도서들의 0점 효과가 관측됩니다. "
        "이는 평점 자체보다 리뷰 건수가 실질적 변별력 지표임을 의미합니다."
    )

    # --- Slide 10: 평점 심층분석 ---
    s10 = dual_card_slide(10, "1.2 [심층분석] 평점 인플레이션 극복 및 0점 돌파 전략",
        "평점 인플레이션의 대용 지표",
        [("소셜 증거(Social Proof)", "독자들은 평점 10점보다, 상세 텍스트 리뷰가 100개 넘게 누적된 책을 강력히 의지합니다."),
         ("리뷰 수로 무게중심 이동", "평점 스펙보다 절대적인 리뷰 볼륨을 키우는 것이 잠재 독자 구매 전환의 핵심입니다.")],
        "신간 도서 0점 신속 탈출 지침",
        [("0점 도서의 치명적 노출 방해", "0점 도서는 미검증 낙인 효과를 주어 독자의 장바구니 결제를 심리적으로 차단합니다."),
         ("초기 1개월 마케팅", "출간 즉시 서평단 유치와 리뷰 마일리지 보상에 집중해 최소 10건 이상의 초기 서평을 확보해야 합니다.")])
    notes(s10,
        "열 번째 슬라이드는 평점 인플레이션 극복 마케팅 지침입니다. "
        "의사결정의 무게추는 절대 평점이 아니라 누적 리뷰 건수입니다. "
        "신간 출시 시 0점 노출에 갇히지 않도록 초기 1개월 골든타임 내에 리뷰 자산을 빌드업해야 합니다."
    )

    # --- Slide 11: 판매지수 분포 ---
    s11 = chart_slide(11, "1.3 판매지수(Sale Index) 상자 그림 분포 분석",
        "03_sale_index_boxplot.png", "판매지수 통계 및 사분위",
        [("중위 판매지수", "1,236점 (현실적 성과 척도)"),
         ("평균 판매지수", "3,026점 (아웃라이어 왜곡 효과)"),
         ("최대 판매지수", "87,480점 (독보적 메가 히트 타이틀)"),
         ("표준편차(std)", "7,077점 (매우 큰 분포 변동성)")],
        "판매지수 박스플롯을 보면 대다수 베스트셀러가 5,000점 이하에 밀집되어 있으나 상위 킬러 타이틀은 최고 8.7만점에 육박합니다. 베스트셀러 내에서도 판매 양극화가 일어나는 전형적인 롱테일 구조입니다.",
        "trend")
    notes(s11,
        "열한 번째 장표는 판매지수 상자 그림 분석입니다. "
        "거의 모든 도서가 5,000점 이하의 바닥권에 수렴해 있으나 상위 점들은 극단적으로 높습니다. "
        "소수의 메가 히트작이 시장 총매출을 리드하는 과점 생태계입니다."
    )

    # --- Slide 12: 판매지수 심층분석 ---
    s12 = dual_card_slide(12, "1.3 [심층분석] 롱테일 시장 구조와 현실적 KPI 수립",
        "도서 시장의 파레토 법칙",
        [("소수 킬러 타이틀 독점", "상위 20% 서적이 총 판매지수의 압도적 지분을 장악하며 유통 플랫폼의 헤드 역할을 담당합니다."),
         ("평균 왜곡의 함정", "단순 평균값(3,026)을 목표로 삼으면 과도한 마케팅 비용이나 예산 불균형을 야기합니다.")],
        "중위수(Median) 벤치마킹 액션",
        [("1차 마일스톤 설정", "현실적 신작 성공 기준은 중위선인 '판매지수 1,236점'을 달성하는 것으로 잡고, 2차 확산 전략을 모색하는 것이 리스크 관리의 지름길입니다."),
         ("업데이트 투자", "상위 아웃라이어 타이틀은 1회성 마케팅이 아닌, 깃허브 피드백 수용과 신속한 개정판 관리의 산물입니다.")])
    notes(s12,
        "열두 번째 장표는 판매지수 왜곡을 고려한 신작 경영 가이드라인입니다. "
        "신간의 현실적 목표를 중위수 1,236점으로 재설정하여 합리적 손익분기 관리를 유도합니다."
    )

    # --- Slide 13: 분철 서비스 ---
    s13 = chart_slide(13, "1.4 분철 서비스 제공 여부(Spring Service) 분석",
        "04_spring_service_count.png", "분철 제공 수량 및 점유비",
        [("분철 미제공 도서(N)", "853권 (85.3% 압도적 지분)"),
         ("분철 제공 도서(Y)", "147권 (14.7% 소수 공급)"),
         ("주요 적용 도서", "수험서, 개발 실무 매뉴얼, 대학 전공서")],
        "분철 지원 서적은 14.7%에 불과합니다. 제작 공정 단가와 물류 편의상 대다수 출판사가 분철을 지양하고 있으나, 실무 전공 서적 독자에게는 잠재 수요가 거대하게 깔려 있습니다.",
        "spring")
    notes(s13,
        "열세 번째 장표는 분철 서비스 제공 비율 차트 분석입니다. "
        "대다수 도서가 일반 제본이지만, 두꺼운 실용서 14.7%에 한해 분철이 제공되고 있습니다."
    )

    # --- Slide 14: 출판사 점유율 ---
    s14 = chart_slide(14, "1.5 상위 30개 출판사별 베스트셀러 점유율",
        "05_top_30_publishers.png", "상위 출판사 점유 통계",
        [("압도적 1위", "한빛미디어 (150권, 15.0% 독점)"),
         ("2~4위권", "길벗(7.4%), 제이펍(5.1%), 이지스퍼블리싱(5.0%)"),
         ("5위 신흥 주자", "골든래빗 (43권, 4.3% 상승)"),
         ("상위 10개사 점유율", "과반(약 50.2%) 돌파")],
        "국내 IT 베스트셀러 시장은 대형 메이저 브랜드의 점유 집중도가 극심합니다. 상위 10개 퍼블리셔가 전체의 과반을 점유해, 신진 브랜드 진입 장벽이 매우 가파릅니다.",
        "building")
    notes(s14,
        "열네 번째 슬라이드는 출판사 점유 현황입니다. "
        "한빛미디어가 15%의 독보적 1위이며, 상위 10개사가 시장 과반을 차지합니다."
    )

    # --- Slide 15: 출판사 심층분석 ---
    s15 = dual_card_slide(15, "1.5 [심층분석] 메이저 출판사 브랜드 락인과 상생 전략",
        "대형 퍼블리셔의 후광 효과",
        [("베타 독자단 인프라", "한빛이나 길벗 등은 수만 명의 IT 전문 베타 리더단과 충성도 높은 뉴스레터 구독자 체계를 보유합니다."),
         ("기획 투고 전략", "무명 저자는 마케팅 인프라가 전무하므로, 대형 출판사의 기획 투고를 통해 브랜드 후광을 타는 것이 성공 기회를 열어줍니다.")],
        "틈새 린(Lean) 론칭 전략",
        [("대형사의 느린 속도 틈새", "대형 퍼블리셔는 내부 결재와 편집 공정이 길어, 신생 프레임워크나 최신 AI API에 민첩하게 반응하기 힘듭니다."),
         ("마이크로 커뮤니티 타겟", "독립 출판이나 소형 기획사는 초단기 트렌드를 선점하고, 개발자 커뮤니티를 활용해 기동성 있게 시장을 우회 침투해야 합니다.")])
    notes(s15,
        "열다섯 번째 장표는 메이저 과점 구도에 대한 상생 및 우회 기획안입니다. "
        "대형사의 거대한 마케팅 리소스를 이기기 힘드므로 파트너십 투고 전략을 우선시하되, "
        "소형 출판사는 대형사보다 민첩하게 최신 AI 트렌드를 포착해 슬림북 형태로 신속 출간해야 합니다."
    )

    # --- Slide 16: PART 2 인트로 ---
    s16 = section_intro(16, 2, "이변량 및 다변량 연계 분석",
        "할인율·분철 서비스 교차 분석, 출판 트렌드 시계열,\n다변량 피어슨 상관관계, TF-IDF 텍스트 마이닝으로\nIT 도서 시장의 실질적인 수요 동력(Driver)을 도출합니다.")
    notes(s16,
        "PART 2 인트로입니다. 검정 배경에 'PART 02'가 80포인트 벽돌색으로 거대하게 배치되었습니다. "
        "이변량 및 다변량 분석을 통해 변수 간 교차 관계와 시장 동력을 검증합니다."
    )

    # --- Slide 17: 할인율 vs 판매지수 ---
    s17 = chart_slide(17, "2.1 할인율(Discount Rate) vs 판매지수 산점도",
        "06_discount_vs_sale_index.png", "할인 정책별 세부 지표",
        [("10% 최대 할인 적용", "851권  |  평균 판매지수: 3,394점"),
         ("0% 무할인(정가)", "123권  |  평균 판매지수: 826점"),
         ("5% 부분 할인", "25권  |  평균 판매지수: 1,276점"),
         ("정가제 상한 도달", "10% 할인이 판매 안착의 필수 기본 조건화")],
        "거의 모든 베스트셀러 도서가 최대 상한인 10% 할인으로 유통됩니다. 무할인 도서는 판매지수가 4분의 1 토막에 그쳐, 10% 할인은 선택이 아닌 필수입니다. 추가 경쟁력은 비금전적 가치에서 발굴해야 합니다.",
        "percent")
    notes(s17,
        "열일곱 번째 장표는 할인율과 판매지수 간의 산점도 분석입니다. "
        "정가 도서는 평균 지수가 현저히 처지며, 10% 할인 제공 도서가 압도적입니다. "
        "가격 깎기는 기본일 뿐, 추가 차별화는 기능 편의에서 모색해야 합니다."
    )

    # --- Slide 18: 분철 vs 판매지수 ---
    s18 = chart_slide(18, "2.2 분철 여부에 따른 판매 성과 교차 분석",
        "07_spring_service_vs_sale_index.png", "분철 여부별 판매지수 비교",
        [("분철 지원 그룹 (Y)", "147권  |  평균: 8,687점"),
         ("분철 미지원 그룹 (N)", "853권  |  평균: 2,050점"),
         ("평균 격차 배율", "약 4.2배 초과 달성"),
         ("중위 지수 비교", "Y: 4,062점  |  N: 1,047점")],
        "분철 서비스를 장착한 도서들의 평균 판매 성과가 미지원 도서 대비 무려 4.2배 압도적으로 높게 도출되었습니다. 단편적 제작비 증가분 이상의 판매 견인 효과를 통계적으로 규명합니다.",
        "spring")
    notes(s18,
        "열여덟 번째 장표는 분철 여부별 판매 성과 교차 분석입니다. "
        "평균치와 중위수 모두 분철 가능한 도서군이 4배 가까이 앞섭니다. "
        "링 제본이 실습 효율을 높여 독자의 장바구니 선택을 결정짓는 주 동인임을 증명합니다."
    )

    # --- Slide 19: 분철 심층분석 ---
    s19 = dual_card_slide(19, "2.2 [심층분석] 분철 가치와 출판 유통사 락인 설계",
        "독자 실용 소비주의 이해",
        [("핸즈프리(Hands-free) 학습 니즈", "IT 독자들은 책을 모니터 옆에 완전히 젖혀두고 키보드를 자유롭게 조작하고 싶어 합니다."),
         ("가치 체감 대비 저렴한 투자", "독자는 1,000~2,000원 소액 추가금을 주더라도 책의 활용성을 높일 수 있다면 흔쾌히 선택합니다.")],
        "출판-유통망 연계 마케팅",
        [("출판 유통의 기본 탑재", "500페이지 이상 기술 서적 기획 시 분철 라인업을 상설 채널로 확보하고 전면 홍보해야 합니다."),
         ("구매 상세페이지 노출 강화", "상세 페이지 최상단에 '스프링 분철 옵션 완비' 표시를 강조하여 타사 도서와의 구매 우위를 확립합니다.")])
    notes(s19,
        "열아홉 번째 장표는 분철 성과 기반 실무 마케팅 액션안입니다. "
        "독자들의 핸즈프리 학습 니즈와 소액 추가금 수용 심리를 활용하여 "
        "온라인 몰 상세 페이지에 '분철 완비'를 전면 기획 노출해야 합니다."
    )

    # --- Slide 20: 출판년도별 추이 ---
    s20 = chart_slide(20, "2.3 출판년도별 도서 수 및 평균 판매지수 추이",
        "08_publish_year_trends.png", "출판 연도별 성과 통계",
        [("2025년 출간 도서", "367권  |  평균: 3,902점"),
         ("2026년 출간 도서", "356권  |  평균: 2,404점"),
         ("2024년 출간 도서", "138권  |  평균: 1,747점"),
         ("2023년 이전 도서", "총 139권  |  급감 흐름")],
        "2025~2026년 신간들이 베스트셀러 진입 권수와 판매지수의 대부분을 장악합니다. 연도가 지날수록 도서 수가 급감하여 IT 분야 서적의 감가상각과 트렌드 이탈 속도가 초고속입니다.",
        "calendar")
    notes(s20,
        "스무 번째 슬라이드는 출판 연도별 성과 분석입니다. "
        "신간이 베스트셀러 진입 권수와 판매지수를 압도적으로 주도하는 지식 감가상각 현상입니다."
    )

    # --- Slide 21: 출판년도 심층분석 ---
    s21 = dual_card_slide(21, "2.3 [심층분석] 시계열 지식 감가상각과 린 출판 대응",
        "IT 도서 지식의 급격한 유통기한",
        [("오픈소스 버전업 리스크", "라이브러리 마이너 버전 패치 하나만으로도 실습 예제가 에러를 발생시키며 소비자 평판 추락으로 이어집니다."),
         ("기존 기획 관행 탈피", "과거의 1년짜리 기획-집필-출판 루프에서 벗어나지 않으면 출간 당일 이미 구닥다리 기술서가 될 위험이 있습니다.")],
        "린 퍼블리싱(Lean Publishing) 도입",
        [("실시간 독자 피드백 결합", "원고 기획서 단계에서부터 깃허브 저장소를 개설하고, 독자와 소통하며 완성 시점에 신속히 출판해야 합니다."),
         ("콘텐츠 슬림화 전략", "트렌드가 민감한 프론트엔드나 생성형 AI 서적은 분량을 200페이지 내외로 축소하여 런칭 주기를 민첩하게 회전해야 합니다.")])
    notes(s21,
        "스물한 번째 장표는 지식 감가상각을 우회하는 린 출판 프로세스 제언입니다. "
        "예제 코드 오류 방지를 위해 오픈소스로 소통하며 집필하고, "
        "두꺼운 바이블 대신 얇고 트렌디한 슬림북 중심의 빠른 개정판 로테이션을 구축해야 합니다."
    )

    # --- Slide 22: 상관관계 히트맵 ---
    s22 = chart_slide(22, "2.4 수치형 변수 간 상관관계 피어슨 상관행렬",
        "09_correlation_heatmap.png", "피어슨 상관계수(r) 결과",
        [("판매지수 vs 리뷰 수", "r = 0.214 (유의미한 정적 상관)"),
         ("판매지수 vs 평점", "r = 0.156 (약한 정적 상관)"),
         ("판매지수 vs 할인율", "r = 0.124 (미약한 정적 상관)"),
         ("할인율 vs 포인트적립", "r = 0.658 (가격 변수 간 수렴)")],
        "상관관계 분석 결과, 판매지수와 가장 높은 연계를 보이는 변수는 '리뷰 수'입니다. 평점 수치보다 리뷰의 절대 수량이 소비자의 구매 안도감을 유도하는 가장 큰 유발 요인입니다.",
        "chart")
    notes(s22,
        "스물두 번째 장표는 수치형 지표들의 피어슨 상관관계 행렬입니다. "
        "판매지수와 가장 강력하게 연동되는 핵심 변수는 리뷰 수(r=0.214)입니다. "
        "판매가격은 매출 성과와 아무런 통계적 연관이 관측되지 않습니다."
    )

    # --- Slide 23: 상관관계 심층분석 ---
    s23 = dual_card_slide(23, "2.4 [심층분석] 리뷰 자산과 소셜 증명의 선순환",
        "별점보다 두터운 리뷰 볼륨",
        [("학습 삽질 비용 회피", "IT 독자가 지불하는 가장 큰 비용은 책값이 아니라 주말 시간입니다. 독자는 실패 리스크를 줄이기 위해 풍부한 리뷰 텍스트를 찾습니다."),
         ("텍스트 피드백의 힘", "구체적인 실습 삽질기나 문제 해결 방안이 기록된 리뷰가 많을수록 구매 결정을 촉진합니다.")],
        "런칭 초기의 리뷰 빌드업 지침",
        [("초기 30일 골든타임", "출간 1개월 이내에 정성스러운 텍스트 및 사진 리뷰가 플랫폼 메인에 노출되도록 서평단을 정교하게 스케줄링해야 합니다."),
         ("인센티브 선순환 구조", "도서 예제 코드 깃허브에 '리뷰 작성 시 추가 소스 제공' 링크를 연계해 자발적 평판 기여를 끊임없이 조장해야 합니다.")])
    notes(s23,
        "스물세 번째 장표는 리뷰 자산 구축 전략입니다. "
        "독자는 학습 피로를 회피하고자 타인의 풍부한 리뷰를 소셜 프루프로 삼습니다. "
        "출간 초기 1개월 골든타임 내에 서평을 집중 런칭하고, 자발적 인증 후기 인센티브를 연계해야 합니다."
    )

    # --- Slide 24: 할인-분철 다변량 ---
    s24 = chart_slide(24, "2.5 [다변량] 할인율-분철 여부별 평균 판매지수",
        "10_discount_spring_vs_sale_index.png", "할인 × 분철 다변량 피봇 결과",
        [("10% 할인 + 분철지원(Y)", "평균: 9,131점 (최상의 성과)"),
         ("10% 할인 + 분철미지원(N)", "평균: 2,321점 (보편 평균)"),
         ("0% 무할인 + 분철지원(Y)", "평균: 3,653점 (할인N 초과)"),
         ("0% 무할인 + 분철미지원(N)", "평균: 603점 (최하 성과)")],
        "10% 할인에 분철 서비스를 결합한 그룹이 압도적 1위 판매지수(9,131)를 기록합니다. 무할인이라도 분철을 장착하면 평균 3,653점으로 분철 없는 할인 도서보다 우수한 성적을 거둡니다.",
        "percent")
    notes(s24,
        "스물네 번째 장표는 할인과 분철을 결합한 다변량 성과 비교입니다. "
        "분철이 제공된다면 무할인이라도 10% 일반 할인 서적의 성과를 가볍게 압도합니다."
    )

    # --- Slide 25: TF-IDF 키워드 ---
    s25 = chart_slide(25, "2.6 도서명 기준 TF-IDF 중요 키워드 분석",
        "11_goods_name_tfidf.png", "상위 TF-IDF 가중치 단어",
        [("핵심 기술 테마", "'AI' (0.059), '파이썬' (0.019), '인공지능' (0.012)"),
         ("AI 보조 도구", "'제미나이' (0.010), '클로드' (0.009), '챗GPT' (0.012)"),
         ("학습 지향 키워드", "'코딩' (0.017), '가이드' (0.016), '입문' (0.011)"),
         ("생산성 키워드", "'활용' (0.011), '실무' (0.009), '실전' (0.009)")],
        "도서 제목 텍스트 마이닝 결과, AI 도구(클로드, 제미나이) 활용 코딩 자동화와 파이썬 기반 데이터 과학 관련 주제가 상위를 석권합니다.",
        "key")
    notes(s25,
        "스물다섯 번째 장표는 도서명 키워드 가중치 마이닝 분석입니다. "
        "시장의 실질 수요는 AI 보조 개발 도구를 이용한 생산성 향상과 파이썬 데이터 과학으로 완전히 전환되었습니다."
    )

    # --- Slide 26: 저자별 판매지수 ---
    s26 = chart_slide(26, "2.7 상위 저자별 도서 수 및 평균 판매지수 분석",
        "12_author_mean_sale_index.png", "저자 세그먼트 성적표",
        [("최다 도서 등록", "놀이교육콘텐츠랩(9권, 평균 1,170점)"),
         ("최고 효율 저자", "오힘찬(3권, 평균 36,848점)"),
         ("스타 저자 쏠림", "소수 저자의 폭발적 성과 집중"),
         ("다작 vs 고효율", "등록 권수와 평균 판매지수 간 반비례")],
        "단순 다작보다 최신 기술 흐름을 명확하게 파헤친 스타 저자의 도서가 압도적인 평균 판매지수 스파이크를 형성합니다. 독자들은 저자의 다작 여부보다 화제성과 실용성에 충성합니다.",
        "users")
    notes(s26,
        "스물여섯 번째 장표는 상위 저자별 성과 분석입니다. "
        "다작 저자들의 지수가 완만한 반면, 최신 생성 AI 트렌드를 저격한 스타 저자의 평균이 수십 배에 달합니다."
    )

    # --- Slide 27: 출판월별 추이 ---
    s27 = chart_slide(27, "2.8 출판 월별 신간 수 및 평균 판매지수 분석",
        "13_publish_month_trends.png", "출판 월별 시즌 성과",
        [("상반기(1~6월)", "신규 진입 도서 집중 및 높은 평균 유지"),
         ("최대 판매지수 월", "12월 (평균 3,985점), 9월 (평균 3,743점)"),
         ("최저 판매지수 월", "7월 (평균 2,006점), 10월 (평균 2,343점)"),
         ("새해/새학기 효과", "방학·신학기 직전 도서 소비 본능 확인")],
        "상반기 발행 도서들의 판매지수 강세가 도드라집니다. 연말연시나 새학기 결심에 맞춰 독자가 지갑을 여는 시즌성 효과입니다. 론칭 일정을 겨울방학이나 신학기 직전에 포커싱해야 합니다.",
        "calendar")
    notes(s27,
        "스물일곱 번째 장표는 출판월별 시즌 패턴 분석입니다. "
        "IT 도서 구매는 새해 다짐이나 방학 공부 등 독자의 결심 시즌에 집중됩니다. "
        "12월이나 2월 직전에 신간 런칭 주기를 맞추는 지혜가 필요합니다."
    )

    # --- Slide 28: 태그 수 vs 판매지수 ---
    s28 = chart_slide(28, "2.9 도서 태그 수(Tag Count)와 판매지수 분석",
        "14_tag_count_vs_sale_index.png", "태그 수에 따른 성과",
        [("적정 태그 수", "2~4개 등록 도서의 판매지수 안정적 우위"),
         ("태그 0개 도서", "685권  |  평균: 1,815점"),
         ("태그 5개 도서", "47권  |  평균: 8,770점"),
         ("검색 엔진 최적화", "정교한 타겟 키워드 매핑 권장")],
        "태그가 없는(0개) 서적은 노출 기회 박탈로 성적이 부진합니다. 2~4개의 정밀한 태그를 매핑한 도서군이 안정적인 흐름을 냅니다. 검색 노출을 위한 최적 태그 관리가 요구됩니다.",
        "tag")
    notes(s28,
        "스물여덟 번째 장표는 태그 수와 판매지수 간 연계입니다. "
        "0~1개 태그를 지양하고, 카테고리 태그 3~4개를 선별 등록하여 검색 최적화를 달성해야 합니다."
    )

    # --- Slide 29: 5대 핵심 전략 (라임 배경) ---
    s29 = prs.slides.add_slide(layout)
    fill_bg(s29, C_LIME)
    # 상단 검정 스트라이프
    add_deco_block(s29, Inches(0), Inches(0), Inches(13.333), Inches(0.25), C_BLACK)
    add_header(s29, "PART 3. IT 도서 시장 활성화를 위한 5대 핵심 전략", "BUSINESS STRATEGY & ROADMAP")
    strategies = [
        ("1. 가격 이원화", "입문 기술서 2.3~2.7만원 표준 가격 + 특수 고난도 전문서 4만원 프리미엄 책정", "price"),
        ("2. 분철 서비스 연계", "500p 이상 실무서 분철 옵션 기본 제공으로 26%+ 추가 매출 동력 확보", "spring"),
        ("3. 지식 감가상각 방어", "트렌드 민감 응용서는 린(Lean) 집필 공정으로 빠르게 출판하고 수명 선제 관리", "calendar"),
        ("4. 소셜프루프 극대화", "초기 서평단 1개월 빌드업 전략 + 자발적 리뷰 축적 인센티브 설계", "review"),
        ("5. 검색 노출 태깅", "무질서한 태깅 지양, 핵심 키워드 3~4개 선별 집중 매핑으로 유입 전환율 향상", "key"),
    ]
    for idx, (title, desc, icon) in enumerate(strategies):
        row = idx // 3
        col = idx % 3
        if row == 1:
            x = Inches(2.35 + (idx - 3) * 4.35)
        else:
            x = Inches(0.6 + col * 4.1)
        y = Inches(1.8 + row * 2.55)
        w = Inches(3.9)
        h = Inches(2.3)
        neo_card(s29, x, y, w, h, fill=C_WHITE)
        neo_icon(s29, icon, x + Inches(0.3), y + Inches(0.3), Inches(0.45))
        tb = s29.shapes.add_textbox(x + Inches(0.3), y + Inches(0.9), w - Inches(0.6), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_TITLE
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = C_BLACK
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = C_DARK_GRAY
        p2.line_spacing = 1.25
        p2.space_before = Pt(6)
    add_footer(s29, 29)
    notes(s29,
        "스물아홉 번째 장표는 이번 EDA의 핵심 액션 아이템인 5대 전략입니다. "
        "라임(#CCFF00) 배경 위에 5개의 흰색 카드가 3+2 그리드로 배치되어 강렬한 대비를 이룹니다. "
        "가격 이원화, 분철 서비스 연계, 린 출판 모델 도입, 초기 30일 서평 유치, "
        "온라인 태그 최적화의 다섯 가지 전략을 각각 아이콘과 함께 제시합니다."
    )

    # --- Slide 30: 결론 및 Q&A (옐로우 배경) ---
    s30 = prs.slides.add_slide(layout)
    fill_bg(s30, C_YELLOW)
    # 데코 블록 (의도적 미스얼라인먼트)
    add_deco_block(s30, Inches(0.5), Inches(0.3), Inches(2.5), Inches(1.0), C_BRICK, rotation=-6)
    add_deco_block(s30, Inches(10.5), Inches(5.5), Inches(2.5), Inches(1.5), C_BLUE, rotation=4)
    # THANK YOU 거대 텍스트
    ty = s30.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12.1), Inches(1.5))
    tf = ty.text_frame
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.name = FONT_TITLE
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = C_BLACK
    # 중앙 카드
    neo_card(s30, Inches(1.5), Inches(2.3), Inches(10.333), Inches(4.2))
    cb = s30.shapes.add_textbox(Inches(2.0), Inches(2.6), Inches(9.3), Inches(3.8))
    tf2 = cb.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_bottom = tf2.margin_right = 0
    # Q&A 타이틀
    p1 = tf2.paragraphs[0]
    p1.text = "Q&A 및 미래 출판 패러다임 예측"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = C_BRICK
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(20)
    infos = [
        "데이터 분석 결과 요약본 배포: PDF 보고서 및 데이터셋 시각화 템플릿 별도 제공",
        "질의응답: 분석 설계 모형 및 비즈니스 의사결정 액션 플랜 세부 질문",
    ]
    for info in infos:
        p = tf2.add_paragraph()
        p.text = info
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_BLACK
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)
    # 맺음말
    p_close_t = tf2.add_paragraph()
    p_close_t.text = "[ 맺음말 ]"
    p_close_t.font.name = FONT_BODY
    p_close_t.font.size = Pt(13)
    p_close_t.font.bold = True
    p_close_t.font.color.rgb = C_DARK_GRAY
    p_close_t.alignment = PP_ALIGN.CENTER
    p_close_t.space_before = Pt(16)
    p_close_t.space_after = Pt(6)
    p_close = tf2.add_paragraph()
    p_close.text = (
        "이번 EDA 분석을 통해, IT 도서 시장은 독자의 편의성(분철)과 기술 트렌드의 "
        "신속성(린 출판), 그리고 소셜 평판(리뷰 건수)에 의해 지배를 받는 특수한 "
        "비즈니스 필드임이 검증되었습니다. 본 리포트의 전략을 바탕으로 더욱 민첩하고 "
        "유연한 비즈니스 로드맵이 수립되기를 기대합니다."
    )
    p_close.font.name = FONT_BODY
    p_close.font.size = Pt(12)
    p_close.font.color.rgb = C_BLACK
    p_close.line_spacing = 1.4
    p_close.alignment = PP_ALIGN.CENTER
    add_footer(s30, 30)
    notes(s30,
        "이상으로 발표를 모두 마치겠습니다. "
        "네오브루탈리즘의 거침없고 명확한 선들처럼, 우리 데이터 역시 명백한 비즈니스 유기성을 나타내고 있습니다. "
        "앞선 5대 전략을 통해 실무적으로 개발 서적의 가치를 전개할 예정이며, "
        "추가 기술 모형 설계나 비즈니스 실행안에 대해 질의해 주시면 감사하겠습니다. "
        "옐로우 배경 위에 'THANK YOU'가 72포인트로 거대하게 배치되어 "
        "네오브루탈리즘 특유의 강렬한 마무리 인상을 남깁니다."
    )

    # =========================================================================
    # 저장
    # =========================================================================
    output = "yes24/docs/EDA_Slide_30Pages_NeoBrutalism_v2.pptx"
    prs.save(output)
    print(f"✅ 네오브루탈리즘 v2 PPTX 생성 완료: {output}")


if __name__ == "__main__":
    create_neobrutalist_v2()
